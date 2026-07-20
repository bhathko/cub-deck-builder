from pathlib import Path
import re
import shutil

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path("my_project")
SRC = ROOT / "source" / "slides.md"
OUT = ROOT / "my_project_review.pptx"
PREVIEW_DIR = ROOT / "qa_preview_v2"
COVER_BG = ROOT / "assets" / "backgrounds" / "cover_bg.png"
CONTENT_BG = ROOT / "assets" / "backgrounds" / "content_bg.png"
LOGO = ROOT / "assets" / "logos" / "cathay_logo.png"

FONT = "Microsoft JhengHei"
DARK = RGBColor(52, 66, 82)
MUTED = RGBColor(104, 114, 126)
GREEN = RGBColor(88, 212, 148)
PURPLE = RGBColor(132, 139, 242)
BLUE = RGBColor(74, 183, 249)
LIGHT = RGBColor(238, 243, 247)
LINE = RGBColor(212, 220, 228)
WHITE = RGBColor(255, 255, 255)

SLIDE_W = 13.333333
SLIDE_H = 7.5


def parse_slides(text):
    pattern = re.compile(r"^## Slide\s+(\d+)\s*$", re.M)
    matches = list(pattern.finditer(text))
    result = {}
    for i, match in enumerate(matches):
        start = match.end()
        end = matches[i + 1].start() if i + 1 < len(matches) else len(text)
        result[int(match.group(1))] = text[start:end]
    return result


def section(block, heading):
    match = re.search(
        r"^###\s+" + re.escape(heading) + r"\s*\n(.*?)(?=^###\s+|^---\s*$|\Z)",
        block,
        re.S | re.M,
    )
    return match.group(1).strip() if match else ""


def split_kv(text):
    for sep in ("：", ":"):
        if sep in text:
            left, right = text.split(sep, 1)
            return left.strip(), right.strip()
    return text.strip(), ""


def required(block):
    data = {}
    for line in section(block, "必須出現的文字").splitlines():
        stripped = line.strip()
        if stripped.startswith("- ") and ("：" in stripped or ":" in stripped):
            k, v = split_kv(stripped[2:].strip())
            data[k] = v
    return data


def nested_items(block, heading="內容"):
    raw = section(block, heading)
    items = []
    current = None
    for line in raw.splitlines():
        if not line.strip().startswith("- "):
            continue
        indent = len(line) - len(line.lstrip(" "))
        text = line.strip()[2:].strip()
        if indent == 0:
            current = {"title": text, "children": []}
            items.append(current)
        elif current is not None:
            current["children"].append(text)
    return items


def toc_items(block):
    items = []
    current = None
    for line in section(block, "必須出現的文字").splitlines():
        stripped = line.strip()
        if stripped.startswith("- 第 "):
            if current:
                items.append(current)
            current = {"title": "", "subtitle": ""}
        elif current and stripped.startswith("- 大標"):
            current["title"] = split_kv(stripped[2:].strip())[1]
        elif current and stripped.startswith("- 副標"):
            current["subtitle"] = split_kv(stripped[2:].strip())[1]
    if current:
        items.append(current)
    return items


def add_bg(slide, path):
    slide.shapes.add_picture(str(path), 0, 0, width=Inches(SLIDE_W), height=Inches(SLIDE_H))


def fmt_text_frame(tf, margin=0.06):
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)


def set_text(shape, text, size=16, bold=False, color=DARK, align=None):
    tf = shape.text_frame
    tf.clear()
    fmt_text_frame(tf)
    for i, line in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if align is not None:
            p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


def textbox(slide, text, x, y, w, h, size=16, bold=False, color=DARK, align=None):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    set_text(shape, text, size=size, bold=bold, color=color, align=align)
    return shape


def rect(slide, x, y, w, h, fill=WHITE, line=LINE, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(0.8)
    return shape


def label(slide, text, x, y, w, h, fill=GREEN, size=15):
    box = rect(slide, x, y, w, h, fill=fill, line=fill)
    set_text(box, text, size=size, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    return box


def card(slide, x, y, w, h, title, body="", title_color=DARK, fill=WHITE):
    rect(slide, x, y, w, h, fill=fill, line=LINE, radius=True)
    textbox(slide, title, x + 0.16, y + 0.12, w - 0.32, 0.32, size=18, bold=True, color=title_color)
    if body:
        textbox(slide, body, x + 0.16, y + 0.55, w - 0.32, h - 0.68, size=14, color=DARK)


def add_logo_and_page(slide, page_no=None):
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(0.32), Inches(6.92), width=Inches(1.34), height=Inches(0.34))
    if page_no is not None:
        textbox(slide, str(page_no), 12.55, 6.86, 0.45, 0.28, size=14, bold=True, color=DARK, align=PP_ALIGN.RIGHT)


def header(slide, title, subtitle, page_no):
    textbox(slide, title, 0.32, 0.26, 11.5, 0.42, size=32, bold=True, color=DARK)
    if subtitle:
        textbox(slide, subtitle, 0.34, 0.83, 11.8, 0.32, size=18, color=MUTED)
    add_logo_and_page(slide, page_no)


def bullets(lines):
    return "\n".join("• " + line for line in lines)


def line(slide, x1, y1, x2, y2, color=LINE, width=1):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    return conn


def polygon(slide, points, fill, line_color=WHITE, line_width=0.6):
    builder = slide.shapes.build_freeform(points[0][0], points[0][1], scale=Inches(1))
    builder.add_line_segments(points[1:], close=True)
    shape = builder.convert_to_shape()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line_color
    shape.line.width = Pt(line_width)
    return shape


def build_deck():
    source = SRC.read_text(encoding="utf-8")
    slides = parse_slides(source)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank = prs.slide_layouts[6]

    # Slide 1
    slide = prs.slides.add_slide(blank)
    add_bg(slide, COVER_BG)
    req = required(slides[1])
    label(slide, "內部簡報", 0.48, 2.07, 1.35, 0.44, fill=DARK, size=15)
    textbox(slide, req["主標題"], 0.34, 2.85, 3.2, 0.52, size=40, bold=True, color=DARK)
    textbox(slide, req["副標題"], 3.30, 2.85, 2.3, 0.52, size=36, bold=True, color=GREEN)
    textbox(slide, f"{req['日期']}  |  {req['報告人']}", 0.34, 3.62, 5.6, 0.34, size=18, bold=True, color=MUTED)
    add_logo_and_page(slide)

    # Slide 2
    slide = prs.slides.add_slide(blank)
    add_bg(slide, CONTENT_BG)
    textbox(slide, "Contents", 0.62, 1.05, 2.8, 0.34, size=14, color=MUTED)
    textbox(slide, "目錄", 0.62, 1.37, 2.1, 0.62, size=32, bold=True, color=DARK)
    for i, item in enumerate(toc_items(slides[2])):
        y = 1.42 + i * 1.10
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4.2), Inches(y), Inches(0.36), Inches(0.36))
        circ.fill.background()
        circ.line.color.rgb = MUTED
        circ.line.width = Pt(0.8)
        set_text(circ, str(i + 1), size=12, color=MUTED, align=PP_ALIGN.CENTER)
        circ.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        textbox(slide, item["title"], 4.78, y - 0.02, 4.2, 0.34, size=20, bold=True, color=DARK)
        if item["subtitle"] != "無":
            textbox(slide, item["subtitle"], 4.79, y + 0.38, 5.8, 0.32, size=14, color=MUTED)
    add_logo_and_page(slide, 2)

    # Slide 3
    slide = prs.slides.add_slide(blank)
    add_bg(slide, CONTENT_BG)
    req = required(slides[3])
    header(slide, req["主標題"], req["副標題"], 3)
    hub = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.18), Inches(2.26), Inches(2.2), Inches(2.2))
    hub.fill.solid()
    hub.fill.fore_color.rgb = DARK
    hub.line.color.rgb = DARK
    set_text(hub, "年度目標\n更快決策\n更穩服務", size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    hub.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    items3 = [
        ("核心任務", "串接流程、資料、服務與人才能力"),
        ("2 項大型專案", "數據治理｜服務體驗升級"),
        ("12 個里程碑", "對齊 Q1-Q4 推進節奏"),
        ("40%", "報表處理時間節省目標"),
        ("30%", "服務回應一致性提升目標"),
        ("年度目標", "管理更快、服務更穩、能力可複製"),
    ]
    positions = [(0.55, 1.65), (0.55, 3.10), (0.55, 4.55), (8.35, 1.65), (8.35, 3.10), (8.35, 4.55)]
    for (title, body), (x, y) in zip(items3, positions):
        card(slide, x, y, 3.75, 0.84, title, body, title_color=GREEN)
        line(slide, x + (3.75 if x < 5 else 0), y + 0.42, 6.28, 3.36, color=LINE, width=0.9)

    # Slide 4
    slide = prs.slides.add_slide(blank)
    add_bg(slide, CONTENT_BG)
    req = required(slides[4])
    header(slide, req["主標題"], req["副標題"], 4)
    content = nested_items(slides[4])
    cols = [(0.50, GREEN), (4.58, DARK), (8.66, BLUE)]
    for (x, color), item in zip(cols, content):
        label(slide, item["title"], x, 1.55, 3.55, 0.42, fill=color, size=17)
        card(slide, x, 2.08, 3.55, 3.88, "", bullets(item["children"]), title_color=color)

    # Slide 5
    slide = prs.slides.add_slide(blank)
    add_bg(slide, CONTENT_BG)
    req = required(slides[5])
    header(slide, req["主標題"], req["副標題"], 5)
    data = nested_items(slides[5])
    label(slide, "Before", 0.70, 1.50, 4.15, 0.42, fill=GREEN, size=18)
    label(slide, "After", 8.48, 1.50, 4.15, 0.42, fill=PURPLE, size=18)
    card(slide, 0.70, 2.08, 4.15, 2.70, "改善前", bullets(data[0]["children"]), title_color=GREEN)
    card(slide, 8.48, 2.08, 4.15, 2.70, "改善後", bullets(data[1]["children"]), title_color=PURPLE)
    arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(5.42), Inches(2.93), Inches(2.38), Inches(0.72))
    arr.fill.solid()
    arr.fill.fore_color.rgb = RGBColor(223, 229, 234)
    arr.line.color.rgb = RGBColor(223, 229, 234)
    label(slide, "關鍵改善目標", 0.70, 5.22, 11.93, 0.42, fill=DARK, size=18)
    for i, txt in enumerate(data[2]["children"]):
        x = 0.92 + i * 3.92
        k, v = split_kv(txt)
        card(slide, x, 5.82, 3.35, 0.70, k, v, title_color=GREEN)

    # Slide 6
    slide = prs.slides.add_slide(blank)
    add_bg(slide, CONTENT_BG)
    req = required(slides[6])
    header(slide, req["主標題"], req["副標題"], 6)
    opts = nested_items(slides[6])
    xs = [0.55, 4.55, 8.55]
    colors = [GREEN, PURPLE, DARK]
    for x, color, item in zip(xs, colors, opts[:3]):
        name, sub = split_kv(item["title"])
        label(slide, name, x, 1.44, 3.45, 0.42, fill=color, size=17)
        textbox(slide, sub, x, 1.93, 3.45, 0.26, size=16, bold=True, color=color, align=PP_ALIGN.CENTER)
        pro = item["children"][0].replace("優點：", "")
        con = item["children"][1].replace("缺點：", "")
        card(slide, x, 2.36, 3.45, 1.35, "優點", pro, title_color=GREEN)
        card(slide, x, 3.92, 3.45, 1.35, "待改善", con, title_color=PURPLE)
    rec = opts[3]["children"] if len(opts) > 3 else []
    label(slide, "建議方向", 0.55, 5.58, 1.55, 0.42, fill=DARK, size=16)
    textbox(slide, "　".join(rec), 2.22, 5.60, 10.4, 0.38, size=16, bold=True, color=DARK)

    # Slide 7
    slide = prs.slides.add_slide(blank)
    add_bg(slide, CONTENT_BG)
    req = required(slides[7])
    header(slide, req["主標題"], req["副標題"], 7)
    story = nested_items(slides[7])
    phases = [
        ("過去", story[1]["children"][0].replace("過去：", ""), PURPLE),
        ("現在", story[1]["children"][1].replace("現在：", ""), GREEN),
        ("未來", story[1]["children"][2].replace("未來：", ""), DARK),
    ]
    for x, (name, body, color) in zip([0.6, 4.55, 8.50], phases):
        label(slide, name, x, 1.52, 3.45, 0.42, fill=color, size=18)
        card(slide, x, 2.12, 3.45, 1.42, "故事", body, title_color=color)
    card(slide, 0.6, 4.03, 3.45, 1.88, "專案背景", bullets(story[0]["children"][:2]), title_color=DARK)
    card(slide, 4.55, 4.03, 3.45, 1.88, "改善重點", bullets(["服務分級", "回應腳本", "知識庫"]), title_color=GREEN)
    card(slide, 8.50, 4.03, 3.45, 1.88, "預期成果", bullets(story[2]["children"]), title_color=BLUE)

    # Slide 8
    slide = prs.slides.add_slide(blank)
    add_bg(slide, CONTENT_BG)
    req = required(slides[8])
    textbox(slide, req["主標題"], 0.32, 0.26, 8.9, 0.42, size=32, bold=True, color=DARK)
    textbox(slide, req["副標題"], 0.34, 0.80, 8.7, 0.28, size=14, color=MUTED)
    add_logo_and_page(slide, 8)
    levels = nested_items(slides[8])
    lane = rect(slide, 0.88, 1.10, 8.48, 5.05, fill=RGBColor(247, 249, 250), line=RGBColor(247, 249, 250))
    lane.line.transparency = 100000

    # Native editable reconstruction of template page 54:
    # large cropped layered pyramid on the left, five horizontal explanation
    # bands in the middle, and a two-section note card on the right.
    pyr = [
        (0.00, 1.78, 1.05, 0.44, GREEN),
        (0.00, 2.72, 1.42, 0.50, RGBColor(52, 183, 112)),
        (0.00, 3.72, 1.98, 0.50, RGBColor(34, 128, 82)),
        (0.00, 4.78, 2.62, 0.54, RGBColor(92, 121, 150)),
        (0.00, 5.90, 3.08, 0.58, DARK),
    ]
    for i, (x, y, w, h, color) in enumerate(pyr):
        polygon(slide, [(x - 0.55, y), (x + w, y), (x + w + 0.46, y + h), (x - 0.32, y + h)], color)
        polygon(slide, [(x + w, y), (x + w + 0.34, y - 0.28), (x + w + 0.78, y + h - 0.28), (x + w + 0.46, y + h)], RGBColor(max(0, color[0] - 28), max(0, color[1] - 28), max(0, color[2] - 28)))
        polygon(slide, [(x - 0.55, y), (x + w, y), (x + w + 0.34, y - 0.28), (x - 0.20, y - 0.28)], RGBColor(min(255, color[0] + 80), min(255, color[1] + 55), min(255, color[2] + 55)), line_color=RGBColor(245, 248, 250))

    label_specs = [
        (0.05, 1.55, 0.95, 0.42, GREEN),
        (0.26, 2.48, 1.20, 0.48, RGBColor(52, 183, 112)),
        (0.52, 3.50, 1.32, 0.48, RGBColor(34, 128, 82)),
        (0.72, 4.62, 1.55, 0.48, RGBColor(92, 121, 150)),
        (1.02, 5.74, 1.78, 0.48, DARK),
    ]
    for (x, y, w, h, color), item in zip(label_specs, reversed(levels[:5])):
        _, name = split_kv(item["title"])
        box = rect(slide, x, y, w, h, fill=color, line=color)
        set_text(box, name, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    row_specs = [
        (0.95, 1.22, 7.95, 0.74, GREEN),
        (1.45, 2.18, 7.80, 0.74, GREEN),
        (2.05, 3.14, 7.40, 0.74, GREEN),
        (2.65, 4.12, 6.80, 0.74, DARK),
        (3.25, 5.22, 6.20, 0.74, DARK),
    ]
    for (x, y, w, h, color), item in zip(row_specs, reversed(levels[:5])):
        rect(slide, x, y, w, h, fill=WHITE, line=WHITE)
        lvl, name = split_kv(item["title"])
        textbox(slide, f"{lvl}｜{name}", x + 0.08, y + 0.08, w - 0.16, 0.24, size=18, bold=True, color=color)
        textbox(slide, item["children"][0], x + 0.08, y + 0.38, w - 0.16, 0.26, size=13, color=DARK)

    rect(slide, 9.05, 1.08, 3.95, 5.28, fill=RGBColor(239, 246, 247), line=LINE, radius=True)
    left_bar = rect(slide, 9.25, 1.38, 0.06, 2.48, fill=GREEN, line=GREEN)
    left_bar.line.transparency = 100000
    textbox(slide, "能力堆疊重點", 9.52, 1.48, 2.7, 0.34, size=22, bold=True, color=GREEN)
    textbox(slide, "由服務標準化開始，逐層累積知識、流程、訓練與持續優化機制。", 9.52, 1.92, 2.78, 0.76, size=14, color=DARK)
    textbox(slide, "# 標準化　# 文件化　# 協作化", 9.52, 3.38, 2.8, 0.28, size=12, color=GREEN)
    dark_bar = rect(slide, 9.25, 4.38, 0.06, 1.84, fill=DARK, line=DARK)
    dark_bar.line.transparency = 100000
    textbox(slide, "年度落地", 9.68, 4.56, 2.2, 0.34, size=22, bold=True, color=DARK)
    textbox(slide, "Q2 建立標準\nQ3 推動訓練\nQ4 固化機制", 9.68, 4.98, 2.25, 0.68, size=14, color=DARK)
    textbox(slide, "# 訓練　# 追蹤　# 優化", 9.68, 5.86, 2.5, 0.28, size=12, color=DARK)

    # Slide 9
    slide = prs.slides.add_slide(blank)
    add_bg(slide, CONTENT_BG)
    req = required(slides[9])
    header(slide, req["主標題"], req["副標題"], 9)
    rows = nested_items(slides[9])
    quarters = ["Q1\n盤點", "Q2\n建置", "Q3\n推廣", "Q4\n固化"]
    x_start = 1.68
    col_w = 2.72
    for i, q in enumerate(quarters):
        label(slide, q, x_start + i * col_w, 1.40, 2.35, 0.52, fill=GREEN if i < 3 else DARK, size=16)
    row_defs = [(rows[0]["title"], rows[0]["children"], 2.20), (rows[1]["title"], rows[1]["children"], 3.72), (rows[2]["title"], rows[2]["children"][:4], 5.24)]
    for title, vals, y in row_defs:
        textbox(slide, title.replace("：", ""), 0.42, y + 0.28, 1.08, 0.55, size=15, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
        for i, val in enumerate(vals[:4]):
            card(slide, x_start + i * col_w, y, 2.35, 1.08, "", val, title_color=GREEN)

    # Slide 10
    slide = prs.slides.add_slide(blank)
    add_bg(slide, COVER_BG)
    textbox(slide, "Thank you", 0.86, 3.28, 4.2, 0.55, size=32, bold=True, color=DARK)

    prs.save(OUT)
    return OUT


def export_preview(pptx_path):
    if PREVIEW_DIR.exists():
        shutil.rmtree(PREVIEW_DIR)
    PREVIEW_DIR.mkdir(parents=True)
    try:
        import win32com.client
    except Exception as exc:
        print(f"PREVIEW_EXPORT_SKIPPED: {exc}")
        return
    app = win32com.client.DispatchEx("PowerPoint.Application")
    app.Visible = 1
    presentation = app.Presentations.Open(str(pptx_path.resolve()), WithWindow=False)
    presentation.Export(str(PREVIEW_DIR.resolve()), "PNG", 1600, 900)
    presentation.Close()
    app.Quit()


if __name__ == "__main__":
    path = build_deck()
    export_preview(path)
    print(path)
    print(PREVIEW_DIR)
