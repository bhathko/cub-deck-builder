#!/usr/bin/env python3
"""build_min_dist — 產生實驗B 的兩個「最小包」zip。

實驗B 的中間頁由 GPT 自由設計,**不使用頁型庫**。但完整的模板包與 tools 會
把整套頁型資訊帶進 GPT 的 Knowledge:

- `template.pptx` 60 頁,每個頁型一頁實頁(模型打得開就看得到所有版面)
- `manifest.page_types` 52 個契約、`bindings.fills` 33 個綁定、`page_map.md`
  86 行對照表
- `tools.zip` 裡的 `make_skeleton.py` 是實驗A 的頁型選版器

instructions 的 Step 0 又要求印出 manifest,等於每次產檔前都打開那份頁型清單。
只擋 `page_types*.md` 沒有用——同樣的資訊從 zip 整包進去了。

本腳本把兩個包裁到實驗B 真正需要的範圍:模板只留 cover/agenda/closing 三頁
實頁與對應契約(結構三頁走工具鏈,少了就產不出來),tools 只留 render_deck /
qa_check 及其相依。裁完全包搜不到任何其他頁型名稱。

**唯一已知殘留**:Knowledge 另外上傳的 `validate_slide_spec_gpts.py` 內含 29 個
頁型名稱。它是結構三頁的閘門(DEPLOY 驗收 ④ 靠它),拿掉 B 就沒有守門;要清就
得為 B 叉一份 engine 程式碼,validator 一改就漂,所以刻意保留。實際影響很小
——B 有名字,但沒有綁定、沒有模板實頁、沒有選版器,產不出那些頁型。

輸入(不改動):
  gpts_experiment_a/dist/template_light.zip   expA 的 content_bg 改版包
  engine/tools/*.py                           工具真相

輸出:
  gpts_experiment_b/dist/template_light.zip
  gpts_experiment_b/dist/tools.zip

用法:python gpts_experiment_b/build_min_dist.py
"""
from __future__ import annotations

import hashlib
import json
import sys
import tempfile
import zipfile
from pathlib import Path

REPO = Path(__file__).resolve().parent.parent
SRC_PACK = REPO / "gpts_experiment_a" / "dist" / "template_light.zip"
TOOLS = REPO / "engine" / "tools"
DIST = Path(__file__).resolve().parent / "dist"

# 結構三頁——需求本身規定走公司模板與工具鏈,不可再裁
KEEP_TYPES = ("cover", "agenda", "closing")

# 實驗B 只跑 render_deck 與 qa_check;這是它們的完整相依閉包
# (make_skeleton/capacity_probe/audit_provenance/inspect_template/run_pipeline
#  與 README_TOOLS.md 都只服務頁型庫流程,不帶)
KEEP_TOOLS = ("render_deck.py", "qa_check.py", "pack_loader.py", "fills_engine.py",
              "fill_helpers.py", "pptx_toolkit.py", "text_tools.py")

RID = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
FREE_LAYOUT = "空白(白底)"


def zip_add(z: zipfile.ZipFile, src: Path, arcname: str) -> None:
    """可重現寫入:固定時戳與權限,內容沒變則 zip sha 不變(同 template_admin)。"""
    zi = zipfile.ZipInfo(arcname, date_time=(1980, 1, 1, 0, 0, 0))
    zi.external_attr = 0o644 << 16
    zi.compress_type = zipfile.ZIP_DEFLATED
    z.writestr(zi, src.read_bytes())


def sha256(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def normalize_pptx(path: Path) -> None:
    """把 pptx(本身也是 zip)內層條目的時戳固定住。

    python-pptx 的 save() 用當下時間寫每個 part,於是「同樣輸入、兩次重打」會
    產出不同 bytes → template_light.zip 的 sha 每次都變 → git 每次都存一顆新的
    9MB blob。條目順序照原樣保留([Content_Types].xml 仍在最前)。
    """
    with zipfile.ZipFile(path) as z:
        items = [(i, z.read(i.filename)) for i in z.infolist()]
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as z:
        for info, data in items:
            zi = zipfile.ZipInfo(info.filename, date_time=(1980, 1, 1, 0, 0, 0))
            zi.external_attr = info.external_attr
            zi.compress_type = zipfile.ZIP_DEFLATED
            z.writestr(zi, data)


def build_template_pack(work: Path) -> Path:
    from pptx import Presentation

    src = work / "full"
    src.mkdir()
    with zipfile.ZipFile(SRC_PACK) as z:
        z.extractall(src)

    manifest = json.loads((src / "manifest.json").read_text(encoding="utf-8-sig"))
    bindings = json.loads((src / "bindings.json").read_text(encoding="utf-8-sig"))
    keep_pages = [manifest["page_types"][k]["template_page"] for k in KEEP_TYPES]
    print(f"  來源 {manifest['template_id']}@{manifest['version']}"
          f"({manifest.get('page_count')} 頁,{len(manifest['page_types'])} 個頁型)")
    print(f"  保留模板頁:{dict(zip(KEEP_TYPES, keep_pages))}")

    out_pptx = work / "template.pptx"
    prs = Presentation(src / "template.pptx")
    n_layouts = len(prs.slide_layouts)
    sld_lst = prs.slides._sldIdLst
    for i, sld in enumerate(list(sld_lst), start=1):
        if i not in keep_pages:
            # drop_rel 才會讓 part 真的不被序列化;只從 sldIdLst 移除會留下孤兒 XML
            prs.part.drop_rel(sld.get(RID))
            sld_lst.remove(sld)
    prs.save(out_pptx)
    normalize_pptx(out_pptx)

    chk = Presentation(out_pptx)
    n_slides = len(chk.slides._sldIdLst)
    layouts = [l.name for l in chk.slide_layouts]
    if n_slides != len(KEEP_TYPES):
        sys.exit(f"✗ 裁切後頁數 {n_slides},應為 {len(KEEP_TYPES)}")
    if FREE_LAYOUT not in layouts:
        sys.exit(f"✗ 自由頁 layout「{FREE_LAYOUT}」不見了")
    with zipfile.ZipFile(out_pptx) as z:
        n_xml = len([n for n in z.namelist()
                     if n.startswith("ppt/slides/slide") and n.endswith(".xml")])
    if n_xml != len(KEEP_TYPES):
        sys.exit(f"✗ pptx 內殘留 {n_xml} 個 slide XML,孤兒 part 沒清乾淨")
    print(f"  template.pptx:60 → {n_slides} 頁(slide XML {n_xml} 個);"
          f"版面配置 {n_layouts} → {len(layouts)},{FREE_LAYOUT} 保留")

    manifest["page_types"] = {k: {"mode": "fill", "template_page": i + 1}
                              for i, k in enumerate(KEEP_TYPES)}
    if isinstance(manifest.get("capacity_overrides"), dict):
        manifest["capacity_overrides"] = {
            k: v for k, v in manifest["capacity_overrides"].items() if k in KEEP_TYPES}
    manifest["page_count"] = len(KEEP_TYPES)
    if not manifest["version"].endswith("-min"):
        manifest["version"] += "-min"
    manifest["template_sha256"] = sha256(out_pptx)
    manifest["_note_min"] = (
        "實驗B 專用最小包(gpts_experiment_b/build_min_dist.py 產生):只含 "
        "cover/agenda/closing。中間頁由 GPT 自由設計、不走頁型庫,故不含其餘頁型"
        "契約、模板實頁與 page_map.md。")

    # render_deck 讀的模板頁碼在 bindings.fills,不是 manifest——兩邊都要重編
    bindings["fills"] = {k: v for k, v in bindings["fills"].items() if k in KEEP_TYPES}
    for i, k in enumerate(KEEP_TYPES):
        bindings["fills"][k]["template_page"] = i + 1

    (work / "manifest.json").write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")
    (work / "bindings.json").write_text(
        json.dumps(bindings, ensure_ascii=False, indent=2), encoding="utf-8")

    out = DIST / "template_light.zip"
    with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as z:
        for n in ("template.pptx", "manifest.json", "bindings.json"):   # 不含 page_map.md
            zip_add(z, work / n, n)
        for f in sorted((src / "assets").rglob("*")):
            if f.is_file():
                zip_add(z, f, f"assets/{f.relative_to(src / 'assets').as_posix()}")
    print(f"  → {out.name}  {out.stat().st_size / 1024 / 1024:.1f} MB"
          f"  {manifest['template_id']}@{manifest['version']}")
    return out


def build_tools() -> Path:
    missing = [n for n in KEEP_TOOLS if not (TOOLS / n).exists()]
    if missing:
        sys.exit(f"✗ engine/tools 缺檔:{missing}")
    out = DIST / "tools.zip"
    with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as z:
        for n in KEEP_TOOLS:
            zip_add(z, TOOLS / n, n)
    print(f"  → {out.name}  {len(KEEP_TOOLS)} 支,{out.stat().st_size / 1024:.0f} KB")
    return out


def audit(paths) -> None:
    """裁完不該再找得到任何頁型名稱(結構三頁除外)。"""
    import re
    pat = re.compile(r"(info|stage|data|phase|cycle|pyramid|vision|structure|evaluation|story)"
                     r"_[a-z_]{4,}")
    hits = set()
    for p in paths:
        with zipfile.ZipFile(p) as z:
            for n in z.namelist():
                if n.endswith((".json", ".md", ".py", ".xml")):
                    hits |= set(pat.findall(z.read(n).decode("utf-8", "ignore")))
    print(f"\n殘留頁型名稱稽核:{'PASS(零殘留)' if not hits else f'FAIL {sorted(hits)}'}")
    if hits:
        sys.exit(1)


def main() -> int:
    if not SRC_PACK.exists():
        sys.exit(f"✗ 找不到來源模板包:{SRC_PACK}")
    DIST.mkdir(exist_ok=True)
    print("打包實驗B 最小 dist")
    with tempfile.TemporaryDirectory() as td:
        pack = build_template_pack(Path(td))
    tools = build_tools()
    audit([pack, tools])
    print()
    for p in (pack, tools):
        print(f"sha256 {sha256(p)}  {p.name}")
    print("\n提醒:zip 改了就要更新 instructions 版本字串,並重跑 DEPLOY 驗收。")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
