#!/usr/bin/env python3
"""render_deck — 通用 renderer:spec(+選配 plan)→ .pptx(冪等,全檔重生)。

核心原則:**修計畫,不修產出**。每次執行都從 light_template.pptx 重新生成整份
簡報;產出不對就改輸入(spec 或 render_plan.json)再跑一次,禁止對產出 pptx
做局部修補(那正是「刪一個壞的再產一個壞的」循環的來源)。

★ 10 種註冊頁型完全自動,不需要 plan:
   builtin(版面內建):cover, agenda, closing, story_chapter_statement,
                        stage_dual_track_roadmap
   fills(模板頁自動填充,見 fills.py):vision_goal_center_balance,
        info_three_column_category, data_two_group_metric_comparison,
        evaluation_option_score_pros_cons, pyramid_layered_maturity_detail
   → 整份 spec 都是註冊頁型時,直接跑本工具即可,LLM 不需要產任何計畫。

★ 只有「未涵蓋頁型」(page_types.md 頁型庫其他頁型)才需要 plan 條目:
   clone 模板頁 + 文字替換編輯清單。plan 也可覆寫自動頁(以 plan 為準)。

用法:
  python render_deck.py --spec /mnt/data/slide_spec.json \
      [--plan /mnt/data/render_plan.json] \
      --template /mnt/data/light_template.pptx \
      --asset-dir /mnt/data --out /mnt/data/deck.pptx

render_plan.json 格式(僅未涵蓋頁需要;內容文字必須逐字取自 spec):
{
  "slides": [
    {"number": 5, "mode": "clone", "template_page": 35,
     "edits": [
       {"match": {"id": 23}, "text": "..."},                  // shape_id 最精準
       {"match": {"contains": "模板現有文字"}, "text": "..."},
       {"match": {"name": "Title 1"}, "text": "...", "nth": 0} // 撞多筆用 nth
     ],
     "delete": [ {"match": {"id": 45}} ],
     "shrink": true}
  ]
}

行為保證:
- 有 UNMATCHED / AMBIGUOUS / FillError → 印出清單、exit 1;檔案仍會存
  (方便檢查)但**不得交付**,修輸入後整檔重跑。
- 頁碼規則自動:內容頁右下 28pt #344252;封面/封底無;clone/fills 頁會先清掉
  模板原本右下角的純數字頁碼框。
- 產出只含 spec 的頁:模板原頁全刪、Section 全清、順序=spec 順序。
"""
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

import fills
import pptx_toolkit as tk
import text_tools as tt

DARK = RGBColor(*tk.COLOR_DARK)
MUTED = RGBColor(*tk.COLOR_MUTED)
GREEN = RGBColor(*tk.COLOR_GREEN)
PURPLE = RGBColor(*tk.COLOR_PURPLE)
LINE = RGBColor(*tk.COLOR_LINE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


# ---------------------------------------------------------------------------
# 基礎繪製(builtin 頁使用;座標移植自 repo fallback/generate_review_deck.py)
# ---------------------------------------------------------------------------
def _tight_margins(tf):
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)


def _textbox(slide, text, x, y, w, h, size, bold=False, color=DARK, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    _tight_margins(tf)
    for i, line in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if align is not None:
            p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = tk.FONT_ZH
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def _rect(slide, x, y, w, h, fill=WHITE, line=LINE, rounded=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(0.8)
    shp.shadow.inherit = False
    return shp


def _set_shape_text(shp, text, size, bold=False, color=DARK, align=PP_ALIGN.CENTER):
    tf = shp.text_frame
    tf.word_wrap = True
    _tight_margins(tf)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, line in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = tk.FONT_ZH
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


def _label(slide, text, x, y, w, h, fill=GREEN, size=15):
    box = _rect(slide, x, y, w, h, fill=fill, line=fill, rounded=False)
    _set_shape_text(box, text, size=size, bold=True, color=WHITE)
    return box


def _card(slide, x, y, w, h, title, body="", title_color=DARK):
    _rect(slide, x, y, w, h)
    if title:
        _textbox(slide, title, x + 0.16, y + 0.10, w - 0.32, 0.32,
                 size=16, bold=True, color=title_color)
        body_y = y + 0.48
    else:
        body_y = y + 0.12
    if body:
        _textbox(slide, body, x + 0.16, body_y, w - 0.32, h - (body_y - y) - 0.10,
                 size=13, color=DARK)


def _bullets(lines):
    return "\n".join("• " + str(x) for x in lines)


def _bg(slide, prs, path: Path):
    slide.shapes.add_picture(str(path), 0, 0, width=prs.slide_width, height=prs.slide_height)


def _logo(slide, asset_dir: Path, rel: str):
    p = asset_dir / rel
    if p.exists():
        slide.shapes.add_picture(str(p), Inches(0.32), Inches(6.92),
                                 width=Inches(1.34), height=Inches(0.34))


def _page_number(slide, number: int):
    _textbox(slide, str(number), 12.30, 6.72, 0.70, 0.50, size=28, bold=True,
             color=DARK, align=PP_ALIGN.RIGHT)


def _header(slide, spec_slide):
    _textbox(slide, spec_slide["title"], 0.32, 0.24, 11.5, 0.48, size=32, bold=True)
    sub = spec_slide["slots"].get("subtitle")
    if sub:
        _textbox(slide, sub, 0.34, 0.84, 11.8, 0.34, size=16, color=MUTED)


def _em(s: str) -> float:
    return sum(1.0 if ord(c) > 0x2E80 else 0.55 for c in str(s))


def _finalize_page_number(slide, spec_slide):
    """清掉模板殘留的右下角純數字頁碼,依 spec 決定補標準頁碼。"""
    for s in list(tt.iter_text_shapes(slide.shapes)):
        txt = tt.shape_text(s).strip()
        if (txt.isdigit() and len(txt) <= 3 and s.top and s.left
                and s.top > Inches(6.3) and s.left > Inches(11.2)):
            s._element.getparent().remove(s._element)
    if spec_slide.get("render_page_number"):
        _page_number(slide, spec_slide["number"])


# ---------------------------------------------------------------------------
# builtin 頁型
# ---------------------------------------------------------------------------
def build_cover(prs, spec_slide, asset_dir):
    slide = tk.add_blank_slide(prs)
    slots = spec_slide["slots"]
    _bg(slide, prs, asset_dir / spec_slide["assets"]["background"])
    title, sub = slots["main_title"], slots["subtitle"]
    title_w = _em(title) * 40 / 72 + 0.25
    sub_w = _em(sub) * 36 / 72 + 0.25
    _textbox(slide, title, 0.34, 2.85, title_w, 0.70, size=40, bold=True)
    if 0.34 + title_w + 0.15 + sub_w <= 12.6:
        _textbox(slide, sub, 0.34 + title_w + 0.15, 2.92, sub_w, 0.62,
                 size=36, bold=True, color=GREEN)
        meta_y = 3.75
    else:
        _textbox(slide, sub, 0.34, 3.60, sub_w, 0.62, size=36, bold=True, color=GREEN)
        meta_y = 4.42
    _textbox(slide, f"{slots['date']}  |  {slots['presenters']}",
             0.34, meta_y, 8.0, 0.40, size=18, bold=True, color=MUTED)
    _logo(slide, asset_dir, spec_slide["assets"].get("logo", "assets/logos/cathay_logo.png"))
    return slide


def build_agenda(prs, spec_slide, asset_dir):
    slide = tk.add_blank_slide(prs)
    slots = spec_slide["slots"]
    _bg(slide, prs, asset_dir / spec_slide["assets"]["background"])
    _textbox(slide, "Contents", 0.62, 1.05, 2.8, 0.34, size=14, color=MUTED)
    _textbox(slide, "目錄", 0.62, 1.37, 2.1, 0.62, size=32, bold=True)
    items = slots["items"]
    spacing = min(1.10, (5.90 - 1.42) / max(len(items) - 1, 1))
    for i, item in enumerate(items):
        y = 1.42 + i * spacing
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4.2), Inches(y),
                                      Inches(0.36), Inches(0.36))
        circ.fill.background()
        circ.line.color.rgb = MUTED
        circ.line.width = Pt(0.8)
        _set_shape_text(circ, item.get("number", str(i + 1)), size=12, color=MUTED)
        _textbox(slide, item["title"], 4.78, y - 0.02, 5.6, 0.36, size=20, bold=True)
        sub = item.get("subtitle")
        if sub and sub != "無":
            _textbox(slide, sub, 4.79, y + 0.38, 6.8, 0.32, size=14, color=MUTED)
    _logo(slide, asset_dir, spec_slide["assets"].get("logo", "assets/logos/cathay_logo.png"))
    _page_number(slide, spec_slide["number"])
    return slide


def build_closing(prs, spec_slide, asset_dir):
    slide = tk.add_blank_slide(prs)
    _bg(slide, prs, asset_dir / spec_slide["assets"]["background"])
    _textbox(slide, spec_slide["slots"].get("main_title", "Thank you"),
             0.86, 3.28, 5.0, 0.60, size=32, bold=True)
    return slide


def build_story_chapter_statement(prs, spec_slide, asset_dir):
    slide = tk.add_blank_slide(prs)
    slots = spec_slide["slots"]
    _bg(slide, prs, asset_dir / spec_slide["assets"]["background"])
    _header(slide, spec_slide)
    colors = [PURPLE, GREEN, DARK]
    for i, phase in enumerate(slots["story"][:3]):
        x = 0.60 + i * 3.95
        _label(slide, phase["phase"], x, 1.50, 3.45, 0.42, fill=colors[i], size=18)
        _card(slide, x, 2.10, 3.45, 1.60, "", phase["text"])
    _card(slide, 0.60, 4.05, 5.65, 2.20, "背景",
          _bullets(slots["background_points"]), title_color=DARK)
    _card(slide, 6.40, 4.05, 5.65, 2.20, "預期成果",
          _bullets(slots["outcomes"]), title_color=GREEN)
    _logo(slide, asset_dir, spec_slide["assets"].get("logo", "assets/logos/cathay_logo.png"))
    _page_number(slide, spec_slide["number"])
    return slide


def build_stage_dual_track_roadmap(prs, spec_slide, asset_dir):
    slide = tk.add_blank_slide(prs)
    slots = spec_slide["slots"]
    _bg(slide, prs, asset_dir / spec_slide["assets"]["background"])
    _header(slide, spec_slide)
    x_start, col_w = 1.72, 2.72
    for i, q in enumerate(slots["quarters"][:4]):
        _label(slide, q, x_start + i * col_w, 1.42, 2.45, 0.50,
               fill=GREEN if i < 3 else DARK, size=15)
    lane_ys = [2.14, 3.68]
    for lane, y in zip(slots["lanes"][:2], lane_ys):
        _textbox(slide, lane["name"], 0.34, y + 0.30, 1.30, 0.90,
                 size=13, bold=True, color=MUTED)
        for i, cell in enumerate(lane["cells"][:4]):
            _card(slide, x_start + i * col_w, y, 2.45, 1.32, "", cell)
    cycle = slots["annual_cycle"]
    _label(slide, "年度循環", 0.34, 5.48, 1.30, 0.42, fill=DARK, size=14)
    total_w = 12.10 - x_start
    cell_w = total_w / len(cycle) - 0.15
    for i, item in enumerate(cycle):
        cx = x_start + i * (cell_w + 0.15)
        chip = _rect(slide, cx, 5.42, cell_w, 0.56)
        _set_shape_text(chip, item, size=13, bold=True, color=GREEN)
    _logo(slide, asset_dir, spec_slide["assets"].get("logo", "assets/logos/cathay_logo.png"))
    _page_number(slide, spec_slide["number"])
    return slide


BUILDERS = {
    "cover": build_cover,
    "agenda": build_agenda,
    "closing": build_closing,
    "story_chapter_statement": build_story_chapter_statement,
    "stage_dual_track_roadmap": build_stage_dual_track_roadmap,
}


# ---------------------------------------------------------------------------
# plan 驅動的 clone 頁(未涵蓋頁型用)
# ---------------------------------------------------------------------------
def _match_shapes(slide, match: dict):
    if "id" in match or ("name" in match and "contains" not in match):
        found = []

        def walk(shapes):
            for s in shapes:
                if "id" in match and s.shape_id == match["id"]:
                    found.append(s)
                elif "name" in match and s.name == match["name"]:
                    found.append(s)
                if getattr(s, "shape_type", None) is not None and hasattr(s, "shapes"):
                    walk(s.shapes)

        walk(slide.shapes)
        return found
    return tt.find_text_shapes(slide, contains=match.get("contains"),
                               name=match.get("name"))


def apply_clone_plan(prs, spec_slide, page_plan, template_page_count, problems):
    tpl = page_plan.get("template_page")
    num = spec_slide["number"]
    if not tpl or not (1 <= tpl <= template_page_count):
        problems.append(f"p{num}: template_page 缺少或超出範圍 1–{template_page_count}")
        return None
    slide = tk.clone_slide(prs, tpl - 1)
    edited = []

    for d in page_plan.get("delete", []):
        hits = _match_shapes(slide, d.get("match", {}))
        if not hits:
            problems.append(f"p{num}: delete UNMATCHED {d.get('match')}")
            continue
        for s in hits:
            s._element.getparent().remove(s._element)

    for e in page_plan.get("edits", []):
        hits = sorted(_match_shapes(slide, e.get("match", {})),
                      key=lambda s: (s.top or 0, s.left or 0))
        nth = e.get("nth")
        if not hits:
            problems.append(f"p{num}: edit UNMATCHED {e.get('match')}")
            continue
        if len(hits) > 1 and nth is None:
            problems.append(f"p{num}: edit AMBIGUOUS x{len(hits)} {e.get('match')} → 加 nth 或改用 id")
            continue
        target = hits[nth or 0] if nth is not None else hits[0]
        if not getattr(target, "has_text_frame", False):
            problems.append(f"p{num}: edit 目標不是文字框 {e.get('match')}")
            continue
        tt.set_text_keep_style(target, e["text"])
        edited.append(target)

    if page_plan.get("shrink", True):
        for s in edited:
            tt.shrink_to_fit(s, min_pt=12)
    _finalize_page_number(slide, spec_slide)
    return slide


# ---------------------------------------------------------------------------
# 主流程
# ---------------------------------------------------------------------------
def main(argv):
    ap = argparse.ArgumentParser()
    ap.add_argument("--spec", required=True)
    ap.add_argument("--plan", help="選配:僅未涵蓋頁型需要")
    ap.add_argument("--template", required=True)
    ap.add_argument("--asset-dir", default="/mnt/data")
    ap.add_argument("--out", required=True)
    args = ap.parse_args(argv)

    # utf-8-sig:容忍 Windows 工具(記事本等)寫入的 BOM,使用者上傳的 JSON 常見
    spec = json.loads(Path(args.spec).read_text(encoding="utf-8-sig"))
    plan_by_num = {}
    if args.plan:
        plan = json.loads(Path(args.plan).read_text(encoding="utf-8-sig"))
        plan_by_num = {p["number"]: p for p in plan.get("slides", [])}
    asset_dir = Path(args.asset_dir)

    spec_by_num = {s["number"]: s for s in spec["slides"]}
    problems = []
    modes = []

    prs = Presentation(args.template)
    n_template = len(prs.slides)

    for num in sorted(spec_by_num):
        spec_slide = spec_by_num[num]
        pt = spec_slide["page_type"]
        page_plan = plan_by_num.get(num)

        if page_plan:  # plan 條目優先(可覆寫自動頁)
            mode = page_plan.get("mode")
            if mode == "builtin":
                if pt not in BUILDERS:
                    problems.append(f"p{num}: builtin 不支援頁型 {pt}")
                    continue
                BUILDERS[pt](prs, spec_slide, asset_dir)
                modes.append(f"p{num}:builtin")
            elif mode == "clone":
                apply_clone_plan(prs, spec_slide, page_plan, n_template, problems)
                modes.append(f"p{num}:clone{page_plan.get('template_page')}")
            else:
                problems.append(f"p{num}: 未知 mode {mode!r}(builtin|clone)")
        elif pt in BUILDERS:
            BUILDERS[pt](prs, spec_slide, asset_dir)
            modes.append(f"p{num}:builtin")
        elif pt in fills.FILLS:
            tpl_page, fn = fills.FILLS[pt]
            slide = tk.clone_slide(prs, tpl_page - 1)
            ctx = fills.Ctx(slide)
            try:
                fn(ctx, spec_slide)
                ctx.shrink_edited(min_pt=12)
            except fills.FillError as e:
                problems.append(f"p{num}: {e}")
            _finalize_page_number(slide, spec_slide)
            modes.append(f"p{num}:auto{tpl_page}")
        else:
            problems.append(
                f"p{num}: 頁型 {pt!r} 無自動支援 → 需要 plan 條目"
                f"(mode=clone + template_page + edits,參考 page_types.md)")

    tk.delete_slides(prs, range(n_template))
    sections_removed = tk.clear_sections(prs)
    prs.save(args.out)

    print(f"產出:{args.out}  頁數:{len(spec_by_num)}  清除Section:{sections_removed}")
    print("模式:" + " ".join(modes))
    if problems:
        print(f"✗ 問題 {len(problems)} 項(修 spec/plan 後整檔重跑,勿手改 pptx):")
        for p in problems:
            print(f"   {p}")
        return 1
    print("結果:OK — 接著跑 qa_check.py")
    return 0


if __name__ == "__main__":
    sys.exit(main(sys.argv[1:]))
