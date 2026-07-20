#!/usr/bin/env python3
"""qa_check — 產檔後自檢:比對產出 pptx 與 slide_spec.json。

驗證器(validate_slide_spec_gpts.py)擋「產檔前的 spec」;本工具擋「產檔後的
pptx」。兩道都過,才交付給使用者。

用法:
  python qa_check.py --spec /mnt/data/slide_spec.json --pptx /mnt/data/deck.pptx

檢查項目:
  FAIL(必修,exit 1):
    - 頁數與 spec / deck.slide_count 不符
    - spec 槽位文字沒有出現在對應頁(內容遺漏或沒替換到)
    - 殘留 PowerPoint Section
  WARN(人工判斷):
    - 字體非 Microsoft JhengHei / Helvetica 的 run
    - 頁碼缺漏或不該有頁碼的頁出現頁碼
    - 文字溢出疑慮(CJK 啟發式估算,列最嚴重前 5)

輸出精簡:只印問題;全部通過就一行 PASS。修法一律是改 render_plan.json /
slide_spec.json 後重跑 render_deck(整檔重生),不要手改 pptx。
"""
from __future__ import annotations

import argparse
import json
import sys
import unicodedata
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches

import text_tools as tt

# 微軟正黑體 = JhengHei 中文名;Noto Sans TC 為模板原生設計(模板優先級最高);
# '+' 開頭為佈景主題字型參照(+mn-ea 等),由模板主題解析,一律放行。
ALLOWED_FONTS = {"Microsoft JhengHei", "微軟正黑體", "Helvetica", "Noto Sans TC"}


def _font_ok(name: str) -> bool:
    return name in ALLOWED_FONTS or name.startswith("+")


def norm(s: str) -> str:
    s = unicodedata.normalize("NFKC", str(s))
    return "".join(ch for ch in s if not ch.isspace())


def collect_strings(val, out):
    """遞迴收集 slots 內所有字串(檢查它們必須出現在產出頁上)。"""
    if isinstance(val, str):
        if val.strip():
            out.append(val)
    elif isinstance(val, list):
        for v in val:
            collect_strings(v, out)
    elif isinstance(val, dict):
        for v in val.values():
            collect_strings(v, out)


def slide_all_text(slide) -> str:
    parts = []
    for shp in tt.iter_text_shapes(slide.shapes):
        parts.append(tt.shape_text(shp))
    for shp in slide.shapes:  # 表格文字
        if getattr(shp, "has_table", False):
            for row in shp.table.rows:
                for cell in row.cells:
                    parts.append(cell.text)
    return norm("".join(parts))


def has_sections(prs) -> bool:
    for extLst in prs._element.findall(qn("p:extLst")):
        for ext in extLst.iter():
            if ext.tag.endswith("}sectionLst"):
                return True
    return False


def main(argv):
    ap = argparse.ArgumentParser()
    ap.add_argument("--spec", required=True)
    ap.add_argument("--pptx", required=True)
    args = ap.parse_args(argv)

    spec = json.loads(Path(args.spec).read_text(encoding="utf-8-sig"))  # 容忍 BOM
    prs = Presentation(args.pptx)
    fails, warns = [], []

    spec_slides = sorted(spec["slides"], key=lambda s: s["number"])
    declared = spec.get("deck", {}).get("slide_count")

    # 頁數
    if len(prs.slides) != len(spec_slides):
        fails.append(f"頁數不符:pptx {len(prs.slides)} 頁,spec {len(spec_slides)} 頁")
    if declared is not None and declared != len(prs.slides):
        fails.append(f"deck.slide_count={declared} 與 pptx {len(prs.slides)} 頁不符")

    # Section
    if has_sections(prs):
        fails.append("殘留 PowerPoint Section 分組(style_guide 禁止)")

    # 逐頁:槽位文字覆蓋 / 頁碼 / 字體 / 溢出
    overflow_all = []
    for i, spec_slide in enumerate(spec_slides):
        if i >= len(prs.slides):
            break
        slide = prs.slides[i]
        num = spec_slide["number"]
        page_text = slide_all_text(slide)

        wanted = []
        collect_strings(spec_slide.get("slots", {}), wanted)
        missing = [w for w in wanted if norm(w) not in page_text]
        for m in missing:
            fails.append(f"p{num}: 內容未出現「{m[:30]}{'…' if len(m) > 30 else ''}」")

        digit_shapes = [s for s in tt.iter_text_shapes(slide.shapes)
                        if tt.shape_text(s).strip().isdigit()
                        and len(tt.shape_text(s).strip()) <= 3
                        and (s.top or 0) > Inches(6.3) and (s.left or 0) > Inches(11.0)]
        if spec_slide.get("render_page_number"):
            if not any(tt.shape_text(s).strip() == str(num) for s in digit_shapes):
                warns.append(f"p{num}: 找不到右下角頁碼 {num}")
        elif digit_shapes:
            warns.append(f"p{num}: 此頁型不應有頁碼,但右下角有數字框")

        for shp in tt.iter_text_shapes(slide.shapes):
            for para in shp.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.name and not _font_ok(run.font.name):
                        warns.append(f"p{num}: 字體 {run.font.name!r}(shape id={shp.shape_id})")
                        break

        for shp in tt.iter_text_shapes(slide.shapes):
            if not tt.shape_text(shp).strip():
                continue
            est = tt.estimate_overflow(shp)
            if not est["fits"]:
                overflow_all.append((est["needed_pt"] / max(est["avail_pt"], 1), num,
                                     shp.shape_id, tt.shape_text(shp)[:20]))

    warns = list(dict.fromkeys(warns))  # 去重(字體警告易重複)
    for ratio, num, sid, txt in sorted(overflow_all, reverse=True)[:5]:
        warns.append(f"p{num}: 溢出疑慮 x{ratio:.1f}(id={sid}「{txt}…」)— 沙箱估算,請開檔確認")

    if fails:
        print(f"✗ FAIL ({len(fails)})")
        for f in fails:
            print(f"   [F] {f}")
    if warns:
        print(f"⚠ WARN ({len(warns)})")
        for w in warns:
            print(f"   [W] {w}")
    if fails:
        print("結果:FAIL — 修 render_plan.json / slide_spec.json 後重跑 render_deck,勿手改 pptx")
        return 1
    print(f"結果:PASS({len(prs.slides)} 頁{',' + str(len(warns)) + ' 個警告' if warns else ''})— 可交付")
    return 0


if __name__ == "__main__":
    sys.exit(main(sys.argv[1:]))
