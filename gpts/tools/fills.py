#!/usr/bin/env python3
"""fills — 已註冊頁型的確定性填充引擎(spec → 模板頁,零 LLM 參與)。

每個函式對應一種註冊頁型:render_deck 先 clone 對應模板頁,再呼叫這裡的
fill_* 把 spec 槽位填進**寫死的 shape id**。id 來自 inspect_template.py 對
light_template.pptx 的盤點(2026-07-20);模板改版時必須重新盤點並更新本檔,
這是本檔與模板之間的硬耦合,詳見 WORKLOG。

涵蓋(FILLS 註冊表):
  vision_goal_center_balance       → 模板 p14
  info_three_column_category       → 模板 p17
  data_two_group_metric_comparison → 模板 p29
  evaluation_option_score_pros_cons→ 模板 p33
  pyramid_layered_maturity_detail  → 模板 p54
其餘註冊頁型(cover/agenda/closing/story/stage)走 render_deck 的 builtin。

已知取捨(修改前先讀):
- evaluation:模板的分數籤/縮圖與契約無此欄位 → 一律刪除;只有 2 個方案時
  刪第三欄、不重新置中(待改進)。recommended/recommendation 以底部一行呈現。
- pyramid:模板固定 5 層,4 層時刪「頂端」一層;模板無副標佔位 → 動態補建。
- 清單超出模板格位時,溢出項以換行併入最後一格,再靠縮字消化。
"""
from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

import text_tools as tt


class FillError(Exception):
    pass


# ---------------------------------------------------------------------------
# 共用小工具
# ---------------------------------------------------------------------------
def index_shapes(slide) -> dict:
    """shape_id → shape(含群組內層)。clone 保留原模板 id,故可寫死引用。"""
    idx = {}

    def walk(shapes):
        for s in shapes:
            idx[s.shape_id] = s
            if s.shape_type == MSO_SHAPE_TYPE.GROUP:
                walk(s.shapes)

    walk(slide.shapes)
    return idx


class Ctx:
    """單頁填充的工作環境:查 id、換字、刪形狀,並記錄動過的框供縮字。"""

    def __init__(self, slide):
        self.slide = slide
        self.idx = index_shapes(slide)
        self.edited = []

    def shape(self, sid):
        if sid not in self.idx:
            raise FillError(f"shape id {sid} 不存在(模板改版?請重新盤點並更新 fills.py)")
        return self.idx[sid]

    def set(self, sid, text):
        shp = self.shape(sid)
        tt.set_text_keep_style(shp, text)
        self.edited.append(shp)

    def delete(self, *sids):
        for sid in sids:
            shp = self.idx.pop(sid, None)
            if shp is not None:
                el = shp._element
                parent = el.getparent()
                if parent is not None:
                    parent.remove(el)

    def shrink_edited(self, min_pt=12.0):
        for shp in self.edited:
            try:
                tt.shrink_to_fit(shp, min_pt=min_pt)
            except Exception:
                pass  # 縮字失敗不致命,交給 qa_check 的溢出警告


def _fill_rows(ctx, points, row_text_ids, sep_ids):
    """把 points 填進固定列(不足刪列與分隔線,超出併入最後一列)。"""
    n_rows = len(row_text_ids)
    for i, sid in enumerate(row_text_ids):
        if i < len(points):
            text = points[i]
            if i == n_rows - 1 and len(points) > n_rows:  # 溢出併入最後一列
                text = "\n".join(points[i:])
                ctx.shape(sid).height = Inches(0.80)      # 併入多列時加高
            ctx.set(sid, text)
        else:
            ctx.delete(sid)
            if i - 1 < len(sep_ids):
                ctx.delete(sep_ids[i - 1])


# ---------------------------------------------------------------------------
# vision_goal_center_balance — 模板 p14
# ---------------------------------------------------------------------------
def fill_vision_goal_center_balance(ctx, spec):
    slots = spec["slots"]
    ctx.set(51, spec["title"])
    ctx.set(52, slots["subtitle"])
    # 中心圓的兩個文字框原始高度只夠一行,加高讓長句換行後仍留在圓內
    s37 = ctx.shape(37)
    s37.top, s37.height = Inches(3.25), Inches(1.00)
    s38 = ctx.shape(38)
    s38.top, s38.height = Inches(4.32), Inches(0.78)
    ctx.set(37, slots["core_mission"])
    ctx.set(38, slots["annual_goal"])
    left = [(13, 12), (15, 14), (17, 16), (19, 18)]     # (文字框, 底圖) 上→下
    right = [(21, 20), (23, 22), (25, 24), (27, 26)]
    for i, (txt_id, bg_id) in enumerate(left):
        if i < len(slots["projects"]):
            ctx.set(txt_id, slots["projects"][i])
        else:
            ctx.delete(txt_id, bg_id)
    kpis = slots["kpis"]
    for i, (txt_id, bg_id) in enumerate(right):
        if i < len(kpis):
            ctx.set(txt_id, f'{kpis[i]["label"]} {kpis[i]["value"]}')
        else:
            ctx.delete(txt_id, bg_id)


# ---------------------------------------------------------------------------
# info_three_column_category — 模板 p17
# ---------------------------------------------------------------------------
# 每欄:標題 id、兩個段落框、(群組, 列點標題, 列點說明) 由上而下
_P17_COLS = [
    {"heading": 6, "paras": [33, 36],
     "groups": [(45, 46, 47), (41, 42, 43), (48, 49, 50)]},
    {"heading": 52, "paras": [55, 57],
     "groups": [(61, 62, 63), (58, 59, 60), (64, 65, 66)]},
    {"heading": 70, "paras": [73, 75],
     "groups": [(79, 80, 81), (76, 77, 78)]},
]


def fill_info_three_column_category(ctx, spec):
    slots = spec["slots"]
    ctx.set(2, spec["title"])
    ctx.set(4, slots["subtitle"])
    ctx.delete(7, 28, 29, 67, 85, 11, 12, 13)  # 類別直排/標籤/右上子導覽:契約無此概念
    for col_def, col in zip(_P17_COLS, slots["columns"]):
        ctx.set(col_def["heading"], col["heading"])
        points = col["points"]
        ctx.set(col_def["paras"][0], points[0])
        if len(points) >= 2:
            ctx.set(col_def["paras"][1], points[1])
        else:
            ctx.delete(col_def["paras"][1])
        rest = points[2:]
        groups = col_def["groups"]
        for gi, (grp, title_id, desc_id) in enumerate(groups):
            if gi < len(rest):
                ctx.set(title_id, rest[gi])
                if gi == len(groups) - 1 and len(rest) > len(groups):
                    ctx.set(desc_id, "\n".join(rest[len(groups):]))
                else:
                    ctx.delete(desc_id)
            else:
                ctx.delete(grp)


# ---------------------------------------------------------------------------
# data_two_group_metric_comparison — 模板 p29
# ---------------------------------------------------------------------------
def fill_data_two_group_metric_comparison(ctx, spec):
    slots = spec["slots"]
    ctx.set(13, spec["title"])
    ctx.set(25, slots["subtitle"])
    ctx.set(11, slots["before"]["heading"])
    ctx.set(50, slots["after"]["heading"])
    _fill_rows(ctx, slots["before"]["points"], [21, 22, 23], [24, 47])
    _fill_rows(ctx, slots["after"]["points"], [58, 59, 60], [61, 62])
    # KPI 群(標題, 數字, 單位),填入順序:左上、右上、左下、右下;單位框一律刪
    clusters = [(14, 15, 16), (51, 52, 53), (17, 18, 19), (54, 55, 56)]
    kpis = slots["kpis"]
    for i, (label_id, value_id, unit_id) in enumerate(clusters):
        ctx.delete(unit_id)
        if i < len(kpis):
            ctx.set(label_id, kpis[i]["label"])
            ctx.set(value_id, kpis[i]["value"])
        else:
            ctx.delete(label_id, value_id)


# ---------------------------------------------------------------------------
# evaluation_option_score_pros_cons — 模板 p33
# ---------------------------------------------------------------------------
# 每方案:(標頭群組, 名稱框, 分數籤, 縮圖, 優點群組, 優點文字, 缺點群組, 缺點文字)
_P33_OPTS = [
    {"hdr": 12, "name": 9, "score": 3, "pic": 7, "pros_g": 29, "pros": 31, "cons_g": 64, "cons": 66},
    {"hdr": 21, "name": 18, "score": 17, "pic": 15, "pros_g": 34, "pros": 36, "cons_g": 68, "cons": 70},
    {"hdr": 22, "name": 27, "score": 26, "pic": 24, "pros_g": 42, "pros": 44, "cons_g": 72, "cons": 74},
]


def fill_evaluation_option_score_pros_cons(ctx, spec):
    slots = spec["slots"]
    ctx.set(10, spec["title"])
    ctx.set(11, slots["subtitle"])
    options = slots["options"]
    for i, od in enumerate(_P33_OPTS):
        if i < len(options):
            opt = options[i]
            ctx.set(od["name"], opt["name"])
            ctx.set(od["pros"], "\n".join(opt["pros"]))
            ctx.set(od["cons"], "\n".join(opt["cons"]))
            ctx.delete(od["score"], od["pic"])  # 契約無分數/縮圖欄位,留著=捏造
        else:
            ctx.delete(od["hdr"], od["pros_g"], od["cons_g"])
    rec_parts = []
    if slots.get("recommended"):
        rec_parts.append(slots["recommended"])
    rec_parts += slots.get("recommendation", [])
    if rec_parts:
        box = ctx.slide.shapes.add_textbox(Inches(0.85), Inches(6.42),
                                           Inches(11.8), Inches(0.36))
        tf = box.text_frame
        tf.word_wrap = True
        run = tf.paragraphs[0].add_run()
        run.text = "建議:" + "、".join(rec_parts)
        run.font.name = "Microsoft JhengHei"
        run.font.size = Pt(14)
        run.font.bold = True


# ---------------------------------------------------------------------------
# pyramid_layered_maturity_detail — 模板 p54
# ---------------------------------------------------------------------------
# 5 層,由上而下:(金字塔橫帶, 帶上標籤, 列標題, 列說明)
_P54_LEVELS = [
    (27, 33, 38, 41),
    (28, 34, 43, 44),
    (30, 35, 49, 52),
    (31, 36, 56, 57),
    (32, 37, 58, 59),
]
# 右側卡:(群組, 標題, 說明框, 多餘說明框, 標籤×3)
_P54_CARDS = [
    (64, 67, 66, 103, (104, 105, 106)),
    (111, 91, 88, 107, (108, 109, 110)),
]


def fill_pyramid_layered_maturity_detail(ctx, spec):
    slots = spec["slots"]
    ctx.set(10, spec["title"])
    # 模板 p54 沒有副標佔位符 → 補建(位置對齊其他頁的副標)
    sub = ctx.slide.shapes.add_textbox(Inches(0.31), Inches(0.86),
                                       Inches(12.7), Inches(0.37))
    tf = sub.text_frame
    tf.word_wrap = True
    run = tf.paragraphs[0].add_run()
    run.text = slots["subtitle"]
    run.font.name = "Microsoft JhengHei"
    run.font.size = Pt(16)

    levels = slots["levels"]  # 4-5 層;4 層時刪模板頂層(帶最短那層)
    level_defs = _P54_LEVELS[len(_P54_LEVELS) - len(levels):]
    for bar, bar_lbl, row_t, row_d in _P54_LEVELS[:len(_P54_LEVELS) - len(levels)]:
        ctx.delete(bar, bar_lbl, row_t, row_d)
    for (bar, bar_lbl, row_t, row_d), lv in zip(level_defs, levels):
        # 帶上標籤框極小(~0.8 吋寬),放不下長 label 時整框刪除,由右列呈現全文
        if len(lv["label"]) <= 6:
            ctx.set(bar_lbl, lv["label"])
        else:
            ctx.delete(bar_lbl)
        ctx.set(row_t, lv["label"])
        ctx.set(row_d, lv["detail"])

    cards = slots.get("side_cards") or []
    for i, (grp, head_id, desc_id, extra_id, tag_ids) in enumerate(_P54_CARDS):
        if i < len(cards):
            ctx.set(head_id, cards[i]["heading"])
            ctx.set(desc_id, "\n".join(cards[i]["points"]))
            ctx.delete(extra_id, *tag_ids)
        else:
            ctx.delete(grp)
    if not cards:
        ctx.delete(63)  # 連右側底板一起刪


# ---------------------------------------------------------------------------
# 註冊表:page_type → (模板頁 1-based, 填充函式)
# ---------------------------------------------------------------------------
FILLS = {
    "vision_goal_center_balance": (14, fill_vision_goal_center_balance),
    "info_three_column_category": (17, fill_info_three_column_category),
    "data_two_group_metric_comparison": (29, fill_data_two_group_metric_comparison),
    "evaluation_option_score_pros_cons": (33, fill_evaluation_option_score_pros_cons),
    "pyramid_layered_maturity_detail": (54, fill_pyramid_layered_maturity_detail),
}
