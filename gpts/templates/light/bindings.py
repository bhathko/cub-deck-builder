#!/usr/bin/env python3
"""light 模板包綁定(grandfather 版:Python 模組,非宣告式 bindings.json)。

2026-07-25 Phase 0 自兩處**原封搬入**(見 gpts/TEMPLATE_PACKS.md §8):
  - BUILDERS 與繪製小工具 ← gpts/tools/render_deck.py(座標/品牌色寫死的
    light 專屬「從零繪製」5 頁型;新模板不得新增 builtin,一律 clone+fill)
  - FILLS 與填充函式  ← 原 gpts/tools/fills.py(共用件已抽出至
    gpts/tools/fill_helpers.py;本檔僅剩 light 頁面特定邏輯)

shape id 來自 inspect_template.py 對 light_template.pptx 的盤點(2026-07-20,
快照見同目錄 inventory.json);模板改版必重新盤點並更新本檔,流程見
gpts/templates/TEMPLATE_LIFECYCLE.md。

涵蓋(FILLS 註冊表):
  vision_goal_center_balance       → 模板 p14
  info_three_column_category       → 模板 p17
  data_two_group_metric_comparison → 模板 p29
  evaluation_option_score_pros_cons→ 模板 p33
  pyramid_layered_maturity_detail  → 模板 p54
其餘註冊頁型(cover/agenda/closing/story/stage)走 BUILDERS。

已知取捨(修改前先讀):
- evaluation:模板的分數籤/縮圖與契約無此欄位 → 一律刪除;只有 2 個方案時
  刪第三欄、不重新置中(待改進)。recommended/recommendation 以底部一行呈現。
- pyramid:模板固定 5 層,4 層時刪「頂端」一層;模板無副標佔位 → 動態補建。
- 清單超出模板格位時,溢出項以換行併入最後一格,再靠縮字消化。
"""
from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

import pptx_toolkit as tk
from fill_helpers import add_styled_textbox, fill_rows

DARK = RGBColor(*tk.COLOR_DARK)
MUTED = RGBColor(*tk.COLOR_MUTED)
GREEN = RGBColor(*tk.COLOR_GREEN)
PURPLE = RGBColor(*tk.COLOR_PURPLE)
LINE = RGBColor(*tk.COLOR_LINE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


# ---------------------------------------------------------------------------
# 基礎繪製(builtin 頁使用;座標移植自 repo fallback/generate_review_deck.py)
# ---------------------------------------------------------------------------
def _tight_margins(tf):
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)


def _textbox(slide, text, x, y, w, h, size, bold=False, color=DARK, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    _tight_margins(tf)
    for i, line in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if align is not None:
            p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = tk.FONT_ZH
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def _rect(slide, x, y, w, h, fill=WHITE, line=LINE, rounded=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(0.8)
    shp.shadow.inherit = False
    return shp


def _set_shape_text(shp, text, size, bold=False, color=DARK, align=PP_ALIGN.CENTER):
    tf = shp.text_frame
    tf.word_wrap = True
    _tight_margins(tf)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, line in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = tk.FONT_ZH
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


def _label(slide, text, x, y, w, h, fill=GREEN, size=15):
    box = _rect(slide, x, y, w, h, fill=fill, line=fill, rounded=False)
    _set_shape_text(box, text, size=size, bold=True, color=WHITE)
    return box


def _card(slide, x, y, w, h, title, body="", title_color=DARK):
    _rect(slide, x, y, w, h)
    if title:
        _textbox(slide, title, x + 0.16, y + 0.10, w - 0.32, 0.32,
                 size=16, bold=True, color=title_color)
        body_y = y + 0.48
    else:
        body_y = y + 0.12
    if body:
        _textbox(slide, body, x + 0.16, body_y, w - 0.32, h - (body_y - y) - 0.10,
                 size=13, color=DARK)


def _bullets(lines):
    return "\n".join("• " + str(x) for x in lines)


def _bg(slide, prs, path):
    slide.shapes.add_picture(str(path), 0, 0, width=prs.slide_width, height=prs.slide_height)


def _logo(slide, asset_dir, rel: str):
    p = asset_dir / rel
    if p.exists():
        slide.shapes.add_picture(str(p), Inches(0.32), Inches(6.92),
                                 width=Inches(1.34), height=Inches(0.34))


def _page_number(slide, number: int):
    _textbox(slide, str(number), 12.30, 6.72, 0.70, 0.50, size=28, bold=True,
             color=DARK, align=PP_ALIGN.RIGHT)


def _header(slide, spec_slide):
    _textbox(slide, spec_slide["title"], 0.32, 0.24, 11.5, 0.48, size=32, bold=True)
    sub = spec_slide["slots"].get("subtitle")
    if sub:
        _textbox(slide, sub, 0.34, 0.84, 11.8, 0.34, size=16, color=MUTED)


def _em(s: str) -> float:
    return sum(1.0 if ord(c) > 0x2E80 else 0.55 for c in str(s))


# ---------------------------------------------------------------------------
# builtin 頁型
# ---------------------------------------------------------------------------
def build_cover(prs, spec_slide, asset_dir):
    slide = tk.add_blank_slide(prs)
    slots = spec_slide["slots"]
    _bg(slide, prs, asset_dir / spec_slide["assets"]["background"])
    title, sub = slots["main_title"], slots["subtitle"]
    title_w = _em(title) * 40 / 72 + 0.25
    sub_w = _em(sub) * 36 / 72 + 0.25
    _textbox(slide, title, 0.34, 2.85, title_w, 0.70, size=40, bold=True)
    if 0.34 + title_w + 0.15 + sub_w <= 12.6:
        _textbox(slide, sub, 0.34 + title_w + 0.15, 2.92, sub_w, 0.62,
                 size=36, bold=True, color=GREEN)
        meta_y = 3.75
    else:
        _textbox(slide, sub, 0.34, 3.60, sub_w, 0.62, size=36, bold=True, color=GREEN)
        meta_y = 4.42
    _textbox(slide, f"{slots['date']}  |  {slots['presenters']}",
             0.34, meta_y, 8.0, 0.40, size=18, bold=True, color=MUTED)
    _logo(slide, asset_dir, spec_slide["assets"].get("logo", "assets/logos/cathay_logo.png"))
    return slide


def build_agenda(prs, spec_slide, asset_dir):
    slide = tk.add_blank_slide(prs)
    slots = spec_slide["slots"]
    _bg(slide, prs, asset_dir / spec_slide["assets"]["background"])
    _textbox(slide, "Contents", 0.62, 1.05, 2.8, 0.34, size=14, color=MUTED)
    _textbox(slide, "目錄", 0.62, 1.37, 2.1, 0.62, size=32, bold=True)
    items = slots["items"]
    spacing = min(1.10, (5.90 - 1.42) / max(len(items) - 1, 1))
    for i, item in enumerate(items):
        y = 1.42 + i * spacing
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4.2), Inches(y),
                                      Inches(0.36), Inches(0.36))
        circ.fill.background()
        circ.line.color.rgb = MUTED
        circ.line.width = Pt(0.8)
        _set_shape_text(circ, item.get("number", str(i + 1)), size=12, color=MUTED)
        _textbox(slide, item["title"], 4.78, y - 0.02, 5.6, 0.36, size=20, bold=True)
        sub = item.get("subtitle")
        if sub and sub != "無":
            _textbox(slide, sub, 4.79, y + 0.38, 6.8, 0.32, size=14, color=MUTED)
    _logo(slide, asset_dir, spec_slide["assets"].get("logo", "assets/logos/cathay_logo.png"))
    _page_number(slide, spec_slide["number"])
    return slide


def build_closing(prs, spec_slide, asset_dir):
    slide = tk.add_blank_slide(prs)
    _bg(slide, prs, asset_dir / spec_slide["assets"]["background"])
    _textbox(slide, spec_slide["slots"].get("main_title", "Thank you"),
             0.86, 3.28, 5.0, 0.60, size=32, bold=True)
    return slide


def build_story_chapter_statement(prs, spec_slide, asset_dir):
    slide = tk.add_blank_slide(prs)
    slots = spec_slide["slots"]
    _bg(slide, prs, asset_dir / spec_slide["assets"]["background"])
    _header(slide, spec_slide)
    colors = [PURPLE, GREEN, DARK]
    for i, phase in enumerate(slots["story"][:3]):
        x = 0.60 + i * 3.95
        _label(slide, phase["phase"], x, 1.50, 3.45, 0.42, fill=colors[i], size=18)
        _card(slide, x, 2.10, 3.45, 1.60, "", phase["text"])
    _card(slide, 0.60, 4.05, 5.65, 2.20, "背景",
          _bullets(slots["background_points"]), title_color=DARK)
    _card(slide, 6.40, 4.05, 5.65, 2.20, "預期成果",
          _bullets(slots["outcomes"]), title_color=GREEN)
    _logo(slide, asset_dir, spec_slide["assets"].get("logo", "assets/logos/cathay_logo.png"))
    _page_number(slide, spec_slide["number"])
    return slide


def build_stage_dual_track_roadmap(prs, spec_slide, asset_dir):
    slide = tk.add_blank_slide(prs)
    slots = spec_slide["slots"]
    _bg(slide, prs, asset_dir / spec_slide["assets"]["background"])
    _header(slide, spec_slide)
    x_start, col_w = 1.72, 2.72
    for i, q in enumerate(slots["quarters"][:4]):
        _label(slide, q, x_start + i * col_w, 1.42, 2.45, 0.50,
               fill=GREEN if i < 3 else DARK, size=15)
    lane_ys = [2.14, 3.68]
    for lane, y in zip(slots["lanes"][:2], lane_ys):
        _textbox(slide, lane["name"], 0.34, y + 0.30, 1.30, 0.90,
                 size=13, bold=True, color=MUTED)
        for i, cell in enumerate(lane["cells"][:4]):
            _card(slide, x_start + i * col_w, y, 2.45, 1.32, "", cell)
    cycle = slots["annual_cycle"]
    _label(slide, "年度循環", 0.34, 5.48, 1.30, 0.42, fill=DARK, size=14)
    total_w = 12.10 - x_start
    cell_w = total_w / len(cycle) - 0.15
    for i, item in enumerate(cycle):
        cx = x_start + i * (cell_w + 0.15)
        chip = _rect(slide, cx, 5.42, cell_w, 0.56)
        _set_shape_text(chip, item, size=13, bold=True, color=GREEN)
    _logo(slide, asset_dir, spec_slide["assets"].get("logo", "assets/logos/cathay_logo.png"))
    _page_number(slide, spec_slide["number"])
    return slide


# ---------------------------------------------------------------------------
# fills 頁型:vision_goal_center_balance — 模板 p14
# ---------------------------------------------------------------------------
def fill_vision_goal_center_balance(ctx, spec):
    slots = spec["slots"]
    ctx.set(51, spec["title"])
    ctx.set(52, slots["subtitle"])
    # 中心圓的兩個文字框原始高度只夠一行,加高讓長句換行後仍留在圓內
    s37 = ctx.shape(37)
    s37.top, s37.height = Inches(3.25), Inches(1.00)
    s38 = ctx.shape(38)
    s38.top, s38.height = Inches(4.32), Inches(0.78)
    ctx.set(37, slots["core_mission"])
    ctx.set(38, slots["annual_goal"])
    left = [(13, 12), (15, 14), (17, 16), (19, 18)]     # (文字框, 底圖) 上→下
    right = [(21, 20), (23, 22), (25, 24), (27, 26)]
    for i, (txt_id, bg_id) in enumerate(left):
        if i < len(slots["projects"]):
            ctx.set(txt_id, slots["projects"][i])
        else:
            ctx.delete(txt_id, bg_id)
    kpis = slots["kpis"]
    for i, (txt_id, bg_id) in enumerate(right):
        if i < len(kpis):
            ctx.set(txt_id, f'{kpis[i]["label"]} {kpis[i]["value"]}')
        else:
            ctx.delete(txt_id, bg_id)


# ---------------------------------------------------------------------------
# info_three_column_category — 模板 p17
# ---------------------------------------------------------------------------
# 每欄:標題 id、兩個段落框、(群組, 列點標題, 列點說明) 由上而下
_P17_COLS = [
    {"heading": 6, "paras": [33, 36],
     "groups": [(45, 46, 47), (41, 42, 43), (48, 49, 50)]},
    {"heading": 52, "paras": [55, 57],
     "groups": [(61, 62, 63), (58, 59, 60), (64, 65, 66)]},
    {"heading": 70, "paras": [73, 75],
     "groups": [(79, 80, 81), (76, 77, 78)]},
]


def fill_info_three_column_category(ctx, spec):
    slots = spec["slots"]
    ctx.set(2, spec["title"])
    ctx.set(4, slots["subtitle"])
    ctx.delete(7, 28, 29, 67, 85, 11, 12, 13)  # 類別直排/標籤/右上子導覽:契約無此概念
    for col_def, col in zip(_P17_COLS, slots["columns"]):
        ctx.set(col_def["heading"], col["heading"])
        points = col["points"]
        ctx.set(col_def["paras"][0], points[0])
        if len(points) >= 2:
            ctx.set(col_def["paras"][1], points[1])
        else:
            ctx.delete(col_def["paras"][1])
        rest = points[2:]
        groups = col_def["groups"]
        for gi, (grp, title_id, desc_id) in enumerate(groups):
            if gi < len(rest):
                ctx.set(title_id, rest[gi])
                if gi == len(groups) - 1 and len(rest) > len(groups):
                    ctx.set(desc_id, "\n".join(rest[len(groups):]))
                else:
                    ctx.delete(desc_id)
            else:
                ctx.delete(grp)


# ---------------------------------------------------------------------------
# data_two_group_metric_comparison — 模板 p29
# ---------------------------------------------------------------------------
def fill_data_two_group_metric_comparison(ctx, spec):
    slots = spec["slots"]
    ctx.set(13, spec["title"])
    ctx.set(25, slots["subtitle"])
    ctx.set(11, slots["before"]["heading"])
    ctx.set(50, slots["after"]["heading"])
    fill_rows(ctx, slots["before"]["points"], [21, 22, 23], [24, 47])
    fill_rows(ctx, slots["after"]["points"], [58, 59, 60], [61, 62])
    # KPI 群(標題, 數字, 單位),填入順序:左上、右上、左下、右下;單位框一律刪
    clusters = [(14, 15, 16), (51, 52, 53), (17, 18, 19), (54, 55, 56)]
    kpis = slots["kpis"]
    for i, (label_id, value_id, unit_id) in enumerate(clusters):
        ctx.delete(unit_id)
        if i < len(kpis):
            ctx.set(label_id, kpis[i]["label"])
            ctx.set(value_id, kpis[i]["value"])
        else:
            ctx.delete(label_id, value_id)


# ---------------------------------------------------------------------------
# evaluation_option_score_pros_cons — 模板 p33
# ---------------------------------------------------------------------------
# 每方案:(標頭群組, 名稱框, 分數籤, 縮圖, 優點群組, 優點文字, 缺點群組, 缺點文字)
_P33_OPTS = [
    {"hdr": 12, "name": 9, "score": 3, "pic": 7, "pros_g": 29, "pros": 31, "cons_g": 64, "cons": 66},
    {"hdr": 21, "name": 18, "score": 17, "pic": 15, "pros_g": 34, "pros": 36, "cons_g": 68, "cons": 70},
    {"hdr": 22, "name": 27, "score": 26, "pic": 24, "pros_g": 42, "pros": 44, "cons_g": 72, "cons": 74},
]


def fill_evaluation_option_score_pros_cons(ctx, spec):
    slots = spec["slots"]
    ctx.set(10, spec["title"])
    ctx.set(11, slots["subtitle"])
    options = slots["options"]
    for i, od in enumerate(_P33_OPTS):
        if i < len(options):
            opt = options[i]
            ctx.set(od["name"], opt["name"])
            ctx.set(od["pros"], "\n".join(opt["pros"]))
            ctx.set(od["cons"], "\n".join(opt["cons"]))
            ctx.delete(od["score"], od["pic"])  # 契約無分數/縮圖欄位,留著=捏造
        else:
            ctx.delete(od["hdr"], od["pros_g"], od["cons_g"])
    rec_parts = []
    if slots.get("recommended"):
        rec_parts.append(slots["recommended"])
    rec_parts += slots.get("recommendation", [])
    if rec_parts:
        add_styled_textbox(ctx.slide, "建議:" + "、".join(rec_parts),
                           0.85, 6.42, 11.8, 0.36,
                           font=tk.FONT_ZH, size_pt=14, bold=True)


# ---------------------------------------------------------------------------
# pyramid_layered_maturity_detail — 模板 p54
# ---------------------------------------------------------------------------
# 5 層,由上而下:(金字塔橫帶, 帶上標籤, 列標題, 列說明)
_P54_LEVELS = [
    (27, 33, 38, 41),
    (28, 34, 43, 44),
    (30, 35, 49, 52),
    (31, 36, 56, 57),
    (32, 37, 58, 59),
]
# 右側卡:(群組, 標題, 說明框, 多餘說明框, 標籤×3)
_P54_CARDS = [
    (64, 67, 66, 103, (104, 105, 106)),
    (111, 91, 88, 107, (108, 109, 110)),
]


def fill_pyramid_layered_maturity_detail(ctx, spec):
    slots = spec["slots"]
    ctx.set(10, spec["title"])
    # 模板 p54 沒有副標佔位符 → 補建(位置對齊其他頁的副標)
    add_styled_textbox(ctx.slide, slots["subtitle"], 0.31, 0.86, 12.7, 0.37,
                       font=tk.FONT_ZH, size_pt=16)

    levels = slots["levels"]  # 4-5 層;4 層時刪模板頂層(帶最短那層)
    level_defs = _P54_LEVELS[len(_P54_LEVELS) - len(levels):]
    for bar, bar_lbl, row_t, row_d in _P54_LEVELS[:len(_P54_LEVELS) - len(levels)]:
        ctx.delete(bar, bar_lbl, row_t, row_d)
    for (bar, bar_lbl, row_t, row_d), lv in zip(level_defs, levels):
        # 帶上標籤框極小(~0.8 吋寬),放不下長 label 時整框刪除,由右列呈現全文
        if len(lv["label"]) <= 6:
            ctx.set(bar_lbl, lv["label"])
        else:
            ctx.delete(bar_lbl)
        ctx.set(row_t, lv["label"])
        ctx.set(row_d, lv["detail"])

    cards = slots.get("side_cards") or []
    for i, (grp, head_id, desc_id, extra_id, tag_ids) in enumerate(_P54_CARDS):
        if i < len(cards):
            ctx.set(head_id, cards[i]["heading"])
            ctx.set(desc_id, "\n".join(cards[i]["points"]))
            ctx.delete(extra_id, *tag_ids)
        else:
            ctx.delete(grp)
    if not cards:
        ctx.delete(63)  # 連右側底板一起刪


# ---------------------------------------------------------------------------
# 註冊表(pack_loader 讀取的介面)
# ---------------------------------------------------------------------------
BUILDERS = {
    "cover": build_cover,
    "agenda": build_agenda,
    "closing": build_closing,
    "story_chapter_statement": build_story_chapter_statement,
    "stage_dual_track_roadmap": build_stage_dual_track_roadmap,
}

# page_type → (模板頁 1-based, 填充函式)
FILLS = {
    "vision_goal_center_balance": (14, fill_vision_goal_center_balance),
    "info_three_column_category": (17, fill_info_three_column_category),
    "data_two_group_metric_comparison": (29, fill_data_two_group_metric_comparison),
    "evaluation_option_score_pros_cons": (33, fill_evaluation_option_score_pros_cons),
    "pyramid_layered_maturity_detail": (54, fill_pyramid_layered_maturity_detail),
}
