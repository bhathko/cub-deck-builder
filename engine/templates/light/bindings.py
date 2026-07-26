#!/usr/bin/env python3
"""light 模板包綁定(Phase 3 起:本檔只剩 BUILDERS,fills 走宣告式 bindings.json)。

- **BUILDERS(本檔)**:cover/agenda/closing/story/stage 五種「從零繪製」頁型,
  座標與品牌色寫死,為 light 專屬 grandfather(builtin 僅 light 允許,
  新模板一律 clone+fill,見 docs/ARCHITECTURE.md §4)。
- **FILLS(同目錄 bindings.json)**:五種模板頁自動填充頁型
  (p14/17/29/33/54),經 engine/tools/fills_engine.py 解譯;Phase 2 已驗證
  與原 Python fills 產出 shape 樹全等(docs/WORKLOG.md §20.3),Phase 3 正式切換。
  pack_loader 合併語意:py 匯出非空 FILLS 才會蓋過 json——本檔刻意不匯出。

繪製座標移植自 repo fallback/generate_review_deck.py;模板改版流程見
engine/templates/TEMPLATE_LIFECYCLE.md。
"""
from __future__ import annotations

from pathlib import Path

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

import pptx_toolkit as tk
from fill_helpers import resolve_asset

_PACK_DIR = Path(__file__).resolve().parent

DARK = RGBColor(*tk.COLOR_DARK)
MUTED = RGBColor(*tk.COLOR_MUTED)
GREEN = RGBColor(*tk.COLOR_GREEN)
PURPLE = RGBColor(*tk.COLOR_PURPLE)
LINE = RGBColor(*tk.COLOR_LINE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


# ---------------------------------------------------------------------------
# 基礎繪製(builtin 頁使用)
# ---------------------------------------------------------------------------
def _tight_margins(tf):
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)


def _textbox(slide, text, x, y, w, h, size, bold=False, color=DARK, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    _tight_margins(tf)
    for i, line in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if align is not None:
            p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = tk.FONT_ZH
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def _rect(slide, x, y, w, h, fill=WHITE, line=LINE, rounded=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(0.8)
    shp.shadow.inherit = False
    return shp


def _set_shape_text(shp, text, size, bold=False, color=DARK, align=PP_ALIGN.CENTER):
    tf = shp.text_frame
    tf.word_wrap = True
    _tight_margins(tf)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, line in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = tk.FONT_ZH
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


def _label(slide, text, x, y, w, h, fill=GREEN, size=15):
    box = _rect(slide, x, y, w, h, fill=fill, line=fill, rounded=False)
    _set_shape_text(box, text, size=size, bold=True, color=WHITE)
    return box


def _card(slide, x, y, w, h, title, body="", title_color=DARK):
    _rect(slide, x, y, w, h)
    if title:
        _textbox(slide, title, x + 0.16, y + 0.10, w - 0.32, 0.32,
                 size=16, bold=True, color=title_color)
        body_y = y + 0.48
    else:
        body_y = y + 0.12
    if body:
        _textbox(slide, body, x + 0.16, body_y, w - 0.32, h - (body_y - y) - 0.10,
                 size=13, color=DARK)


def _bullets(lines):
    return "\n".join("• " + str(x) for x in lines)


def _bg(slide, prs, asset_dir, rel: str):
    # 素材解析:asset_dir 優先、包目錄兜底(light 的 asset_resolution 宣告)
    p = resolve_asset(rel, asset_dir, _PACK_DIR)
    slide.shapes.add_picture(str(p), 0, 0, width=prs.slide_width, height=prs.slide_height)


def _logo(slide, asset_dir, rel: str):
    p = resolve_asset(rel, asset_dir, _PACK_DIR)
    if p.exists():
        slide.shapes.add_picture(str(p), Inches(0.32), Inches(6.92),
                                 width=Inches(1.34), height=Inches(0.34))


def _page_number(slide, number: int):
    _textbox(slide, str(number), 12.30, 6.72, 0.70, 0.50, size=28, bold=True,
             color=DARK, align=PP_ALIGN.RIGHT)


def _header(slide, spec_slide):
    _textbox(slide, spec_slide["title"], 0.32, 0.24, 11.5, 0.48, size=32, bold=True)
    sub = spec_slide["slots"].get("subtitle")
    if sub:
        _textbox(slide, sub, 0.34, 0.84, 11.8, 0.34, size=16, color=MUTED)


def _em(s: str) -> float:
    return sum(1.0 if ord(c) > 0x2E80 else 0.55 for c in str(s))


# ---------------------------------------------------------------------------
# builtin 頁型
# ---------------------------------------------------------------------------
def build_cover(prs, spec_slide, asset_dir):
    slide = tk.add_blank_slide(prs)
    slots = spec_slide["slots"]
    _bg(slide, prs, asset_dir, spec_slide["assets"]["background"])
    title, sub = slots["main_title"], slots["subtitle"]
    title_w = _em(title) * 40 / 72 + 0.25
    sub_w = _em(sub) * 36 / 72 + 0.25
    _textbox(slide, title, 0.34, 2.85, title_w, 0.70, size=40, bold=True)
    if 0.34 + title_w + 0.15 + sub_w <= 12.6:
        _textbox(slide, sub, 0.34 + title_w + 0.15, 2.92, sub_w, 0.62,
                 size=36, bold=True, color=GREEN)
        meta_y = 3.75
    else:
        _textbox(slide, sub, 0.34, 3.60, sub_w, 0.62, size=36, bold=True, color=GREEN)
        meta_y = 4.42
    _textbox(slide, f"{slots['date']}  |  {slots['presenters']}",
             0.34, meta_y, 8.0, 0.40, size=18, bold=True, color=MUTED)
    _logo(slide, asset_dir, spec_slide["assets"].get("logo", "assets/logos/cathay_logo.png"))
    return slide


def build_agenda(prs, spec_slide, asset_dir):
    slide = tk.add_blank_slide(prs)
    slots = spec_slide["slots"]
    _bg(slide, prs, asset_dir, spec_slide["assets"]["background"])
    _textbox(slide, "Contents", 0.62, 1.05, 2.8, 0.34, size=14, color=MUTED)
    _textbox(slide, "目錄", 0.62, 1.37, 2.1, 0.62, size=32, bold=True)
    items = slots["items"]
    spacing = min(1.10, (5.90 - 1.42) / max(len(items) - 1, 1))
    for i, item in enumerate(items):
        y = 1.42 + i * spacing
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4.2), Inches(y),
                                      Inches(0.36), Inches(0.36))
        circ.fill.background()
        circ.line.color.rgb = MUTED
        circ.line.width = Pt(0.8)
        _set_shape_text(circ, item.get("number", str(i + 1)), size=12, color=MUTED)
        _textbox(slide, item["title"], 4.78, y - 0.02, 5.6, 0.36, size=20, bold=True)
        sub = item.get("subtitle")
        if sub and sub != "無":
            _textbox(slide, sub, 4.79, y + 0.38, 6.8, 0.32, size=14, color=MUTED)
    _logo(slide, asset_dir, spec_slide["assets"].get("logo", "assets/logos/cathay_logo.png"))
    _page_number(slide, spec_slide["number"])
    return slide


def build_closing(prs, spec_slide, asset_dir):
    slide = tk.add_blank_slide(prs)
    _bg(slide, prs, asset_dir, spec_slide["assets"]["background"])
    _textbox(slide, spec_slide["slots"].get("main_title", "Thank you"),
             0.86, 3.28, 5.0, 0.60, size=32, bold=True)
    return slide


def build_story_chapter_statement(prs, spec_slide, asset_dir):
    slide = tk.add_blank_slide(prs)
    slots = spec_slide["slots"]
    _bg(slide, prs, asset_dir, spec_slide["assets"]["background"])
    _header(slide, spec_slide)
    colors = [PURPLE, GREEN, DARK]
    for i, phase in enumerate(slots["story"][:3]):
        x = 0.60 + i * 3.95
        _label(slide, phase["phase"], x, 1.50, 3.45, 0.42, fill=colors[i], size=18)
        _card(slide, x, 2.10, 3.45, 1.60, "", phase["text"])
    _card(slide, 0.60, 4.05, 5.65, 2.20, "背景",
          _bullets(slots["background_points"]), title_color=DARK)
    _card(slide, 6.40, 4.05, 5.65, 2.20, "預期成果",
          _bullets(slots["outcomes"]), title_color=GREEN)
    _logo(slide, asset_dir, spec_slide["assets"].get("logo", "assets/logos/cathay_logo.png"))
    _page_number(slide, spec_slide["number"])
    return slide


def build_stage_dual_track_roadmap(prs, spec_slide, asset_dir):
    slide = tk.add_blank_slide(prs)
    slots = spec_slide["slots"]
    _bg(slide, prs, asset_dir, spec_slide["assets"]["background"])
    _header(slide, spec_slide)
    x_start, col_w = 1.72, 2.72
    for i, q in enumerate(slots["quarters"][:4]):
        _label(slide, q, x_start + i * col_w, 1.42, 2.45, 0.50,
               fill=GREEN if i < 3 else DARK, size=15)
    lane_ys = [2.14, 3.68]
    for lane, y in zip(slots["lanes"][:2], lane_ys):
        _textbox(slide, lane["name"], 0.34, y + 0.30, 1.30, 0.90,
                 size=13, bold=True, color=MUTED)
        for i, cell in enumerate(lane["cells"][:4]):
            _card(slide, x_start + i * col_w, y, 2.45, 1.32, "", cell)
    cycle = slots["annual_cycle"]
    _label(slide, "年度循環", 0.34, 5.48, 1.30, 0.42, fill=DARK, size=14)
    total_w = 12.10 - x_start
    cell_w = total_w / len(cycle) - 0.15
    for i, item in enumerate(cycle):
        cx = x_start + i * (cell_w + 0.15)
        chip = _rect(slide, cx, 5.42, cell_w, 0.56)
        _set_shape_text(chip, item, size=13, bold=True, color=GREEN)
    _logo(slide, asset_dir, spec_slide["assets"].get("logo", "assets/logos/cathay_logo.png"))
    _page_number(slide, spec_slide["number"])
    return slide


# ---------------------------------------------------------------------------
# 註冊表(pack_loader 讀取的介面;本檔刻意不匯出 FILLS → fills 由 bindings.json 生效)
# ---------------------------------------------------------------------------
BUILDERS = {
    "cover": build_cover,
    "agenda": build_agenda,
    "closing": build_closing,
    "story_chapter_statement": build_story_chapter_statement,
    "stage_dual_track_roadmap": build_stage_dual_track_roadmap,
}
