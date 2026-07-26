#!/usr/bin/env python3
"""fit_capacity — 量測模板版位真正裝得下多少內容,寫入該包的 capacity_overrides。

## 為什麼需要這支工具

契約裡的字數上限如果是憑感覺填的,結果是**閘門說 PASS、版面卻是壞的**。
2026-07-26 的實例:light 的 `core_mission` 宣告 60 字,實際版位是 2.10×0.50 吋
24pt(模板原文「核心目標」4 字),塞 60 字時渲染器靜默把字縮到 12pt——
32 頁 golden 有 99 個框被縮,所有自動檢查全綠,設計師目檢才發現。

設計原則(設計師 2026-07-26 定調):

> **字級是被設計過的。內容塞不下時正確做法是改寫更短或換頁型,不是用小一號的字體。**

所以上限必須等於「版位在**設計字級**下真正裝得下的量」,由本工具量測後寫進
manifest 的 `capacity_overrides`,讓閘門在產檔前就擋下來要求改稿。

## 四個收斂訊號

單看「文字裝不裝得進自己的框」不夠——設計師的原話是
「有時候不是 autofit 沒超過,是 autofit 後跑到了不該出現字的位置」。
所以每輪 golden 之後檢查四件事:

1. **有框被縮字** → 字級被動了(非 autofit 框)。
2. **文字佔用範圍侵入鄰欄且比模板原本更糟** → `wrap="none"` 往右長、
   autofit 往下長,壓到隔欄。
3. **autofit 框長得比「設計師自己那份原文需要的高度」還高** → 即使旁邊是
   空的、沒撞到人,版面也已經不是設計師排的樣子,所以也算壞。
4. **`add_textbox` 新增的框裝不進自己宣告的幾何** → 它沒有模板對照、沒有
   設計師基準可放寬,x/y/w/h 是綁定自己寫的,就必須自己裝得下。

四個都歸零才算收斂。**注意收斂條件的實作是「本輪沒有任何提案」**,
與「訊號全滅」不完全等價:量到的上限不比現行更緊時就不會產生提案。
所以 `fit` 回 0 之後,發版前仍要照 REGRESSION R12 獨立驗一次。

另外兩個先於主迴圈的前置修正(幾何可直接算出、不需迭代):

- **平行欄位格位數對稱**:同一契約清單綁到多組平行版位時(三欄、左右兩側),
  上限取各組格位數的**最小值**。light p17 三欄是 6/6/4 格,上限開到 5 只有
  第三欄會併格 → 三欄長得不一樣(設計師目檢 p4 抓到)。分組鍵**必須保留
  索引**,正規化掉會把三欄併成一組、檢查永遠不觸發。
- **`add_textbox` 的跨槽位預算**:prefix + 多個槽位用 join 串成一條字串放進
  一個框,限制是「總長度」。逐槽位算容量必然算錯。

## 已知的量測陷阱(全部已在本檔處理,改動前先讀)

1. 清單型槽位的契約路徑結尾是 `.item`,拿最後一段去比對槽位名會全部失敗
   → 用「最後一個非 item 段」。
2. 綁定的槽位路徑有索引與切片兩種語意:`points[0]` 指單一項目(→ `.item`)、
   `points[2:]` 指清單本身(→ 不加後綴)。搞混會把整個清單長度算到併格框上。
3. 併格框(`overflow.merge_into_id` / `rows` 的 `merge_height_in`)裝的是
   「溢出的那幾項」,不是整個清單。清單新長度 = 原長度 − n + k
   (n=原本要它裝幾項、k=實際裝得下幾項),不是 k。
4. 一個框可能承載多個欄位(`template: "{label}\\n{detail}"`),要全部收緊,
   不能只收第一個。
5. `add_textbox` 產生的框在模板裡沒有對應形狀,但一樣要納入。
6. 探測字元要依**槽位語意**選,不能看模板佔位字的字種:模板寫 `Allance`
   (半形 0.55em)就用半形算,使用者寫中文會爆掉近一倍。只有 `value`/
   `number`/`values` 這種結構上保證是數字的才用半形。
7. 群組內子形狀的 `left/top` 是群組子座標系的值,必須用 `chOff/chExt` 換算
   成絕對座標(`text_tools.walk_absolute`),否則會算出不可能的跨欄重疊。
8. 碰撞判準是「**不得比模板原本更侵入鄰欄**」,不是絕對零重疊——模板本身
   就有刻意疊在一起的設計(light p47 的徽章與副字、p29 的 44pt 數字),
   絕對零重疊會把模型誤差與設計意圖一起罰掉。
9. 收緊對象只能選「佔用範圍會隨文字長度變大」的那一方:`wrap="square"` 的
   框寬度固定,縮它的字數只會減行數,救不了橫向碰撞。
10. 鄰居的容許重疊量必須跟基準用**同一種算法**(文字範圍 vs 文字範圍)。
    拿「文字範圍 vs 框」去比會比基準嚴,永遠收斂不了——light p47 的
    name/caption 兩個框本來就在垂直方向重疊。
11. `wrap="none"` 的框行數恆等於硬換行數,所以**單行時縱向檢查完全不會擋**。
    橫向檢查(`fits_w`)不能漏,否則那種槽位等於沒量測過(旁邊剛好是空白
    就會一路放行到搜尋上界)。
12. 量 footprint 時要**傳設計字級**。預設會讀「框裡現在的字級」,而訊號①
    的定義就是「字已經被縮小了」——在最需要精準的情境下面積會偏小、
    可容字數高估近一倍。
13. **本工具只會收緊、不會放寬**(避免永久卡在「1 字」)。所以:
    設計師把版位改大之後,**必須加 `--reset` 重量**,不加的話會直接回報
    收斂、上限維持舊的緊值。
14. golden 的 qa 紅是**預期的**(溢出與碰撞正是本工具要消掉的東西),
    要求 golden exit 0 才繼續會死鎖;但 validator/render 紅必須停,
    那代表 spec 本身不合法。
15. `--reset` 會先清空再量,所以它**必須等確認這個包渲染得出來才動手**,
    否則任何一步失敗都會把包留在「閘門上限幾乎全失、status 仍 registered」
    的危險狀態。也因此 `--reset` 與 `--dry-run` 不可併用。

用法:

    python engine/release/fit_capacity.py --id light            # 量測並寫入
    python engine/release/fit_capacity.py --id light --dry-run   # 只印不寫
    python engine/release/fit_capacity.py --id light --reset     # 先清空再量
                                                                 # (改過版位大小必須加)

收斂後會**自動重生** `engine/rules/page_types_registry.md` 的容量表——
那張表是給 LLM 看的,手維護必然與 manifest 漂移(2026-07-26 漂了 18 處:
模型照表寫卻被閘門退回,或反過來過度自我審查)。

需要 python-pptx(本機沒裝:`uv run --with python-pptx python …`)。
"""
from __future__ import annotations

import argparse
import json
import re
import subprocess
import sys
from pathlib import Path

REPO = Path(__file__).resolve().parent.parent.parent
sys.path.insert(0, str(REPO / "engine" / "tools"))
sys.path.insert(0, str(REPO / "engine" / "rules"))

MIN_CHARS = 6           # 一個清單格位至少要能寫 6 個中文字,否則該格位沒有意義
COLLIDE_FLOOR = 0.03    # 平方吋:小於此視為量測誤差
COLLIDE_SLACK = 1.10    # 允許比模板基準多 10%(行距估算誤差)
MAX_ROUNDS = 8

_NUMERIC_SLOTS = {"value", "values", "number"}


# ---------------------------------------------------------------------------
# 契約路徑
# ---------------------------------------------------------------------------
def norm_slot(slot_path: str) -> str:
    """綁定的 slot 路徑 → 契約路徑。

    `$.slots.columns[0].points[0]`  → `columns.item.points.item`(索引=單一項目)
    `$.slots.columns[0].points[2:]` → `columns.item.points`     (切片=清單本身)
    """
    p = slot_path[len("$.slots."):] if slot_path.startswith("$.slots.") else slot_path
    return re.sub(r"\[([^\]]*)\]", lambda m: "" if ":" in m.group(1) else ".item", p)


def node_at(slots: dict, path: str):
    node = {"kind": "object", "fields": slots}
    for seg in path.split("."):
        if seg == "item":
            node = node.get("item", {})
        elif node.get("kind") == "object":
            node = node["fields"].get(seg, {})
        else:
            return None
        if not node:
            return None
    return node


def parent_list(slots: dict, path: str):
    """path 形如 a.item / a.b.item → 回傳 (清單路徑, 清單節點)。"""
    if not path.endswith(".item"):
        return None, None
    lp = path[: -len(".item")]
    return lp, node_at(slots, lp)


def probe_char(path: str) -> str:
    """量容量用的探測字元。預設中文(最寬)——契約只數字元數不分全形半形,
    用半形量會讓上限虛胖近一倍。數字型槽位例外。"""
    tail = [x for x in str(path).split(".") if x != "item"]
    return "0" if (tail and tail[-1] in _NUMERIC_SLOTS) else "測"


def slot_shape_map(ops: list, slots: dict) -> dict:
    """{契約路徑: [(shape_id, 這個框要裝幾項)]}。"""
    out: dict = {}

    def add(p, sid, n=1):
        out.setdefault(p, []).append((sid, n))

    for op in ops:
        kind = op["op"]
        if kind == "set" and op.get("slot", "").startswith("$.slots."):
            path = norm_slot(op["slot"])
            nd = node_at(slots, path)
            if op.get("join") and nd and nd.get("kind") == "list":
                add(path + ".item", op["id"], nd.get("max", 1))   # 整個清單併入一框
            else:
                add(path, op["id"])
        elif kind == "rows":
            base = norm_slot(op["slot"])
            rid = op["row_ids"]
            nd = node_at(slots, base)
            extra = (max(0, (nd.get("max", len(rid)) if nd else len(rid)) - len(rid))
                     if op.get("merge_height_in") else 0)
            for j, sid in enumerate(rid):
                add(base + ".item", sid, 1 + (extra if j == len(rid) - 1 else 0))
        elif kind == "list":
            base = norm_slot(op["slot"])
            nd = node_at(slots, base)
            ov = op.get("overflow") or {}
            # 切片 [2:] 表示前 2 項由別的 op 負責,溢出量要扣掉這個位移
            m = re.search(r"\[(\d+):\]", op["slot"])
            offset = int(m.group(1)) if m else 0
            extra = (max(0, (nd.get("max", 0) if nd else 0) - offset - len(op["items"]))
                     if ov else 0)
            for it in op["items"]:
                for st in it.get("sets", []):
                    s, sid = st["slot"], st["id"]
                    if st.get("template"):
                        fields = re.findall(r"\{(\w+)\}", st["template"])
                        for f in fields:
                            add(f"{base}.item.{f}", sid, len(fields))
                    elif st.get("join"):
                        sub = base + ".item" if s == "@" else base + ".item." + s[2:]
                        sn = node_at(slots, sub)
                        add(sub + ".item", sid,
                            sn.get("max", 1) if sn and sn.get("kind") == "list" else 1)
                    elif s == "@":
                        add(base + ".item", sid)
                    else:
                        add(base + ".item." + s[2:], sid)
            if ov.get("merge_into_id"):
                add(base + ".item", ov["merge_into_id"], 1 + extra)
        elif kind == "add_textbox":
            for sl in op.get("slots", []):
                path = norm_slot(sl)
                nd = node_at(slots, path)
                if nd and nd.get("kind") == "list":
                    path += ".item"
                add(path, None)          # shape_id 事前不可知,由文字反推
    return out


# ---------------------------------------------------------------------------
# 幾何判準
# ---------------------------------------------------------------------------
def _geo(slide, shape_id, tt):
    for s, L, T, W, H in tt.walk_absolute(slide.shapes):
        if getattr(s, "has_text_frame", False) and s.shape_id == shape_id:
            return s, L, T, W, H
    return None


def _neighbour_footprints(slide, shape_id, base_pairs, tt):
    """同頁其他文字的實際佔用範圍 + 各自的容許重疊量(取自模板基準)。

    必須與基準用同一種算法(範圍 vs 範圍)。拿「範圍 vs 框」比會比基準嚴,
    永遠收斂不了——light p47 的 name/caption 兩個框本來就垂直重疊。
    """
    out = []
    for s, L, T, W, H in tt.walk_absolute(slide.shapes):
        if not getattr(s, "has_text_frame", False) or not s.text_frame.text.strip():
            continue
        if s.shape_id == shape_id:
            continue
        fp = tt.text_footprint(s, L, T, W, H)
        if fp is None:
            continue
        ref = base_pairs.get(frozenset((shape_id, s.shape_id)), 0.0)
        out.append((fp, max(ref * COLLIDE_SLACK, COLLIDE_FLOOR)))
    return out


def _clean_cap(shape, geom, n, ch, size, neighbours, tt, tpl_shape=None):
    """二分搜最大的每項字數 L,使得同時成立:

    - 文字範圍不比模板基準更侵入鄰欄
    - 需要的高度不超過容許量(框高;autofit 框放寬到模板原文所需)
    """
    _, L, T, W, H = geom

    def ok(length):
        original = shape.text_frame.text
        try:
            tt.set_text_keep_style(shape, "\n".join([ch * length] * n))
            fp = tt.text_footprint(shape, L, T, W, H, size_pt=size)
            est = tt.estimate_overflow(shape, size_pt=size)
        finally:
            tt.set_text_keep_style(shape, original)
        # 縱向:行數超出容許高度
        if est["lines"] > 1:
            allow = est["avail_pt"]
            if tpl_shape is not None and tt.has_autofit(shape):
                allow = max(allow, tt.estimate_overflow(tpl_shape)["needed_pt"])
            if est["needed_pt"] > allow * COLLIDE_SLACK:
                return False
        # 橫向:wrap="none" 的框文字往旁邊長。**這條不能漏**——那種框的
        # lines 恆等於硬換行數,單行時上面的縱向檢查完全不會擋,只剩碰撞
        # 判準;旁邊剛好是空白就會一路放行到搜尋上界(等於沒量測)。
        if not est["fits_w"]:
            allow_w = est["avail_w_pt"]
            if tpl_shape is not None:
                allow_w = max(allow_w, tt.estimate_overflow(tpl_shape)["needed_w_pt"])
            if est["needed_w_pt"] > allow_w * COLLIDE_SLACK:
                return False
        if fp is None:
            return True
        return all(tt.overlap_area(fp, nb) <= lim for nb, lim in neighbours)

    if not ok(1):
        return 0
    lo, hi = 1, 300
    if ok(hi):
        return hi
    while lo + 1 < hi:
        mid = (lo + hi) // 2
        lo, hi = (mid, hi) if ok(mid) else (lo, mid)
    return lo


# ---------------------------------------------------------------------------
# 主流程
# ---------------------------------------------------------------------------
REGISTRY = REPO / "engine" / "rules" / "page_types_registry.md"
_TABLE_HEAD = "\n---\n\n## light 模板的實際容量"
_TABLE_TAIL = "\n---\n\n## 頁型選擇原則"


def sync_registry_table(manifest: dict, verbose=True) -> None:
    """把 registry 的「實際容量」表由 manifest 重新產生。

    那張表是給 LLM 看的閘門上限。手維護必然漂移——2026-07-26 漂了 18 處
    (12 個數字不符 + 6 條指向已不存在的 override),症狀是模型照表寫卻被
    閘門退回,或反過來過度自我審查。所以由 `fit` 收斂後自動重生。
    只處理 light(其他包若也要出這張表,再擴 pack id 判斷)。
    """
    if manifest.get("template_id") != "light" or not REGISTRY.exists():
        return
    from validate_slide_spec_gpts import PAGE_TYPES, apply_capacity_overrides
    merged = apply_capacity_overrides(PAGE_TYPES, manifest.get("capacity_overrides"))
    rows = []

    def walk(node, base, pt, path=""):
        if node.get("kind") == "text":
            if node["max_chars"] != base["max_chars"]:
                rows.append((pt, path, f"≤{base['max_chars']} 字", f"≤{node['max_chars']} 字"))
        elif node.get("kind") == "list":
            if node["max"] != base["max"]:
                rows.append((pt, path, f"{base['min']}–{base['max']} 項",
                             f"{node['min']}–{node['max']} 項"))
            walk(node["item"], base["item"], pt, path + "[]")
        elif node.get("kind") == "object":
            for k, v in node["fields"].items():
                walk(v, base["fields"][k], pt, f"{path}.{k}" if path else k)

    for pt in merged:
        for k, v in merged[pt]["slots"].items():
            walk(v, PAGE_TYPES[pt]["slots"][k], pt, k)
    body = [
        "", "---", "", "## light 模板的實際容量(**以本表為準**)", "",
        "上方各節的字數是「語意契約的預設值」。**實際可寫多少,由選定的模板包決定**——",
        "版位大小是設計師定的,字級也是設計過的,塞不下時請**改寫更短或換頁型**,",
        "系統不會偷偷縮小字級來遷就,也不會讓文字疊到隔欄。下表列出 light 包與預設值",
        "不同的槽位;沒列出的沿用上方數字。閘門依本表擋,寫超過會被退回。", "",
        "| 頁型 | 槽位 | 預設 | light 實際 |", "| --- | --- | --- | --- |",
    ]
    for r in sorted(rows):
        body.append(f"| {r[0]} | `{r[1]}` | {r[2]} | **{r[3]}** |")
    body += ["", f"(共 {len(rows)} 個槽位。**本表由 `template_admin.py fit` 自動重生,",
             "不要手改**——手維護必然與 manifest 漂移。量測判準:設計字級下不縮字、",
             "文字不比模板原本更侵入鄰欄、平行欄位格位數一致。)", ""]
    text = REGISTRY.read_text(encoding="utf-8")
    if _TABLE_HEAD in text and _TABLE_TAIL in text:
        head = text[: text.index(_TABLE_HEAD)]
        tail = text[text.index(_TABLE_TAIL):]
    else:
        head, tail = text.rstrip() , ""
    REGISTRY.write_text(head + "\n".join(body) + tail, encoding="utf-8")
    if verbose:
        print(f"registry 容量表已重生:{len(rows)} 列")


def run(pack_id: str, packs_root: Path, dry_run=False, reset=False, verbose=True) -> int:
    from pptx import Presentation
    import text_tools as tt
    from validate_slide_spec_gpts import PAGE_TYPES, apply_capacity_overrides

    pack_dir = Path(pack_id) if Path(pack_id).is_dir() else packs_root / pack_id
    mpath = pack_dir / "manifest.json"
    bpath = pack_dir / "bindings.json"
    if not mpath.exists():
        print(f"✗ 找不到 manifest:{mpath}")
        return 2
    manifest = json.loads(mpath.read_text(encoding="utf-8-sig"))
    if reset and dry_run:
        # --reset 會清空並寫回 manifest,--dry-run 又不寫入量測值 → 併用等於
        # 「把所有閘門上限無聲蒸發」,而 lint/golden 事後仍會全綠。直接擋掉。
        print("✗ --reset 與 --dry-run 不可併用(會清空上限卻不寫回新值)")
        return 2
    if reset:
        # 先確認這個包渲染得出來,再清空——否則任何一步失敗就把包留在
        # 「閘門上限幾乎全失、status 仍是 registered」的危險狀態。
        probe = subprocess.run(
            [sys.executable, str(REPO / "engine/release/template_admin.py"),
             "golden", "--id", pack_id]
            + (["--packs-root", str(packs_root)]
               if packs_root and Path(packs_root) != REPO / "engine" / "templates" else []),
            capture_output=True, text=True, cwd=str(REPO))
        if "失敗於「validator」" in probe.stdout or "失敗於「render」" in probe.stdout:
            print("✗ 這個包目前連渲染都過不了,不清空上限(先修好綁定再來):")
            print(probe.stdout[-1200:])
            return 1
        backup = dict(manifest.get("capacity_overrides") or {})
        manifest["capacity_overrides"] = {}
        mpath.write_text(json.dumps(manifest, ensure_ascii=False, indent=2) + "\n",
                         encoding="utf-8")
        if verbose:
            print(f"已清空 capacity_overrides({len(backup)} 條),從零量測")

    fills = json.loads(bpath.read_text(encoding="utf-8-sig")).get("fills", {})
    page_of = {pt: e["template_page"] for pt, e in manifest["page_types"].items()
               if e.get("mode") == "fill"}
    tpl = Presentation(str(pack_dir / manifest.get("template_file", "template.pptx")))

    # 模板側:每頁的形狀字級、與自身的碰撞基準
    tpl_shapes, tpl_base = {}, {}
    for pt, pg in page_of.items():
        if pg > len(tpl.slides):
            continue
        sl = tpl.slides[pg - 1]
        tpl_shapes[pt] = {s.shape_id: s for s in tt.iter_text_shapes(sl.shapes)}
        tpl_base[pt] = {frozenset((a.shape_id, b.shape_id)): area
                        for area, a, b in tt.text_collisions(sl)}

    # shape_id → 契約路徑(一個框可能對多個欄位)
    sid2paths = {}
    for pt, ent in fills.items():
        table: dict = {}
        for path, pairs in slot_shape_map(ent["ops"], PAGE_TYPES[pt]["slots"]).items():
            for sid, _ in pairs:
                if sid is not None:
                    table.setdefault(sid, []).append(path)
        sid2paths[pt] = table
    # add_textbox 的槽位:靠文字內容反推(變體文字帶槽位名)
    addbox = {}
    for pt, ent in fills.items():
        for op in ent["ops"]:
            if op["op"] == "add_textbox":
                for sl in op.get("slots", []):
                    p = norm_slot(sl)
                    nd = node_at(PAGE_TYPES[pt]["slots"], p)
                    if nd and nd.get("kind") == "list":
                        p += ".item"
                    addbox.setdefault(pt, []).append(p)

    # ---- 平行欄位的格位數對稱 ----
    # 同一個契約清單常被綁到多組平行版位(三欄、左右兩側),各組的實體格位數
    # 可能不一樣:light p17 三欄分別有 6/6/4 格。上限若開到 5,只有第三欄會
    # 觸發併格 → 三欄長得不一樣(設計師目檢 p4 抓到的就是這個)。
    # 這個不變式**可以從綁定算出來**,所以由工具守,不靠文件提醒:
    # 上限 = 各平行組格位數的最小值。
    sym = {}
    for pt, ent in fills.items():
        slots = PAGE_TYPES[pt]["slots"]
        counts = {}      # 契約清單路徑 → {平行組 key: 格位數}
        for op in ent["ops"]:
            raw = op.get("slot", "")
            if not raw.startswith("$.slots."):
                continue
            path = norm_slot(raw)
            if op["op"] == "set":
                nd = node_at(slots, path)
                if nd and nd.get("kind") == "text" and path.endswith(".item"):
                    lp = path[: -len(".item")]
                    # 分組鍵要**保留索引**——索引就是「第幾組平行版位」的身分。
                    # 正規化掉會把三欄併成一組,對稱檢查永遠不觸發。
                    key = raw.rsplit("[", 1)[0]          # columns[0].points
                    counts.setdefault(lp, {}).setdefault(key, 0)
                    counts[lp][key] += 1
            elif op["op"] == "list":
                nd = node_at(slots, path)
                if nd and nd.get("kind") == "list":
                    key = re.sub(r"\[[^\]]*:\]$", "", raw)   # 去掉尾端切片,留索引
                    counts.setdefault(path, {}).setdefault(key, 0)
                    counts[path][key] += len(op["items"])
        for lp, groups in counts.items():
            if len(groups) < 2:
                continue        # 只有一組 → 沒有對稱問題
            sym[f"{pt}.slots.{lp}.max"] = min(groups.values())
    if sym:
        cur = dict(manifest.get("capacity_overrides") or {})
        wrote = []
        for key, n in sym.items():
            pt = key.split(".slots.")[0]
            path = key[len(f"{pt}.slots."):-len(".max")]
            nd = node_at(PAGE_TYPES[pt]["slots"], path)
            if not (nd and nd.get("kind") == "list"):
                continue
            n = max(n, nd["min"])
            if cur.get(key, nd["max"]) > n:
                cur[key] = n
                wrote.append((path, n))
        if wrote and not dry_run:
            manifest["capacity_overrides"] = dict(sorted(cur.items()))
            mpath.write_text(json.dumps(manifest, ensure_ascii=False, indent=2) + "\n",
                             encoding="utf-8")
        if wrote and verbose:
            print("平行欄位格位數對稱:" + ", ".join(f"{p}→{v} 項" for p, v in wrote))

    # ---- add_textbox 的跨槽位預算 ----
    # 這種 op 把 prefix + 多個槽位用 join 串成**一條字串**放進一個自己宣告
    # 幾何的框(x/y/w/h/pt)。逐槽位算容量必然算錯:限制是「總長度」。
    # 幾何不依賴渲染結果,所以在迴圈外一次算完。
    box_budget = {}
    for pt, ent in fills.items():
        for op in ent["ops"]:
            if op["op"] != "add_textbox":
                continue
            size = float(op.get("pt", 14))
            # 一行放得下幾個中文字(扣左右內邊界各 0.1 吋)
            total = int(max((float(op["w"]) - 0.2) * 72 / size, 1))
            units, paths = 0, []
            for sl in op.get("slots", []):
                path = norm_slot(sl)
                nd = node_at(PAGE_TYPES[pt]["slots"], path)
                if nd and nd.get("kind") == "list":
                    path += ".item"
                    nd = nd.get("item")
                    units += node_at(PAGE_TYPES[pt]["slots"], path[:-5]).get("max", 1)
                else:
                    units += 1
                paths.append(path)
            overhead = len(op.get("prefix", "")) + len(op.get("join", "")) * max(units - 1, 0)
            per = max((total - overhead) // max(units, 1), 1)
            for path in paths:
                box_budget[f"{pt}.slots.{path}.max_chars"] = per
    if box_budget:
        cur = dict(manifest.get("capacity_overrides") or {})
        wrote = []
        for key, per in box_budget.items():
            pt = key.split(".slots.")[0]
            path = key[len(f"{pt}.slots."):-len(".max_chars")]
            nd = node_at(PAGE_TYPES[pt]["slots"], path)
            if nd and nd.get("kind") == "text" and cur.get(key, nd["max_chars"]) > per:
                cur[key] = per
                wrote.append((path, per))
        if wrote and not dry_run:
            manifest["capacity_overrides"] = dict(sorted(cur.items()))
            mpath.write_text(json.dumps(manifest, ensure_ascii=False, indent=2) + "\n",
                             encoding="utf-8")
        if wrote and verbose:
            print("add_textbox 跨槽位預算:"
                  + ", ".join(f"{p}→{v}" for p, v in wrote))

    order = [pt for pt in manifest["page_types"] if pt in fills]
    seq = [pt for pt in order for _ in ("min", "max")]

    for rnd in range(1, MAX_ROUNDS + 1):
        cmd = [sys.executable, str(REPO / "engine/release/template_admin.py"),
               "golden", "--id", pack_id]
        if packs_root and Path(packs_root) != REPO / "engine" / "templates":
            cmd += ["--packs-root", str(packs_root)]   # 否則子行程會去量預設包
        r = subprocess.run(cmd, capture_output=True, text=True, cwd=str(REPO))
        # 量測階段只需要「渲染出來的檔案」。qa 紅是**預期的**——溢出與碰撞
        # 正是本工具要消掉的東西,要求 qa 先綠就成了雞生蛋。
        # 但 validator 紅代表 spec 本身不合法(容量互相矛盾),那必須停。
        if "golden 失敗於「validator」" in r.stdout or "golden 失敗於「render」" in r.stdout:
            print(f"✗ 第 {rnd} 輪連渲染都過不了,無法量測:")
            print(r.stdout[-1500:])
            return 1
        made = REPO / "ppt_out" / f"golden_{pack_id}.pptx"
        if not made.exists():
            print(f"✗ 找不到驗收檔 {made}:")
            print(r.stdout[-1200:])
            return 1
        out = Presentation(str(made))
        overrides = dict(manifest.get("capacity_overrides") or {})
        live = apply_capacity_overrides(PAGE_TYPES, overrides)
        changes, unmapped, stuck = [], [], []

        for idx, slide in enumerate(out.slides):
            if idx >= len(seq):
                break
            pt = seq[idx]
            # 用**已套 override** 的契約讀現值。拿基準契約的 max 去算清單長度
            # 公式(原長 − n + k)會算出比現行還鬆的目標 → 永遠縮不下去 →
            # 直接掉進 stuck 並誤報「該版位放不了有意義的內容」。
            slots = live[pt]["slots"]
            base = tpl_base.get(pt, {})
            suspects = set()

            # 訊號 ①:被縮字(字級是設計過的,不該被動)
            # 訊號 ③:autofit 框長得比「設計師自己那份原文需要的高度」還高。
            #   即使旁邊是空的、沒撞到人,版面也已經不是設計師排的樣子。
            #   容許量取 max(渲染後框高, 模板原文所需)——框可能被綁定的
            #   resize 刻意加高,那才是最終尺寸。
            for s in tt.iter_text_shapes(slide.shapes):
                o = tpl_shapes.get(pt, {}).get(s.shape_id)
                if o is None:
                    # 訊號 ④:add_textbox 新增的框。沒有模板對照就沒有設計師
                    # 基準可放寬——它的幾何是綁定自己宣告的,必須裝得進去。
                    if s.text_frame.text.strip() and not tt.estimate_overflow(s)["fits"]:
                        suspects.add(s.shape_id)
                    continue
                design = tt._first_run_size_pt(o)
                if tt._first_run_size_pt(s) < design - 0.5:
                    suspects.add(s.shape_id)
                    continue
                est = tt.estimate_overflow(s, size_pt=design)
                if est["lines"] <= 1:
                    continue
                allow = est["avail_pt"]
                if tt.has_autofit(s):
                    allow = max(allow, tt.estimate_overflow(o)["needed_pt"])
                if est["needed_pt"] > allow * COLLIDE_SLACK:
                    suspects.add(s.shape_id)
            # 訊號 ②:侵入鄰欄且比模板更糟;只收「範圍會隨文字變大」的那一方
            for area, s1, s2 in tt.text_collisions(slide):
                ref = base.get(frozenset((s1.shape_id, s2.shape_id)), 0.0)
                if area <= max(ref * COLLIDE_SLACK, COLLIDE_FLOOR):
                    continue
                for cand in (s1, s2):
                    g = _geo(slide, cand.shape_id, tt)
                    if g is None:
                        continue
                    cs, cL, cT, cW, cH = g
                    fp = tt.text_footprint(cs, cL, cT, cW, cH)
                    if fp and (fp[2] - (cL + cW) > 0.05 or fp[3] - (cT + cH) > 0.05):
                        suspects.add(cand.shape_id)

            for sid in sorted(suspects):
                paths = sid2paths.get(pt, {}).get(sid)
                g = _geo(slide, sid, tt)
                if g is None:
                    continue
                shape = g[0]
                if not paths:
                    # add_textbox 產生的框:從文字反推槽位名。一個框常同時承載
                    # 多個槽位(prefix + 兩個 slots 用 join 串起來),所以要比對
                    # 文字裡出現的**所有**識別字,不能只取第一個。
                    found = set(re.findall(r"[a-z_]{3,}", shape.text_frame.text))
                    paths = [p for p in addbox.get(pt, [])
                             if [x for x in p.split(".") if x != "item"][-1] in found] or None
                    if not paths:
                        unmapped.append((idx + 1, sid, shape.text_frame.text[:20]))
                        continue
                o = tpl_shapes.get(pt, {}).get(sid)
                design = tt._first_run_size_pt(o if o is not None else shape)
                ch = probe_char(paths[0])
                n = len(shape.text_frame.text.split("\n"))
                neighbours = _neighbour_footprints(slide, sid, base, tt)
                cap = _clean_cap(shape, g, n, ch, design, neighbours, tt, o)
                if cap < 2 and n <= 1:
                    continue        # 搜不到可行解 → 不是這一方能解決的
                if cap < MIN_CHARS and n > 1:
                    # 這個框裝不下這麼多「項目」:減字數救不了,要減清單長度。
                    # **本輪只改清單長度、不寫字數上限**——本工具只會收緊不會放寬,
                    # 若先按 n 項算出的小字數寫下去就永久卡住(曾寫出「1 字」)。
                    # 下一輪 golden 會用新的清單長度重新量,自然得到合理字數。
                    k = n
                    while k > 1 and _clean_cap(shape, g, k, ch, design, neighbours, tt, o) < MIN_CHARS:
                        k -= 1
                    shrunk_list = False
                    for path in paths:
                        lp, ln = parent_list(slots, path)
                        if not (ln and ln.get("kind") == "list"):
                            continue
                        # 清單總長 = 原長 − n + k(其他格位仍在,不是直接等於 k)
                        tgt = max(ln["max"] - n + k, ln["min"])
                        key = f"{pt}.slots.{lp}.max"
                        if overrides.get(key, ln["max"]) > tgt:
                            overrides[key] = tgt
                            changes.append((pt, lp, "清單", tgt))
                            shrunk_list = True
                    if shrunk_list:
                        continue
                    # 清單已縮到下限還是裝不下 → 這個版位放不了有意義的內容。
                    # 但本輪別急著喊:同一清單可能因為別的格位而在本輪被縮短,
                    # 下一輪重量就解決了。只有整輪都無變更時才是真的卡住。
                    stuck.append((idx + 1, pt, paths[0], sid))
                    continue
                for path in paths:
                    nd = node_at(slots, path)
                    if not nd or nd.get("kind") != "text":
                        continue
                    key = f"{pt}.slots.{path}.max_chars"
                    if overrides.get(key, nd["max_chars"]) > max(cap, 1):
                        overrides[key] = max(cap, 1)
                        changes.append((pt, path, "字數", max(cap, 1)))

        if unmapped and verbose:
            print(f"  ⚠ 第 {rnd} 輪有 {len(unmapped)} 個問題框對應不到槽位"
                  f"(綁定可能用了本工具不認識的 op 組合):")
            for pg, sid, txt in unmapped[:6]:
                print(f"      golden p{pg} shape id={sid} 「{txt}」")
        if stuck and not changes and verbose:
            print("  ⚠ 以下版位在任何字數下都放不了有意義的內容"
                  f"(至少 {MIN_CHARS} 個中文字)——建議該頁型降級 clone,"
                  "或請設計師把版位改大:")
            for pg, pt, path, sid in stuck:
                print(f"      {pt}.{path}(golden p{pg} shape id={sid})")
        if not changes:
            sync_registry_table(manifest, verbose=verbose)
            if verbose:
                print(f"✓ 第 {rnd} 輪:零縮字、零新增侵入 → 收斂"
                      + ("" if not unmapped else "(但有未對應框,見上)"))
            return 0 if not (unmapped or stuck) else 1
        if verbose:
            print(f"第 {rnd} 輪:收緊 {len(changes)} 項  "
                  + ", ".join(f"{a}.{b.split('.')[-1]}({c}→{d})" for a, b, c, d in changes[:4])
                  + (" …" if len(changes) > 4 else ""))
        manifest["capacity_overrides"] = dict(sorted(overrides.items()))
        if dry_run:
            print("(--dry-run:不寫入 manifest,停在第一輪提案)")
            for a, b, c, d in changes:
                print(f"    {a}.{b} {c} → {d}")
            return 0
        mpath.write_text(json.dumps(manifest, ensure_ascii=False, indent=2) + "\n",
                         encoding="utf-8")

    print(f"✗ {MAX_ROUNDS} 輪仍未收斂——多半是某個版位在任何字數下都會撞到鄰欄,"
          f"該頁型應降級 clone,或請設計師把該版位改大")
    return 1


def main(argv=None):
    ap = argparse.ArgumentParser(description="量測模板版位容量並寫入 capacity_overrides")
    ap.add_argument("--id", required=True, help="模板包 id 或目錄")
    ap.add_argument("--packs-root", help="模板包根目錄(預設 engine/templates)")
    ap.add_argument("--dry-run", action="store_true", help="只印第一輪提案,不寫入")
    ap.add_argument("--reset", action="store_true", help="先清空現有 overrides 再量")
    a = ap.parse_args(argv)
    root = Path(a.packs_root) if a.packs_root else REPO / "engine" / "templates"
    return run(a.id, root, dry_run=a.dry_run, reset=a.reset)


if __name__ == "__main__":
    sys.exit(main())
