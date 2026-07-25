#!/usr/bin/env python3
"""template_admin — 模板包註冊工具鏈單一入口(維護者/註冊 skill 專用,不入 tools.zip)。

子命令(TEMPLATE_PACKS.md §5):
  new --pptx <路徑> --id <id> --name <名>   建包骨架(status=draft)
  freeze --id <id>                          重建 inventory.json + template_sha256
  lint [--id <id> | --all]                  manifest/綁定靜態檢查(含全覆蓋原則)
  golden --id <id> [--page-types a,b]       黃金驗收:min/max 渲染+qa+冪等雙跑
  golden --regen-specs                      從 PAGE_TYPES 重派生 engine/golden/(契約改版用)
  register --id <id>                        原子性註冊:lint→golden→light 回歸→isolation
  pack [--id <id> | --tools]                打 template_<id>.zip / tools.zip(正斜線+檢查)
  isolation                                 git diff 對照模板隔離白名單
  list                                      列出所有包

直譯器:new/freeze/golden/register 需 python-pptx(本機無則用
`uv run --with python-pptx python`);lint/pack/isolation/list 純標準庫。
命令一律單行 python、相對路徑,三種 shell 原樣可跑。
"""
from __future__ import annotations

import argparse
import hashlib
import json
import re
import shutil
import subprocess
import sys
import tempfile
from datetime import datetime, timezone
from pathlib import Path

REPO = Path(__file__).resolve().parent.parent.parent
ENGINE = REPO / "engine"
TOOLS = ENGINE / "tools"
RULES = ENGINE / "rules"
GOLDEN_DIR = ENGINE / "golden"
DEFAULT_PACKS = ENGINE / "templates"
DIST = REPO / "gpts" / "dist"  # GPTs 前端的上傳產物(zips)
sys.path.insert(0, str(TOOLS))
sys.path.insert(0, str(RULES))

from validate_slide_spec_gpts import PAGE_TYPES, apply_capacity_overrides  # noqa: E402

ID_RE = re.compile(r"^[a-z][a-z0-9_]{2,31}$")
MODES = {"builtin", "fill", "clone", "unsupported"}
OPS = {"set", "delete", "keep", "rows", "list", "add_textbox", "resize"}


def sha256_file(p: Path) -> str:
    return hashlib.sha256(p.read_bytes()).hexdigest()


def load_manifest(pack_dir: Path) -> dict:
    return json.loads((pack_dir / "manifest.json").read_text(encoding="utf-8-sig"))


def save_manifest(pack_dir: Path, m: dict):
    with open(pack_dir / "manifest.json", "w", encoding="utf-8") as f:
        json.dump(m, f, ensure_ascii=False, indent=2)
        f.write("\n")


def pack_dir_of(args) -> Path:
    root = Path(args.packs_root) if getattr(args, "packs_root", None) else DEFAULT_PACKS
    return Path(args.id) if Path(args.id).is_dir() else root / args.id


def semantic_lib() -> set:
    """語意頁型全集 = validator PAGE_TYPES ∪ page_types.md 的 ### 條目名。"""
    names = set(PAGE_TYPES)
    pt_md = RULES / "page_types.md"
    if pt_md.exists():
        for line in pt_md.read_text(encoding="utf-8").splitlines():
            m = re.match(r"^###\s+([a-z][a-z0-9_]+)\s*$", line)
            if m:
                names.add(m.group(1))
    return names


# ---------------------------------------------------------------------------
# new / freeze
# ---------------------------------------------------------------------------
def cmd_new(args) -> int:
    if not ID_RE.match(args.id):
        print(f"✗ id 不合法(^[a-z][a-z0-9_]{{2,31}}$):{args.id}")
        return 2
    src = Path(args.pptx)
    if not src.exists():
        print(f"✗ 找不到模板檔:{src}")
        return 2
    root = Path(args.packs_root) if args.packs_root else DEFAULT_PACKS
    pack_dir = root / args.id
    if (pack_dir / "manifest.json").exists():
        print(f"✗ 模板包 {args.id} 已存在:{pack_dir}")
        return 2
    from pptx import Presentation
    n_pages = len(Presentation(str(src)).slides)
    (pack_dir / "assets_src").mkdir(parents=True, exist_ok=True)
    (pack_dir / "examples").mkdir(exist_ok=True)
    shutil.copy2(src, pack_dir / "template.pptx")
    manifest = {
        "template_id": args.id,
        "display_name": args.name,
        "version": datetime.now(timezone.utc).astimezone().strftime("%Y-%m-%d") + ".1",
        "status": "draft",
        "template_file": "template.pptx",
        "template_sha256": sha256_file(pack_dir / "template.pptx"),
        "page_count": n_pages,
        "aspect": "16:9",
        "language": "zh-TW",
        "style": {"font_zh": "Microsoft JhengHei", "font_en": "Helvetica",
                  "allowed_fonts": [], "colors": {}},
        "asset_defaults": {},
        "page_number": {"box_in": [12.30, 6.72, 0.70, 0.50], "size_pt": 28,
                        "color": "344252",
                        "clear_zone_in": {"left": 11.2, "top": 6.3},
                        "detect_zone_in": {"left": 11.0, "top": 6.3}},
        "page_types": {},
        "capacity_overrides": {},
        "registered_by": "", "registered_at": "",
    }
    save_manifest(pack_dir, manifest)
    with open(pack_dir / "registration_state.json", "w", encoding="utf-8") as f:
        json.dump({"confirmed_pages": {}, "note": "映射確認進度(register skill 用)"},
                  f, ensure_ascii=False, indent=2)
    print(f"已建包骨架:{pack_dir}({n_pages} 頁,status=draft)。"
          f"下一步:freeze → 映射確認 → 寫 bindings.json → lint/golden")
    return 0


def cmd_freeze(args) -> int:
    pack_dir = pack_dir_of(args)
    m = load_manifest(pack_dir)
    tpl = pack_dir / m.get("template_file", "template.pptx")
    if not tpl.exists():  # light 特例已於 Phase 1 移入包內,此兜底防呆
        print(f"✗ 找不到模板檔:{tpl}")
        return 2
    from pptx import Presentation
    from inspect_template import _shape_info
    prs = Presentation(str(tpl))
    fill_pages = sorted({e["template_page"] for e in m.get("page_types", {}).values()
                         if e.get("mode") == "fill"})
    pages = fill_pages or range(1, len(prs.slides) + 1)  # 尚無 fill 宣告 → 全頁快照
    inv = {"template_sha256": sha256_file(tpl),
           "generated_at": datetime.now(timezone.utc).astimezone().strftime("%Y-%m-%d"),
           "pages": {str(pg): [_shape_info(s) for s in prs.slides[pg - 1].shapes]
                     for pg in pages}}
    with open(pack_dir / "inventory.json", "w", encoding="utf-8") as f:
        json.dump(inv, f, ensure_ascii=False, separators=(",", ":"))
        f.write("\n")
    m["template_sha256"] = inv["template_sha256"]
    m["page_count"] = len(prs.slides)
    save_manifest(pack_dir, m)
    print(f"freeze 完成:{len(inv['pages'])} 頁快照,sha={inv['template_sha256'][:12]}…")
    return 0


# ---------------------------------------------------------------------------
# lint
# ---------------------------------------------------------------------------
def _binding_ids(ent: dict) -> set:
    ids = set()
    for op in ent.get("ops", []):
        ids.update(op.get("ids", []))
        for key in ("id",):
            if key in op:
                ids.add(op[key])
        ids.update(op.get("row_ids", []))
        ids.update(op.get("sep_ids", []))
        if "overflow" in op and isinstance(op["overflow"], dict):
            ids.add(op["overflow"].get("merge_into_id"))
        for item in op.get("items", []):
            ids.update(s["id"] for s in item.get("sets", []))
            ids.update(item.get("delete_always", []))
            ids.update(item.get("delete_on_missing", []))
        ids.update(op.get("delete_when_empty", []))
    ids.discard(None)
    return ids


def _flatten_shapes(shapes, acc):
    for s in shapes:
        acc[s["id"]] = s
        _flatten_shapes(s.get("ch", []), acc)
    return acc


def _cover(shapes, covered: set):
    """referenced 集合的覆蓋閉包:群組被引用(刪/留)則其子形狀視為已覆蓋。"""
    out = set(covered)
    def walk(items, parent_covered):
        for s in items:
            c = parent_covered or s["id"] in out
            if c:
                out.add(s["id"])
            walk(s.get("ch", []), c)
    walk(shapes, False)
    return out


def lint_pack(pack_dir: Path) -> list:
    errs = []
    m = load_manifest(pack_dir)
    pid = m.get("template_id", pack_dir.name)
    for field in ("template_id", "display_name", "version", "status", "template_file",
                  "template_sha256", "page_count", "page_types"):
        if not m.get(field) and m.get(field) != 0:
            errs.append(f"manifest 缺欄位 {field!r}")
    if m.get("template_id") and not ID_RE.match(m["template_id"]):
        errs.append(f"template_id 不合法:{m['template_id']}")
    if m.get("status") not in ("draft", "registered"):
        errs.append(f"status 不合法:{m.get('status')!r}")
    lib = semantic_lib()
    n_pages = m.get("page_count", 0)
    for pt, e in m.get("page_types", {}).items():
        mode = e.get("mode")
        if mode not in MODES:
            errs.append(f"page_types.{pt}: mode 不合法 {mode!r}")
            continue
        if mode == "builtin" and pid != "light":
            errs.append(f"page_types.{pt}: builtin 僅 light 允許(新模板一律 fill/clone)")
        if pt not in lib:
            errs.append(f"page_types.{pt}: 不存在於語意契約或 page_types.md 語意庫")
        if mode in ("fill", "clone"):
            pg = e.get("template_page")
            if not isinstance(pg, int) or not (1 <= pg <= n_pages):
                errs.append(f"page_types.{pt}: template_page {pg!r} 超出 1–{n_pages}")
        if mode == "fill" and pt not in PAGE_TYPES:
            errs.append(f"page_types.{pt}: fill 級必須是已註冊語意頁型(PAGE_TYPES)")
    try:
        apply_capacity_overrides(PAGE_TYPES, m.get("capacity_overrides"))
    except ValueError as e:
        errs.append(str(e))
    # 素材:有頁型要求素材鍵時,asset_defaults 指到的檔必須存在(包內或 assets_src)
    ad = m.get("asset_defaults") or {}
    need_assets = any((e.get("assets") if e.get("assets") is not None
                       else PAGE_TYPES.get(pt, {}).get("assets", []))
                      for pt, e in m.get("page_types", {}).items()
                      if e.get("mode") in ("builtin", "fill"))
    if need_assets:
        for key, rel in ad.items():
            cands = [pack_dir / rel, pack_dir / rel.replace("assets/", "assets_src/", 1)]
            if not any(c.exists() for c in cands):
                errs.append(f"asset_defaults.{key}: 檔案不存在({rel})")
    # 綁定
    inv_path = pack_dir / "inventory.json"
    bj = pack_dir / "bindings.json"
    if (pack_dir / "bindings.py").exists():
        print(f"  [i] {pid}: 有 bindings.py(BUILDERS/選配 FILLS grandfather,僅 light);"
              f"fills 以 bindings.json 為準(py 匯出非空 FILLS 才覆蓋)")
    if bj.exists():
        if not inv_path.exists():
            errs.append("有 bindings.json 但缺 inventory.json(先跑 freeze)")
        else:
            inv = json.loads(inv_path.read_text(encoding="utf-8-sig"))
            if inv.get("template_sha256") != m.get("template_sha256"):
                errs.append("inventory sha 與 manifest 不符(模板改版?重跑 freeze)")
            data = json.loads(bj.read_text(encoding="utf-8-sig"))
            if data.get("engine") != 1:
                errs.append(f"bindings.json engine 版本不符:{data.get('engine')!r}")
            fills_m = {pt for pt, e in m.get("page_types", {}).items()
                       if e.get("mode") == "fill"}
            fills_b = set(data.get("fills", {}))
            for pt in sorted(fills_m - fills_b):
                errs.append(f"manifest fill 頁型 {pt} 缺 bindings 條目")
            for pt in sorted(fills_b - fills_m):
                errs.append(f"bindings 有 {pt} 但 manifest 未宣告為 fill")
            for pt, ent in data.get("fills", {}).items():
                mpg = m.get("page_types", {}).get(pt, {}).get("template_page")
                if mpg is not None and ent.get("template_page") != mpg:
                    errs.append(f"{pt}: bindings template_page {ent.get('template_page')}"
                                f" ≠ manifest {mpg}")
                snap = inv.get("pages", {}).get(str(ent.get("template_page")))
                if snap is None:
                    errs.append(f"{pt}: inventory 缺第 {ent.get('template_page')} 頁快照"
                                f"(重跑 freeze)")
                    continue
                flat = _flatten_shapes(snap, {})
                # chart/SmartArt 頁禁 fill(靜默捏造源,TEMPLATE_PACKS §3)
                if any(s.get("cht") for s in flat.values()):
                    errs.append(f"{pt}: 模板第 {ent.get('template_page')} 頁含 chart,"
                                f"禁止註冊為 fill(改 clone 或 unsupported)")
                if any(s.get("smart") for s in flat.values()):
                    errs.append(f"{pt}: 模板頁含 SmartArt,一律 unsupported")
                for i, op in enumerate(ent.get("ops", [])):
                    if op.get("op") not in OPS:
                        errs.append(f"{pt} ops[{i}]: 未知 op {op.get('op')!r}")
                refd = _binding_ids(ent)
                missing = sorted(i for i in refd if i not in flat)
                if missing:
                    errs.append(f"{pt}: 綁定引用不存在的 shape id {missing}")
                covered = _cover(snap, refd)
                # 引擎豁免:清除窗內純數字 ≤3 字 = 模板頁碼,渲染時
                # _finalize_page_number 一律清掉,綁定不必覆蓋
                zone = (m.get("page_number") or {}).get("clear_zone_in",
                                                        {"left": 11.2, "top": 6.3})

                def engine_covered(s):
                    t = (s.get("txt") or "").strip()
                    return (t.isdigit() and len(t) <= 3
                            and (s.get("x") or 0) > zone["left"]
                            and (s.get("y") or 0) > zone["top"])
                naked = sorted(i for i, s in flat.items()
                               if s.get("txt") and i not in covered
                               and not engine_covered(s))
                if naked:
                    errs.append(f"{pt}: 全覆蓋原則違反——含文字 shape 未被"
                                f" set/rows/list/delete/keep 覆蓋:{naked}")
    elif not (pack_dir / "bindings.py").exists():
        if any(e.get("mode") == "fill" for e in m.get("page_types", {}).values()):
            errs.append("宣告了 fill 頁型但缺 bindings(.json/.py)")
    return errs


def cmd_lint(args) -> int:
    root = Path(args.packs_root) if args.packs_root else DEFAULT_PACKS
    dirs = ([p for p in sorted(root.iterdir()) if (p / "manifest.json").exists()]
            if args.all else [pack_dir_of(args)])
    bad = 0
    for d in dirs:
        errs = lint_pack(d)
        tag = load_manifest(d).get("template_id", d.name)
        if errs:
            bad += 1
            print(f"✗ {tag}: {len(errs)} 項")
            for e in errs:
                print(f"   {e}")
        else:
            print(f"✓ {tag}: lint OK")
    return 1 if bad else 0


# ---------------------------------------------------------------------------
# golden
# ---------------------------------------------------------------------------
def _variant_text(name, max_chars, variant):
    if name == "value":
        return "99.9%"
    if name == "number":
        return "01"
    if variant == "min":
        t = f"【{name}】待填"
        return t if len(t) <= max_chars else "待填"[:max_chars]
    return (name + "測" * max_chars)[:max_chars]


def _variant_slot(name, slot, variant, page_type):
    from make_skeleton import FIXED
    kind = slot["kind"]
    if kind == "text":
        fixed = FIXED.get((page_type, name))
        return fixed if fixed else _variant_text(name, slot["max_chars"], variant)
    if kind == "list":
        n = slot["min"] if variant == "min" else slot["max"]
        return [_variant_slot(name, slot["item"], variant, page_type)
                for _ in range(max(n, slot["min"]))]
    if kind == "object":
        return {fn: _variant_slot(fn, fs, variant, page_type)
                for fn, fs in slot["fields"].items()
                if variant == "max" or fs.get("required", True)}
    raise ValueError(kind)


def derive_golden_spec(pack_id, contracts, page_types_entries, asset_defaults,
                       only=None) -> dict:
    slides = []
    fill_pts = [pt for pt, e in page_types_entries.items() if e.get("mode") == "fill"
                and (not only or pt in only)]
    cover_bg = asset_defaults.get("background_cover", "assets/backgrounds/cover_bg.png")
    content_bg = asset_defaults.get("background_content", "assets/backgrounds/content_bg.png")
    logo = asset_defaults.get("logo", "assets/logos/cathay_logo.png")
    for pt in fill_pts:
        contract = contracts[pt]
        for variant in ("min", "max"):
            num = len(slides) + 1
            required_assets = page_types_entries[pt].get("assets")
            if required_assets is None:
                required_assets = contract["assets"]
            assets = {}
            for key in required_assets:
                if key == "logo":
                    assets[key] = logo
                else:
                    assets[key] = cover_bg if pt in ("cover", "closing") else content_bg
            slides.append({
                "number": num,
                "page_type": pt,
                "title": _variant_text("title", 30, variant),
                "render_page_number": contract["page_number"] == "required",
                "assets": assets,
                "slots": {sn: _variant_slot(sn, ss, variant, pt)
                          for sn, ss in contract["slots"].items()
                          if variant == "max" or ss.get("required", True)},
            })
    return {"deck": {"deck_name": f"golden_{pack_id}", "language": "Traditional Chinese",
                     "slide_count": len(slides), "template": pack_id},
            "slides": slides}


def _run(label, cmd) -> int:
    print(f"  $ {' '.join(str(c) for c in cmd[1:])}", flush=True)
    rc = subprocess.run([sys.executable] + [str(c) for c in cmd[1:]]).returncode
    if rc != 0:
        print(f"✗ golden 失敗於「{label}」(exit {rc})")
    return rc


def run_golden(pack_dir: Path, only=None, out_dir: Path | None = None) -> int:
    m = load_manifest(pack_dir)
    pid = m["template_id"]
    contracts = apply_capacity_overrides(PAGE_TYPES, m.get("capacity_overrides"))
    spec = derive_golden_spec(pid, contracts, m.get("page_types", {}),
                              m.get("asset_defaults") or {}, only)
    if not spec["slides"]:
        print(f"✗ {pid}: 無 fill 頁型可驗(--page-types 過濾後為空?)")
        return 2
    out_dir = out_dir or (REPO / "ppt_out")
    out_dir.mkdir(exist_ok=True)
    work = Path(tempfile.mkdtemp(prefix=f"golden_{pid}_"))
    # 素材沙箱:包 assets/(或 assets_src 映射)→ work/assets
    src_assets = pack_dir / "assets"
    if not src_assets.is_dir():
        src_assets = pack_dir / "assets_src"
    if src_assets.is_dir():
        shutil.copytree(src_assets, work / "assets", dirs_exist_ok=True)
    spec_path = work / "golden_spec.json"
    with open(spec_path, "w", encoding="utf-8") as f:
        json.dump(spec, f, ensure_ascii=False, indent=1)
    shutil.copy2(spec_path, out_dir / f"golden_{pid}.spec.json")

    no_slides = work / "run.no-slides"
    if _run("validator", [sys.executable, RULES / "validate_slide_spec_gpts.py",
                          "--spec", spec_path, "--slides", no_slides,
                          "--asset-dir", work, "--template-pack", pack_dir,
                          "--allow-draft"]):
        return 1
    deck1, deck2 = work / "g1.pptx", work / "g2.pptx"
    for deck in (deck1, deck2):  # 連跑兩次:冪等實證
        if _run("render", [sys.executable, TOOLS / "render_deck.py", "--spec", spec_path,
                           "--template-pack", pack_dir, "--asset-dir", work,
                           "--out", deck]):
            return 1
    s1, s2 = work / "g1.json", work / "g2.json"
    for deck, sj in ((deck1, s1), (deck2, s2)):
        if _run("inspect", [sys.executable, TOOLS / "inspect_template.py",
                            "--pptx", deck, "--all", "--out", sj]):
            return 1
    if json.loads(s1.read_text()) != json.loads(s2.read_text()):
        print("✗ 冪等實證失敗:連跑兩次 shape 樹不一致(渲染層有隨機性?)")
        return 1
    if _run("qa", [sys.executable, TOOLS / "qa_check.py", "--spec", spec_path,
                   "--pptx", deck1, "--template-pack", pack_dir]):
        return 1
    final = out_dir / f"golden_{pid}.pptx"
    shutil.copy2(deck1, final)
    print(f"golden PASS:{len(spec['slides'])} 頁(每 fill 頁型 min+max)"
          f",冪等雙跑一致 → 目檢檔 {final}")
    return 0


def cmd_golden(args) -> int:
    if args.regen_specs:
        GOLDEN_DIR.mkdir(exist_ok=True)
        n = 0
        for pt, contract in PAGE_TYPES.items():
            for variant in ("min", "max"):
                spec = derive_golden_spec(
                    "light", PAGE_TYPES, {pt: {"mode": "fill", "template_page": 0}},
                    {}, only=[pt])
                keep = [s for s in spec["slides"]
                        if (variant == "min") == (s is spec["slides"][0])]
                spec["slides"] = [dict(keep[0], number=1)]
                spec["deck"]["slide_count"] = 1
                spec["deck"]["deck_name"] = f"golden_{pt}_{variant}"
                with open(GOLDEN_DIR / f"{pt}.{variant}.json", "w", encoding="utf-8") as f:
                    json.dump(spec, f, ensure_ascii=False, indent=1)
                    f.write("\n")
                n += 1
        print(f"已重派生 {n} 份 golden fixtures → {GOLDEN_DIR}"
              f"(基準契約;各包實跑時依 merged 契約即時派生)")
        return 0
    only = [t.strip() for t in args.page_types.split(",")] if args.page_types else None
    return run_golden(pack_dir_of(args), only,
                      Path(args.out) if args.out else None)


# ---------------------------------------------------------------------------
# register / pack / isolation / list
# ---------------------------------------------------------------------------
def cmd_register(args) -> int:
    pack_dir = pack_dir_of(args)
    m = load_manifest(pack_dir)
    pid = m["template_id"]
    print(f"== register {pid}:lint ==")
    errs = lint_pack(pack_dir)
    if errs:
        for e in errs:
            print(f"   {e}")
        print("✗ lint 未過,status 維持 draft")
        return 1
    print(f"== register {pid}:golden all(不信任舊戳記)==")
    if run_golden(pack_dir):
        print("✗ golden 未過,status 維持 draft")
        return 1
    if pid != "light":
        print("== register:light 回歸(隔離驗證)==")
        light = DEFAULT_PACKS / "light"
        if run_golden(light):
            print("✗ light 回歸紅 → 共用檔疑似被動到,本次註冊掛起(回報維護者)")
            return 1
    print("== register:isolation ==")
    if cmd_isolation(args, advisory_pack=pid):
        print("⚠ isolation 有越界檔(見上);模板目錄外的改動請拆 commit")
    m["status"] = "registered"
    m["registered_at"] = datetime.now(timezone.utc).astimezone().strftime("%Y-%m-%d")
    m["golden"] = {"last_pass": datetime.now(timezone.utc).astimezone().isoformat(timespec="seconds")}
    save_manifest(pack_dir, m)
    state = pack_dir / "registration_state.json"
    if state.exists():
        state.unlink()  # 註冊完成,draft 進度檔移除
    fills = [pt for pt, e in m["page_types"].items() if e.get("mode") in ("builtin", "fill")]
    clones = [pt for pt, e in m["page_types"].items() if e.get("mode") == "clone"]
    unsup = [pt for pt, e in m["page_types"].items() if e.get("mode") == "unsupported"]
    print(f"★ registered:{pid}@{m['version']} 全自動 {len(fills)} / 半自動 {len(clones)}"
          f" / 不支援 {len(unsup)}")
    print("發佈清單(人工,見 gpts/README.md 多模板發佈 checklist):")
    print(f"  1. python engine/release/template_admin.py pack --id {pid}")
    print(f"  2. engine/templates/INDEX.md 加列 + instructions.md roster/版本字串")
    print(f"  3. Builder 刪舊傳新 template_{pid}.zip;light 抽測一份")
    return 0


PACK_WHITELIST = ("template.pptx", "manifest.json", "bindings.py", "bindings.json",
                  "page_map.md")


def cmd_pack(args) -> int:
    import zipfile
    if args.tools:
        out = DIST / "tools.zip"
        names = sorted(f.name for f in TOOLS.iterdir() if f.suffix in (".py", ".md"))
        with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as z:
            for n in names:
                z.write(TOOLS / n, arcname=n)
        print(f"tools.zip:{len(names)} 檔")
    else:
        pack_dir = pack_dir_of(args)
        pid = load_manifest(pack_dir)["template_id"]
        out = DIST / f"template_{pid}.zip"
        with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as z:
            for n in PACK_WHITELIST:
                if (pack_dir / n).exists():
                    z.write(pack_dir / n, arcname=n)
            assets = pack_dir / "assets"
            src_map = "assets"
            if not assets.is_dir():
                assets, src_map = pack_dir / "assets_src", "assets_src"
            if assets.is_dir():
                for f in sorted(assets.rglob("*")):
                    if f.is_file():
                        rel = f.relative_to(assets).as_posix()
                        z.write(f, arcname=f"assets/{rel}")
        print(f"{out.name} 打包完成(素材源:{src_map})")
    bad = [i.orig_filename for i in zipfile.ZipFile(out).infolist()
           if "\\" in i.orig_filename]
    if bad:
        print(f"✗ zip 內含反斜線路徑:{bad}")
        return 1
    print(f"sha256 {sha256_file(out)}  {out.name}(REGRESSION R7 基準需同步)")
    return 0


def cmd_isolation(args, advisory_pack=None) -> int:
    r = subprocess.run(["git", "-C", str(REPO), "status", "--porcelain"],
                       capture_output=True, text=True)
    paths = [l[3:].split(" -> ")[-1].strip('"') for l in r.stdout.splitlines() if l.strip()]
    touched_ids = {p.split("/")[2] for p in paths
                   if p.startswith("engine/templates/") and len(p.split("/")) > 3}
    zips = {re.match(r"gpts/dist/template_([a-z0-9_]+)\.zip", p).group(1)
            for p in paths if re.match(r"gpts/dist/template_[a-z0-9_]+\.zip", p)}
    ids = (touched_ids | zips) - {"INDEX.md", "TEMPLATE_LIFECYCLE.md"}
    if advisory_pack:
        ids = {advisory_pack}
    if len(ids) != 1:
        print(f"[i] isolation:涉及模板 {sorted(ids) or '無'}(單一模板 commit 才適用白名單)")
        return 0
    (tid,) = ids
    allow = [f"engine/templates/{tid}/", f"gpts/dist/template_{tid}.zip",
             "engine/templates/INDEX.md", "gpts/instructions.md"]
    viol = [p for p in paths if not any(p.startswith(a) for a in allow)]
    if viol:
        print(f"✗ 模板 {tid} 的變更越界({len(viol)} 檔;共用檔/他包改動請拆 commit):")
        for p in viol:
            print(f"   {p}")
        return 1
    print(f"✓ isolation OK:變更僅及模板 {tid} 白名單")
    return 0


def cmd_list(args) -> int:
    root = Path(args.packs_root) if args.packs_root else DEFAULT_PACKS
    for p in sorted(root.iterdir()):
        if (p / "manifest.json").exists():
            m = load_manifest(p)
            modes = [e.get("mode") for e in m.get("page_types", {}).values()]
            print(f"{m.get('template_id', p.name):<12} {m.get('version','?'):<14} "
                  f"{m.get('status','?'):<11} 全自動 {sum(x in ('builtin','fill') for x in modes)}"
                  f" / 半自動 {sum(x == 'clone' for x in modes)}"
                  f" / 不支援 {sum(x == 'unsupported' for x in modes)}  {m.get('display_name','')}")
    return 0


def main(argv):
    ap = argparse.ArgumentParser(description=__doc__)
    sub = ap.add_subparsers(dest="cmd", required=True)

    def common(p, need_id=True):
        if need_id:
            p.add_argument("--id", required=True, help="模板包 id 或目錄路徑")
        p.add_argument("--packs-root")

    p = sub.add_parser("new"); common(p)
    p.add_argument("--pptx", required=True); p.add_argument("--name", required=True)
    common(sub.add_parser("freeze"))
    p = sub.add_parser("lint"); p.add_argument("--id"); p.add_argument("--all", action="store_true")
    p.add_argument("--packs-root")
    p = sub.add_parser("golden"); p.add_argument("--id"); p.add_argument("--packs-root")
    p.add_argument("--page-types"); p.add_argument("--out")
    p.add_argument("--regen-specs", action="store_true")
    common(sub.add_parser("register"))
    p = sub.add_parser("pack"); p.add_argument("--id"); p.add_argument("--tools", action="store_true")
    p.add_argument("--packs-root")
    p = sub.add_parser("isolation"); p.add_argument("--packs-root")
    p = sub.add_parser("list"); p.add_argument("--packs-root")

    args = ap.parse_args(argv)
    if args.cmd == "lint" and not (args.id or args.all):
        print("✗ lint 需要 --id 或 --all")
        return 2
    if args.cmd == "golden" and not (args.id or args.regen_specs):
        print("✗ golden 需要 --id 或 --regen-specs")
        return 2
    if args.cmd == "pack" and not (args.id or args.tools):
        print("✗ pack 需要 --id 或 --tools")
        return 2
    return {"new": cmd_new, "freeze": cmd_freeze, "lint": cmd_lint,
            "golden": cmd_golden, "register": cmd_register, "pack": cmd_pack,
            "isolation": cmd_isolation, "list": cmd_list}[args.cmd](args)


if __name__ == "__main__":
    sys.exit(main(sys.argv[1:]))
