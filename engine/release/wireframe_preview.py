#!/usr/bin/env python3
"""wireframe_preview — 把 pptx 畫成線框示意圖(本機目檢輔助,WORKLOG §8 選配)。

緩解「沙箱無法產可靠預覽」的痛點:每頁輸出一張 PNG(形狀外框 + 文字),
輔助 golden 目檢與設計師確認版面結構;**線框僅示意位置與文字,字級/配色/
換行以 PowerPoint 開檔為準**(qa 的溢出 WARN 仍要人工開檔看)。

用法(需 python-pptx + pillow;本機:uv run --with python-pptx --with pillow python):
  python engine/release/wireframe_preview.py --pptx ppt_out/golden_light.pptx --out ppt_out/wf
  # 產出 ppt_out/wf/page_01.png … 與 overview.png(全頁縮圖網格)

中文字型:自動偵測系統字型(macOS PingFang / Windows msjh / Linux Noto);
都找不到時退回 PIL 內建字型(CJK 會變豆腐,線框仍可看結構)。
"""
from __future__ import annotations

import argparse
import sys
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

SCALE = 100  # px per inch → 16:9 = 1333x750
EMU_PER_IN = 914400

FONT_CANDIDATES = [
    "/System/Library/Fonts/PingFang.ttc",            # macOS
    "/System/Library/Fonts/STHeiti Light.ttc",
    "C:/Windows/Fonts/msjh.ttc",                     # Windows 微軟正黑
    "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",  # Linux
]


def load_font(size):
    for p in FONT_CANDIDATES:
        if Path(p).exists():
            try:
                return ImageFont.truetype(p, size)
            except OSError:
                continue
    return ImageFont.load_default()


def _px(emu):
    return int((emu or 0) / EMU_PER_IN * SCALE)


def draw_shape(d, shp, font, ox=0, oy=0):
    x, y = _px(shp.left) + ox, _px(shp.top) + oy
    w, h = _px(shp.width), _px(shp.height)
    st = shp.shape_type
    if st == MSO_SHAPE_TYPE.GROUP:
        d.rectangle([x, y, x + w, y + h], outline=(200, 160, 60), width=1)
        for c in shp.shapes:
            draw_shape(d, c, font)  # 群組子形狀座標已是絕對值(pptx 幾何換算簡化)
        return
    if st == MSO_SHAPE_TYPE.PICTURE:
        d.rectangle([x, y, x + w, y + h], outline=(120, 120, 220), width=2)
        d.line([x, y, x + w, y + h], fill=(120, 120, 220), width=1)
        d.line([x + w, y, x, y + h], fill=(120, 120, 220), width=1)
        return
    color = (60, 60, 60) if getattr(shp, "has_text_frame", False) else (170, 170, 170)
    d.rectangle([x, y, x + w, y + h], outline=color, width=1)
    if getattr(shp, "has_text_frame", False):
        txt = shp.text_frame.text.strip().replace("\n", " / ")
        if txt:
            d.text((x + 3, y + 2), txt[:40], fill=(20, 20, 20), font=font)


def render_page(slide, size, font):
    img = Image.new("RGB", size, (250, 250, 250))
    d = ImageDraw.Draw(img)
    d.rectangle([0, 0, size[0] - 1, size[1] - 1], outline=(0, 0, 0), width=1)
    for shp in slide.shapes:
        draw_shape(d, shp, font)
    return img


def main(argv):
    ap = argparse.ArgumentParser()
    ap.add_argument("--pptx", required=True)
    ap.add_argument("--out", required=True, help="輸出目錄(page_NN.png + overview.png)")
    ap.add_argument("--cols", type=int, default=4, help="overview 每列頁數")
    args = ap.parse_args(argv)

    prs = Presentation(args.pptx)
    size = (_px(prs.slide_width), _px(prs.slide_height))
    out = Path(args.out)
    out.mkdir(parents=True, exist_ok=True)
    font = load_font(14)
    pages = []
    for i, slide in enumerate(prs.slides, 1):
        img = render_page(slide, size, font)
        img.save(out / f"page_{i:02d}.png")
        pages.append(img)
    # overview 網格
    if pages:
        cols = max(1, args.cols)
        rows = (len(pages) + cols - 1) // cols
        tw, th = size[0] // 3, size[1] // 3
        grid = Image.new("RGB", (tw * cols, th * rows), (255, 255, 255))
        for i, img in enumerate(pages):
            grid.paste(img.resize((tw, th)), ((i % cols) * tw, (i // cols) * th))
        grid.save(out / "overview.png")
    print(f"線框示意:{len(pages)} 頁 → {out}/(page_NN.png + overview.png);"
          f"字級/溢出仍以 PowerPoint 開檔為準")
    return 0


if __name__ == "__main__":
    sys.exit(main(sys.argv[1:]))
