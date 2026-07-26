#!/usr/bin/env python3
"""regress — 引擎級回歸 runner:REGRESSION.md 全部案例一鍵跑完。

用法(repo 根目錄執行;三種 shell 通用):
  python3 engine/release/regress.py               # 全部案例(R0–R13 + light R-L0/R-L1)
  python3 engine/release/regress.py --list        # 列案例
  python3 engine/release/regress.py --only R2,R12 # 只跑指定(相依案例自動帶入)
  python3 engine/release/regress.py --skip R9     # 跳過指定(會列為 SKIP 並使結果非綠)

為什麼要有這支:REGRESSION.md 的案例本來是「散文夾貼指令」,要人逐段複製貼上
比對預期——摩擦高到必然被跳過(R9 曾紅 11 個 commit 沒人發現,就是因為沒有人
真的跑過它)。本 repo 的教訓是「凡是靠人記得的不變式,在自動化工具面前一定會
被覆蓋」;回歸本身也是不變式,所以收進工具。REGRESSION.md 仍是各案例的
**規格與解說**(為什麼這樣測、預期輸出的由來),本檔只是它的可執行形式;
新增案例要兩邊一起動。

設計約束:
- 純標準庫 + subprocess,無 shell 專屬語法 → Windows(禁 WSL)原樣可跑,
  解掉 REGRESSION.md「命令為 POSIX shell 語法」的限制。
- 渲染類步驟自動偵測 python-pptx;沒有就用 `uv run --with python-pptx python`
  前綴(同 skill 的 prepare_env 慣例),兩者皆無 → 該案例 FAIL 並說明環境缺件。
- 會暫時污染工作區的案例(R10 種 bindings.py、R13 改 INDEX.md)一律
  try/finally 還原;R11 失敗時**刻意不還原**——留下的 diff 就是修法
  (把重派生結果一起 commit),還原反而湮滅證據。
- R7 的 sha 基準與檔數**從 REGRESSION.md 解析**,不在本檔重複——基準只有
  一份真相,漂移(改了 zip 沒更新 R7)在這裡變成紅燈。
"""
from __future__ import annotations

import argparse
import hashlib
import json
import re
import shutil
import subprocess
import sys
import tempfile
import time
import zipfile
from pathlib import Path

REPO = Path(__file__).resolve().parent.parent.parent
ENGINE = REPO / "engine"
RULES = ENGINE / "rules"
TOOLS = ENGINE / "tools"
EXAMPLES = ENGINE / "examples"
LIGHT = ENGINE / "templates" / "light"
ADMIN = ENGINE / "release" / "template_admin.py"
REG_MD = ENGINE / "REGRESSION.md"

TIMEOUT = 1200  # 單一子程序上限(register 演練含多輪 golden,放寬)


def sha256_file(p: Path) -> str:
    return hashlib.sha256(p.read_bytes()).hexdigest()


class Env:
    """跨案例共享狀態:$RT 沙箱、pptx 直譯器、R2 產物。"""

    def __init__(self):
        self.rt: Path | None = None
        self.spec02_fixed: Path | None = None
        self.deck02: Path | None = None
        self._pptx_py: list[str] | None = None
        self.tempdirs: list[Path] = []

    def mktemp(self, prefix) -> Path:
        d = Path(tempfile.mkdtemp(prefix=prefix))
        self.tempdirs.append(d)
        return d

    def pptx_py(self) -> list[str]:
        """渲染用直譯器:本機有 python-pptx 用 python3,否則 uv 前綴。"""
        if self._pptx_py is None:
            probe = subprocess.run([sys.executable, "-c", "import pptx"],
                                   capture_output=True)
            if probe.returncode == 0:
                self._pptx_py = [sys.executable]
            elif shutil.which("uv"):
                self._pptx_py = ["uv", "run", "--with", "python-pptx", "python"]
            else:
                raise RuntimeError("需要 python-pptx 或 uv(渲染類案例無法執行)")
        return self._pptx_py


def run(cmd, timeout=TIMEOUT) -> tuple[int, str]:
    """跑子程序,回 (exit, stdout+stderr)。一律 repo 根目錄、不經 shell。"""
    r = subprocess.run([str(c) for c in cmd], cwd=REPO, timeout=timeout,
                       capture_output=True, text=True, encoding="utf-8")
    return r.returncode, (r.stdout or "") + (r.stderr or "")


def run_py_snippet(env: Env, code: str, need_pptx=False) -> tuple[int, str]:
    """把多行 python 存成暫存檔執行(REGRESSION.md 的 heredoc 等價物,Windows 可跑)。"""
    d = env.mktemp("regress_snip_")
    f = d / "snippet.py"
    f.write_text(code, encoding="utf-8")
    interp = env.pptx_py() if need_pptx else [sys.executable]
    return run(interp + [f])


# ---------------------------------------------------------------------------
# 案例實作:每個回 (ok, detail);detail 在 FAIL 時完整印出
# ---------------------------------------------------------------------------
def r0(env: Env):
    """archive 完整性與環境建置(鏡射 GPTs Step 0 佈局)。"""
    rt = env.mktemp("regress_rt_")
    lines = []
    for zname, dest in (("tools.zip", rt / "tools"),
                        ("template_light.zip", rt / "templates" / "light")):
        zp = REPO / "gpts" / "dist" / zname
        zf = zipfile.ZipFile(zp)
        bad = [i.orig_filename for i in zf.infolist() if "\\" in i.orig_filename]
        if bad:
            return False, f"{zname} 含反斜線路徑:{bad}"
        dest.mkdir(parents=True, exist_ok=True)
        zf.extractall(dest)
        lines.append(f"{zname} OK")
    shutil.copy2(RULES / "validate_slide_spec_gpts.py", rt)
    need = [rt / "tools" / "run_pipeline.py", rt / "tools" / "README_TOOLS.md",
            rt / "templates" / "light" / "template.pptx",
            rt / "templates" / "light" / "manifest.json",
            rt / "templates" / "light" / "bindings.json",
            rt / "templates" / "light" / "page_map.md",
            rt / "templates" / "light" / "assets",
            rt / "validate_slide_spec_gpts.py"]
    missing = [str(p.relative_to(rt)) for p in need if not p.exists()]
    if missing:
        return False, f"$RT 佈局缺:{missing}"
    forbidden = [rt / "templates" / "light" / "bindings.py",   # builtin 載體,不得再出現
                 rt / "assets", rt / "light_template.pptx"]    # 舊佈局,素材已隨包出貨
    ghosts = [str(p.relative_to(rt)) for p in forbidden if p.exists()]
    if ghosts:
        return False, f"$RT 出現不該有的檔:{ghosts}"
    env.rt = rt
    return True, "\n".join(lines) + f"\n$RT = {rt}"


def r1(env: Env):
    """examples 四份 validator 預期 exit。"""
    expect = {"01_minimal_4p": 0, "02_full_8p": 0,
              "03_advanced_unregistered_6p": 0, "04_broken_should_fail": 1}
    out = []
    for name, want in expect.items():
        rc, o = run([sys.executable, env.rt / "validate_slide_spec_gpts.py",
                     "--spec", EXAMPLES / f"{name}.json", "--asset-dir", env.rt])
        out.append(f"{name} exit={rc}(預期 {want})")
        if rc != want:
            return False, "\n".join(out) + "\n--- 輸出 ---\n" + o
    return True, "\n".join(out)


def r2(env: Env):
    """稽核閘門 + strict 全流程:(a) 原樣 02 擋 deck_name;(b) 修後 4/4 PASS。"""
    py = env.pptx_py()
    rc, o = run(py + [env.rt / "tools" / "run_pipeline.py",
                      "--spec", EXAMPLES / "02_full_8p.json",
                      "--slides", EXAMPLES / "02_full_8p.source_slides.md",
                      "--asset-dir", env.rt, "--out", env.rt / "deck02.pptx"])
    if rc == 0 or "deck_name「my_project」" not in o:
        return False, f"(a) 預期擋 deck_name(exit≠0),實得 exit={rc}\n{o}"
    fixed = env.rt / "spec02_fixed.json"
    s = json.loads((EXAMPLES / "02_full_8p.json").read_text(encoding="utf-8-sig"))
    s["deck"]["deck_name"] = "年度工作總覽"
    fixed.write_text(json.dumps(s, ensure_ascii=False, indent=2), encoding="utf-8")
    rc, o = run(py + [env.rt / "tools" / "run_pipeline.py", "--spec", fixed,
                      "--slides", EXAMPLES / "02_full_8p.source_slides.md",
                      "--asset-dir", env.rt, "--out", env.rt / "deck02.pptx"])
    if rc != 0 or "管線結果" not in o or "4/4" not in o:
        return False, f"(b) 預期 PASS(4/4 階段),實得 exit={rc}\n{o}"
    env.spec02_fixed, env.deck02 = fixed, env.rt / "deck02.pptx"
    return True, "(a) deck_name 閘門有效;(b) 修後 4/4 PASS"


def r3(env: Env):
    """QA 先 WARN 後 PASS(刪一個頁碼框,exit 仍 0)。"""
    code = f"""
from pptx import Presentation
from pptx.util import Inches
p = Presentation({str(env.deck02)!r}); done = False
for sl in p.slides:
    if done: break
    for shp in list(sl.shapes):
        if shp.has_text_frame and shp.text_frame.text.strip().isdigit() \\
           and (shp.top or 0) > Inches(6.3) and (shp.left or 0) > Inches(11.0):
            shp._element.getparent().remove(shp._element); done = True; break
p.save({str(env.rt / "deck02_warn.pptx")!r})
"""
    rc, o = run_py_snippet(env, code, need_pptx=True)
    if rc != 0:
        return False, f"改檔腳本失敗 exit={rc}\n{o}"
    rc, o = run(env.pptx_py() + [env.rt / "tools" / "qa_check.py",
                                 "--spec", env.spec02_fixed,
                                 "--pptx", env.rt / "deck02_warn.pptx"])
    if rc != 0 or "找不到右下角頁碼" not in o or "結果:PASS" not in o:
        return False, f"預期 exit=0 + 頁碼 WARN + PASS,實得 exit={rc}\n{o}"
    return True, "頁碼缺失 WARN 不擋交付,exit=0"


def r4(env: Env):
    """title 注入 → 稽核擋。"""
    s = json.loads(env.spec02_fixed.read_text(encoding="utf-8"))
    s["slides"][5]["title"] = "來源沒有的注入標題"
    inj = env.rt / "spec02_inject.json"
    inj.write_text(json.dumps(s, ensure_ascii=False), encoding="utf-8")
    rc, o = run([sys.executable, env.rt / "tools" / "audit_provenance.py",
                 "--spec", inj, "--slides", EXAMPLES / "02_full_8p.source_slides.md"])
    if rc != 1 or "未逐字出現在該頁來源區塊" not in o:
        return False, f"預期 exit=1 + 逐字錯誤,實得 exit={rc}\n{o}"
    return True, "注入標題被稽核硬擋"


def r5(env: Env):
    """精確數字 token(95%→9%)→ 稽核擋。"""
    raw = env.spec02_fixed.read_text(encoding="utf-8").replace('"95%"', '"9%"')
    tok = env.rt / "spec02_token.json"
    tok.write_text(raw, encoding="utf-8")
    rc, o = run([sys.executable, env.rt / "tools" / "audit_provenance.py",
                 "--spec", tok, "--slides", EXAMPLES / "02_full_8p.source_slides.md"])
    if rc != 1 or "無精確對應" not in o:
        return False, f"預期 exit=1 + token 錯誤,實得 exit={rc}\n{o}"
    return True, "數字 token 竄改被稽核硬擋"


def r6(env: Env):
    """fixture 05 純淨度(真未切頁、無頁型指示)。"""
    text = (EXAMPLES / "05_outline_to_ppt_source.md").read_text(encoding="utf-8")
    n = len(re.findall(r"## Slide|page_type", text))
    return (n == 0), f"命中 {n} 處(預期 0)"


def r7(env: Env):
    """Knowledge 清單與 archive hash——基準解析自 REGRESSION.md,漂移即紅。"""
    md = REG_MD.read_text(encoding="utf-8")
    sec = md.split("## R7", 1)[1].split("\n## ", 1)[0]
    baselines = {name: sha for sha, name
                 in re.findall(r"^([0-9a-f]{64})\s+(\S+)$", sec, re.M)}
    counts = re.findall(r"#\s*預期\s*(\d+)", sec)
    if len(baselines) != 2 or len(counts) < 2:
        return False, "無法從 REGRESSION.md R7 解析基準(sha×2 + 預期檔數×2)"
    lines = []
    n_rules = len([p for p in RULES.iterdir() if p.name != "__pycache__"])
    n_dist = len(list((REPO / "gpts" / "dist").iterdir()))
    for got, want, label in ((n_rules, int(counts[0]), "engine/rules 散檔"),
                             (n_dist, int(counts[1]), "gpts/dist zip")):
        lines.append(f"{label}:{got}(基準 {want})")
        if got != want:
            return False, "\n".join(lines) + "\n檔數與 R7 基準不符(加了檔要更新基準與 Knowledge 清單)"
    for zname, want in baselines.items():
        got = sha256_file(REPO / "gpts" / "dist" / zname)
        lines.append(f"{zname} sha {'一致' if got == want else '不符'}")
        if got != want:
            return False, ("\n".join(lines) + f"\n{zname} 實際 {got}\n"
                           "zip 內容變了但 R7 基準沒更新(或 zip 忘了重打)")
    return True, "\n".join(lines)


def r8(env: Env):
    """直供 JSON 模式全流程(追溯關)。"""
    spec = env.rt / "spec01.json"
    shutil.copy2(EXAMPLES / "01_minimal_4p.json", spec)
    rc, o = run(env.pptx_py() + [env.rt / "tools" / "run_pipeline.py", "--spec", spec,
                                 "--asset-dir", env.rt, "--out", env.rt / "deck01.pptx"])
    if rc != 0 or "來源追溯" not in o or "3/3" not in o:
        return False, f"預期 exit=0 + 追溯關 + 3/3 PASS,實得 exit={rc}\n{o}"
    return True, "直供模式 3/3 PASS(來源追溯關)"


def r9(env: Env):
    """多模板:lightcopy 端到端註冊演練 + 同 spec 換包產兩份。"""
    py = env.pptx_py()
    pr = env.mktemp("regress_pr_")
    rc, o = run(py + [ADMIN, "new", "--pptx", LIGHT / "template.pptx",
                      "--id", "lightcopy", "--name", "淺色複本測試", "--packs-root", pr])
    if rc != 0:
        return False, f"new 失敗 exit={rc}\n{o}"
    src = json.loads((LIGHT / "manifest.json").read_text(encoding="utf-8-sig"))
    mp = pr / "lightcopy" / "manifest.json"
    m = json.loads(mp.read_text(encoding="utf-8-sig"))
    m["style"], m["asset_defaults"], m["page_number"] = (
        src["style"], src["asset_defaults"], src["page_number"])
    m["page_types"] = {pt: e for pt, e in src["page_types"].items()
                       if e["mode"] == "fill"}
    # capacity_overrides 是版位量測結果,不是可選設定——漏抄它 golden 必 FAIL
    # (會用契約上限產 max 變體,遠超版位實際容量;R9 曾因此紅 11 個 commit)
    m["capacity_overrides"] = src["capacity_overrides"]
    mp.write_text(json.dumps(m, ensure_ascii=False, indent=2), encoding="utf-8")
    shutil.copy2(LIGHT / "bindings.json", pr / "lightcopy" / "bindings.json")
    shutil.copytree(LIGHT / "assets_src", pr / "lightcopy" / "assets_src",
                    dirs_exist_ok=True)
    for step in ("freeze", "register"):
        rc, o = run(py + [ADMIN, step, "--id", "lightcopy", "--packs-root", pr])
        if rc != 0:
            return False, f"{step} 失敗 exit={rc}\n{o}"
    # 同一份 spec 換 deck.template 產兩份,各自 render+qa
    spec_l = EXAMPLES / "01_minimal_4p.json"
    s = json.loads(spec_l.read_text(encoding="utf-8-sig"))
    s["deck"]["template"] = "lightcopy"
    spec_lc = pr / "spec01_lightcopy.json"
    spec_lc.write_text(json.dumps(s, ensure_ascii=False, indent=2), encoding="utf-8")
    for spec, extra, tag in ((spec_l, [], "模板包:light@"),
                             (spec_lc, ["--packs-root", pr], "模板包:lightcopy@")):
        deck = pr / f"deck_{Path(spec).stem}.pptx"
        rc, o = run(py + [TOOLS / "render_deck.py", "--spec", spec,
                          "--asset-dir", env.rt, "--out", deck] + extra)
        if rc != 0 or tag not in o:
            return False, f"render({tag})exit={rc},預期 0 且含「{tag}」\n{o}"
        rc, o = run(py + [TOOLS / "qa_check.py", "--spec", spec, "--pptx", deck] + extra)
        if rc != 0:
            return False, f"qa({tag})exit={rc}\n{o}"
    return True, "register exit=0;同 spec 以 light / lightcopy 各 render+qa PASS"


def r10(env: Env):
    """全包 lint + bindings.py 後門反向測試(兩道防線都必須紅)。"""
    rc, o = run([sys.executable, ADMIN, "lint", "--all"])
    if rc != 0:
        return False, f"lint --all 預期 0,實得 {rc}\n{o}"
    planted = LIGHT / "bindings.py"
    if planted.exists():
        return False, f"{planted} 已存在?工作區不乾淨,先處理再跑"
    try:
        planted.write_text("BUILDERS = {}\n", encoding="utf-8")
        rc1, o1 = run([sys.executable, ADMIN, "lint", "--id", "light"])
        rc2, o2 = run(env.pptx_py() + [TOOLS / "render_deck.py",
                                       "--spec", EXAMPLES / "01_minimal_4p.json",
                                       "--asset-dir", LIGHT / "assets_src",
                                       "--out", env.rt / "_guard.pptx"])
    finally:
        planted.unlink(missing_ok=True)
    if rc1 != 1 or "包內不得有 bindings.py" not in o1:
        return False, f"lint 防線失效:exit={rc1}(預期 1)\n{o1}"
    if rc2 != 2 or "內有 bindings.py" not in o2:
        return False, f"載入期防線失效:exit={rc2}(預期 2)\n{o2}"
    return True, "lint --all 綠;種入 bindings.py 後 lint=1、render=2(後門仍封死)"


def r11(env: Env):
    """golden 契約快照與現行契約同步。失敗時刻意不還原——留下的 diff 即修法。"""
    rc, o = run(["git", "status", "--porcelain", "--", "engine/golden"])
    if o.strip():
        return False, f"engine/golden/ 工作區不乾淨,無法判定同步:\n{o}"
    rc, o = run([sys.executable, ADMIN, "golden", "--regen-specs"])
    if rc != 0:
        return False, f"--regen-specs exit={rc}\n{o}"
    rc, o = run(["git", "status", "--porcelain", "--", "engine/golden"])
    if o.strip():
        return False, ("契約改了但快照沒重派生(把以下重派生結果一起 commit):\n" + o)
    return True, "快照與 PAGE_TYPES 同步"


def r12(env: Env):
    """誠實容量:contract 上限不得觸發縮字/侵入鄰欄;含碰撞偵測反向測試。"""
    py = env.pptx_py()
    rc, o = run(py + [ADMIN, "golden", "--id", "light"])
    if rc != 0:
        return False, f"golden --id light exit={rc}\n{o}"
    verify = """
import sys, json; sys.path.insert(0,'engine/tools'); sys.path.insert(0,'engine/rules')
from pptx import Presentation
import text_tools as tt
m = json.load(open('engine/templates/light/manifest.json'))
fills = list(json.load(open('engine/templates/light/bindings.json'))["fills"])
order = [pt for pt in m["page_types"] if pt in fills]
seq = [pt for pt in order for _ in ("min", "max")]
page_of = {pt: e["template_page"] for pt, e in m["page_types"].items()
           if e.get("mode") == "fill"}
tpl = Presentation('engine/templates/light/template.pptx')
g = Presentation('ppt_out/golden_light.pptx')
assert len(g.slides) == len(seq), f"頁數不符:golden {len(g.slides)} vs 預期 {len(seq)}"
T = {pt: {s.shape_id: s for s in tt.iter_text_shapes(tpl.slides[pg-1].shapes)}
     for pt, pg in page_of.items()}
BASE = {pt: {frozenset((a.shape_id, b.shape_id)): ar
             for ar, a, b in tt.text_collisions(tpl.slides[pg-1])}
        for pt, pg in page_of.items()}
shr = worse = 0
for i, sl in enumerate(g.slides, 1):
    pt = seq[i-1]
    for s in tt.iter_text_shapes(sl.shapes):
        o = T[pt].get(s.shape_id)
        if o is not None and tt._first_run_size_pt(s) < tt._first_run_size_pt(o) - 0.5:
            shr += 1; print(f"  縮字 p{i} {pt} id{s.shape_id}")
    for ar, a, b in tt.text_collisions(sl):
        ref = BASE[pt].get(frozenset((a.shape_id, b.shape_id)), 0.0)
        if ar > max(ref * 1.10, 0.03):
            worse += 1; print(f"  侵入 p{i} {pt} id{a.shape_id}x{b.shape_id} {ar:.2f}in2")
print(f"RESULT shr={shr} worse={worse}")
"""
    rc, o = run_py_snippet(env, verify, need_pptx=True)
    mres = re.search(r"RESULT shr=(\d+) worse=(\d+)", o)
    if rc != 0 or not mres or mres.group(1) != "0" or mres.group(2) != "0":
        return False, f"縮字/侵入檢查未歸零(exit={rc}):\n{o}"
    # 反向測試:刻意塞超長文字直接渲染,qa 必須 FAIL——exit=0 代表偵測器壞了
    s = json.loads((REPO / "ppt_out" / "golden_light.spec.json").read_text(encoding="utf-8"))
    for sl in s["slides"]:
        if sl["page_type"] == "data_three_number_kpis":
            for k in sl.get("slots", {}).get("kpis", []):
                k["detail"] = "刻意超長的說明文字用來驗證碰撞偵測會不會擋下來真的很長"
    s["deck"].pop("template", None)
    neg = REPO / "ppt_out" / "_neg.json"
    neg.write_text(json.dumps(s, ensure_ascii=False, indent=2), encoding="utf-8")
    rc, o = run(py + [TOOLS / "render_deck.py", "--spec", neg, "--template-pack", LIGHT,
                      "--asset-dir", LIGHT / "assets_src",
                      "--out", REPO / "ppt_out" / "_neg.pptx"])
    if rc != 0:
        return False, f"反向測試 render exit={rc}\n{o}"
    rc, o = run(py + [TOOLS / "qa_check.py", "--spec", neg,
                      "--pptx", REPO / "ppt_out" / "_neg.pptx", "--template-pack", LIGHT])
    if rc != 1 or "文字壓到別的元素" not in o:
        return False, (f"碰撞偵測器疑似壞了:qa exit={rc}(預期 1 + [F] 碰撞)\n{o}\n"
                       "偵測器壞掉比破版更嚴重——後續驗收都會假綠")
    return True, "golden PASS、縮字/侵入歸零;反向測試 qa 正確 FAIL"


def r13(env: Env):
    """sync-docs 派生檔同步 + pack 閘門反向測試(INDEX 改壞要打不出 zip)。"""
    rc, o = run([sys.executable, ADMIN, "sync-docs"])
    if rc != 0 or "派生檔與機器真相一致" not in o:
        return False, f"sync-docs exit={rc}(預期 0)\n{o}"
    idx = ENGINE / "templates" / "INDEX.md"
    orig = idx.read_text(encoding="utf-8")
    if "registered" not in orig:
        return False, "INDEX.md 無 registered 字樣,反向測試前置不成立"
    tools_zip = REPO / "gpts" / "dist" / "tools.zip"
    sha_before = sha256_file(tools_zip)
    try:
        idx.write_text(orig.replace("registered", "draft", 1), encoding="utf-8")
        rc, o = run([sys.executable, ADMIN, "pack", "--tools"])
    finally:
        idx.write_text(orig, encoding="utf-8")
        if sha256_file(tools_zip) != sha_before:
            run(["git", "checkout", "--", "gpts/dist/tools.zip"])
    if rc != 1:
        return False, f"pack 閘門失效:INDEX 漂移下 pack exit={rc}(預期 1)\n{o}"
    if sha256_file(tools_zip) != sha_before:
        return False, "pack 在漂移下仍改寫了 tools.zip(已從 git 還原)"
    return True, "sync-docs 綠;INDEX 漂移時 pack 拒打包且 zip 未被改寫"


def rl0(env: Env):
    """light 包完整性(sha 對 manifest、page_types 自洽)。"""
    m = json.loads((LIGHT / "manifest.json").read_text(encoding="utf-8-sig"))
    sha = sha256_file(LIGHT / "template.pptx")
    if sha != m["template_sha256"]:
        return False, "template.pptx sha 與 manifest 不符:模板改版未重跑盤點(TEMPLATE_LIFECYCLE.md)"
    from collections import Counter
    c = Counter(v.get("mode") for v in m["page_types"].values())
    if len(m["page_types"]) != sum(c.values()):
        return False, f"page_types 自洽檢查不符:{dict(c)}"
    return True, f"sha OK;page_types {len(m['page_types'])} 筆 {dict(c)}"


def rl1(env: Env):
    """light smoke spec 直供模式全流程。"""
    rc, o = run(env.pptx_py() + [env.rt / "tools" / "run_pipeline.py",
                                 "--spec", LIGHT / "examples" / "smoke_spec.json",
                                 "--asset-dir", env.rt,
                                 "--out", env.rt / "deck_light_smoke.pptx"])
    if rc != 0 or "3/3" not in o:
        return False, f"預期 exit=0 + 3/3 PASS,實得 exit={rc}\n{o}"
    return True, "smoke spec 3/3 PASS"


CASES = [
    # (id, 標題, fn, needs)
    ("R0", "archive 完整性與環境建置", r0, []),
    ("R1", "examples 四份 validator 預期 exit", r1, ["R0"]),
    ("R2", "稽核閘門 + strict 全流程", r2, ["R0"]),
    ("R3", "QA 先 WARN 後 PASS", r3, ["R2"]),
    ("R4", "title 注入 → 稽核擋", r4, ["R2"]),
    ("R5", "精確數字 token → 稽核擋", r5, ["R2"]),
    ("R6", "fixture 05 純淨度", r6, []),
    ("R7", "Knowledge 清單與 archive hash", r7, []),
    ("R8", "直供 JSON 模式全流程", r8, ["R0"]),
    ("R9", "多模板 lightcopy 註冊演練", r9, ["R0"]),
    ("R10", "全包 lint + bindings.py 反向測試", r10, ["R0"]),
    ("R11", "golden 契約快照同步", r11, []),
    ("R12", "誠實容量 + 碰撞反向測試", r12, []),
    ("R13", "sync-docs + pack 閘門反向測試", r13, []),
    ("R-L0", "light 包完整性", rl0, []),
    ("R-L1", "light smoke spec 全流程", rl1, ["R0"]),
    # R-L2(clone 抽測)是人工案例,見 light REGRESSION.md——不自動化
]


def expand_only(only: set[str]) -> set[str]:
    """--only 自動帶入相依案例(R3 需要 R2 的產物、R2 需要 R0 的 $RT)。"""
    by_id = {cid: needs for cid, _, _, needs in CASES}
    out, stack = set(), list(only)
    while stack:
        cid = stack.pop()
        if cid in out:
            continue
        out.add(cid)
        stack.extend(by_id.get(cid, []))
    return out


def main(argv):
    ap = argparse.ArgumentParser(description="引擎級回歸 runner(案例規格見 engine/REGRESSION.md)")
    ap.add_argument("--only", help="逗號分隔案例 id(自動帶入相依案例)")
    ap.add_argument("--skip", help="逗號分隔案例 id(列為 SKIP,結果視為未全綠)")
    ap.add_argument("--list", action="store_true")
    args = ap.parse_args(argv)

    if args.list:
        for cid, title, _, needs in CASES:
            dep = f"(需 {','.join(needs)})" if needs else ""
            print(f"{cid:<6} {title}{dep}")
        print("R-L2   clone 抽測——人工案例,見 engine/templates/light/REGRESSION.md")
        return 0

    only = expand_only({t.strip() for t in args.only.split(",")}) if args.only else None
    skip = {t.strip() for t in args.skip.split(",")} if args.skip else set()
    env = Env()
    results = []  # (cid, title, status, secs, detail)
    for cid, title, fn, needs in CASES:
        if only is not None and cid not in only:
            continue
        if cid in skip or any(s in skip for s in needs):
            results.append((cid, title, "SKIP", 0.0, "手動跳過(或相依案例被跳過)"))
            print(f"— {cid} {title}:SKIP")
            continue
        if any(r[0] in needs and r[2] != "PASS" for r in results):
            results.append((cid, title, "SKIP", 0.0, "相依案例未過"))
            print(f"— {cid} {title}:SKIP(相依案例未過)")
            continue
        t0 = time.time()
        try:
            ok, detail = fn(env)
        except Exception as e:  # noqa: BLE001 — runner 不得因單一案例炸掉整輪
            ok, detail = False, f"案例執行例外:{type(e).__name__}: {e}"
        secs = time.time() - t0
        status = "PASS" if ok else "FAIL"
        mark = "✓" if ok else "✗"
        print(f"{mark} {cid} {title}({secs:.1f}s)")
        if not ok:
            print("  " + "\n  ".join(detail.strip().splitlines()))
        results.append((cid, title, status, secs, detail))

    n_pass = sum(1 for r in results if r[2] == "PASS")
    n_fail = sum(1 for r in results if r[2] == "FAIL")
    n_skip = sum(1 for r in results if r[2] == "SKIP")
    print(f"\n結果:{n_pass} 綠 / {n_fail} 紅 / {n_skip} 跳過(共 {len(results)})")
    if n_fail:
        print("紅燈案例:" + ", ".join(r[0] for r in results if r[2] == "FAIL"))
    if n_fail == 0 and n_skip == 0 and env.rt is not None:
        for d in env.tempdirs:
            shutil.rmtree(d, ignore_errors=True)
    elif env.rt is not None:
        print(f"(暫存目錄保留供診斷:{env.rt})")
    return 1 if (n_fail or n_skip) else 0


if __name__ == "__main__":
    sys.exit(main(sys.argv[1:]))
