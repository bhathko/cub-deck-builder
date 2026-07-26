#!/usr/bin/env python3
"""text_tools — 文字框的搜尋、替換(保留樣式)與 CJK 溢出估算。

設計重點:
- 替換文字時 deepcopy 原本第一個 run 的 rPr / 段落的 pPr,字級、顏色、粗細、
  對齊全部保留 → 「模板改字」不會把樣式改跑。
- 沙箱沒有中文字體,無法真量測;estimate_overflow 用「CJK 字寬 ≈ 1em、
  半形 ≈ 0.55em、行高 ≈ 1.3em」啟發式估算,寬鬆容忍(超過 8% 才判溢出)。
- 所有函式都吃 shape 物件,群組請先用 iter_text_shapes 展開。
"""
from __future__ import annotations

import copy
import math

from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.util import Pt

EMU_PER_PT = 12700
EMU_PER_IN = 914400
DEFAULT_FONT_PT = 18.0
LINE_SPACING = 1.3
OVERFLOW_TOLERANCE = 1.08  # 估算誤差容忍


def iter_text_shapes(shapes):
    """遞迴走訪(含群組內層),yield 所有帶 text_frame 的 shape。"""
    for shp in shapes:
        if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_text_shapes(shp.shapes)
        elif getattr(shp, "has_text_frame", False):
            yield shp


def shape_text(shape) -> str:
    return shape.text_frame.text if getattr(shape, "has_text_frame", False) else ""


def find_text_shapes(slide, contains=None, name=None, shape_id=None):
    """依條件找文字框,回傳 list。條件可複合(AND)。"""
    out = []
    for shp in iter_text_shapes(slide.shapes):
        if shape_id is not None and shp.shape_id != shape_id:
            continue
        if name is not None and shp.name != name:
            continue
        if contains is not None and contains not in shape_text(shp):
            continue
        out.append(shp)
    return out


def set_text_keep_style(shape, text) -> None:
    """替換整個 text_frame 的文字,沿用原第一段/第一 run 的樣式。支援 \\n 多行。"""
    tf = shape.text_frame
    first_p = tf.paragraphs[0]._p
    pPr_tpl = first_p.find(qn("a:pPr"))
    rPr_tpl = None
    first_r = first_p.find(qn("a:r"))
    if first_r is not None:
        rPr_tpl = first_r.find(qn("a:rPr"))
    # 沒 run 時退而求其次:用 endParaRPr 當樣式模板
    if rPr_tpl is None:
        end_pr = first_p.find(qn("a:endParaRPr"))
        if end_pr is not None:
            rPr_tpl = copy.deepcopy(end_pr)
            rPr_tpl.tag = qn("a:rPr")

    tf.clear()
    lines = str(text).split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if pPr_tpl is not None:
            old = p._p.find(qn("a:pPr"))
            if old is not None:
                p._p.remove(old)
            p._p.insert(0, copy.deepcopy(pPr_tpl))
        run = p.add_run()
        run.text = line
        if rPr_tpl is not None:
            run._r.insert(list(run._r).index(run._r.find(qn("a:t"))), copy.deepcopy(rPr_tpl))


# 數學英數符號區(U+1D400–U+1D7FF):粗體/花體的英數字,字面像 CJK 但字寬同半形。
# 模板用它排步驟編號徽章(如 p35 的 𝟭𝟮𝟯𝟰),誤判成全形會讓 0.34 吋的框假性溢出。
_NARROW_RANGES = ((0x1D400, 0x1D7FF),)


def _char_units(s: str) -> float:
    """以 em 為單位的估算寬度:CJK≈1,半形與數學英數符號≈0.55。"""
    def unit(c):
        o = ord(c)
        if any(lo <= o <= hi for lo, hi in _NARROW_RANGES):
            return 0.55
        return 1.0 if o > 0x2E80 else 0.55
    return sum(unit(c) for c in s)


def _first_run_size_pt(shape) -> float:
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.size is not None:
                return run.font.size.pt
    return DEFAULT_FONT_PT


def estimate_overflow(shape, text=None, size_pt=None):
    """估算 text(預設取現有文字)在 shape 內是否放得下。

    回傳 dict:{fits, lines, needed_pt, avail_pt, size_pt}
    """
    tf = shape.text_frame
    text = shape_text(shape) if text is None else str(text)
    size_pt = size_pt or _first_run_size_pt(shape)

    # 直排文字(vert 屬性)無法用橫排邏輯估算,一律視為放得下
    if tf._bodyPr.get("vert") not in (None, "horz"):
        return {"fits": True, "lines": 1, "needed_pt": 0.0, "avail_pt": 0.0,
                "size_pt": size_pt}

    ml = tf.margin_left if tf.margin_left is not None else 91440
    mr = tf.margin_right if tf.margin_right is not None else 91440
    mt = tf.margin_top if tf.margin_top is not None else 45720
    mb = tf.margin_bottom if tf.margin_bottom is not None else 45720

    avail_w_pt = max((shape.width - ml - mr) / EMU_PER_PT, 1.0)
    avail_h_pt = max((shape.height - mt - mb) / EMU_PER_PT, 1.0)

    cap_em = avail_w_pt / size_pt  # 一行放得下幾個 em
    # 窄框(一行約一個字)是刻意的直式堆疊設計(如「優點」側標、步驟編號徽章)。
    # **不能因此就當作一定放得下**——舊版在這裡無條件 return fits=True,結果
    # 0.34 吋的編號框塞「01」會折成兩行破版,偵測器卻完全看不到(偽陰性)。
    # 正確做法是照樣算行數(窄框每行約一字),只是由框高決定放不放得下。
    # wrap="none" = 框設定為不自動換行,文字只會橫向超出框線,不會折行堆高。
    # 這類框(模板用於年份標籤、英文徽章)行數固定等於硬換行數,不能按框寬折算,
    # 否則短字串在窄框會被誤判成多行溢出。
    no_wrap = tf._bodyPr.get("wrap") == "none"
    lines = 0
    for hard_line in text.split("\n"):
        units = _char_units(hard_line)
        if no_wrap:
            lines += 1
            continue
        # 折行數上限 = 字元數:窄到一行放不滿一個字的框(直排側標「優點」),
        # units/cap_em 會算出比字數還多的行數。一個字最多佔一行。
        lines += max(1, min(math.ceil(units / max(cap_em, 0.1)), len(hard_line) or 1))
    needed_pt = lines * size_pt * LINE_SPACING
    # 橫向:wrap="none" 的框不會折行,文字直接往旁邊長出去壓到鄰居
    # (light p30 的 name 徽章溢出 1.58 吋壓到隔欄說明、p22 三段 detail 互疊)。
    # 會折行的框沒有橫向問題(超出就換行),故只對 no_wrap 判斷。
    needed_w_pt = max((_char_units(l) for l in text.split("\n")), default=0.0) * size_pt
    fits_w = (not no_wrap) or needed_w_pt <= avail_w_pt * OVERFLOW_TOLERANCE
    # 單行文字在 PowerPoint 不會被裁切(框高偏小只是視覺貼邊),不算縱向溢出;
    # 真正的問題是換行後行數超出框高 → 文字溢到卡片外。
    fits_h = lines <= 1 or needed_pt <= avail_h_pt * OVERFLOW_TOLERANCE
    return {
        "fits": fits_h and fits_w,
        "fits_h": fits_h,
        "fits_w": fits_w,
        "lines": lines,
        "needed_pt": round(needed_pt, 1),
        "avail_pt": round(avail_h_pt, 1),
        "needed_w_pt": round(needed_w_pt, 1),
        "avail_w_pt": round(avail_w_pt, 1),
        "size_pt": size_pt,
    }


def has_autofit(shape) -> bool:
    """框是否設了「自動調整」:spAutoFit=框長高貼合文字、normAutofit=PowerPoint 自己縮字。

    這兩種框由 PowerPoint 自行消化溢出,我們不該再動字級——模板本身就靠它排版
    (light p47 的說明框原文就要 3 行、比框高多 56%)。字級是設計過的,塞不下
    要改稿或換版面,不是偷偷縮小一號。容量由 capacity_overrides 在閘門擋住。
    """
    bp = shape.text_frame._bodyPr
    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    return bp.find(ns + "spAutoFit") is not None or bp.find(ns + "normAutofit") is not None


def shrink_to_fit(shape, min_pt: float = 14.0) -> float:
    """溢出時逐階(-2pt)縮小全部 run 字級,直到放得下或到 min_pt。回傳最終字級。

    autofit 框跳過(見 has_autofit)。
    """
    if has_autofit(shape):
        return _first_run_size_pt(shape)
    size = _first_run_size_pt(shape)
    while size > min_pt and not estimate_overflow(shape, size_pt=size)["fits"]:
        size -= 2
    if size != _first_run_size_pt(shape):
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(max(size, min_pt))
    return max(size, min_pt)

# ---------------------------------------------------------------------------
# 絕對座標與文字佔用範圍(碰撞偵測用)
#
# 為什麼需要:「文字裝得進自己的框」不等於「版面沒壞」。實際的破版是文字跑到
# 不該出現的位置——wrap="none" 的框往右長、autofit 的框往下長,結果壓到鄰欄。
# 設計師目檢 light golden 時就是這樣抓到 p10/p22/p30 的(2026-07-26)。
# 群組內子形狀的 left/top 是「群組子座標系」的值,必須用 chOff/chExt 換算,
# 否則會算出不可能的跨欄重疊。
# ---------------------------------------------------------------------------
_NS_A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def walk_absolute(shapes, tf=(0.0, 0.0, 1.0, 1.0)):
    """展平所有非群組形狀,回傳 (shape, left, top, width, height),單位吋(絕對座標)。"""
    dx, dy, sx, sy = tf
    for s in shapes:
        L = (s.left or 0) / EMU_PER_IN * sx + dx
        T = (s.top or 0) / EMU_PER_IN * sy + dy
        W = (s.width or 0) / EMU_PER_IN * sx
        H = (s.height or 0) / EMU_PER_IN * sy
        if s.shape_type == MSO_SHAPE_TYPE.GROUP:
            xfrm = s._element.find(f".//{_NS_A}xfrm")
            nsx = nsy = 1.0
            cdx = cdy = 0.0
            if xfrm is not None:
                ext, ch_ext = xfrm.find(f"{_NS_A}ext"), xfrm.find(f"{_NS_A}chExt")
                off, ch_off = xfrm.find(f"{_NS_A}off"), xfrm.find(f"{_NS_A}chOff")
                if ext is not None and ch_ext is not None:
                    nsx = int(ext.get("cx")) / (int(ch_ext.get("cx")) or 1)
                    nsy = int(ext.get("cy")) / (int(ch_ext.get("cy")) or 1)
                if off is not None and ch_off is not None:
                    cdx = (int(off.get("x")) - int(ch_off.get("x")) * nsx) / EMU_PER_IN
                    cdy = (int(off.get("y")) - int(ch_off.get("y")) * nsy) / EMU_PER_IN
            yield from walk_absolute(s.shapes, (dx + cdx * sx, dy + cdy * sy, sx * nsx, sy * nsy))
        else:
            yield s, L, T, W, H


def text_footprint(shape, left, top, width, height):
    """文字實際佔用的矩形 (x0, y0, x1, y1),單位吋;無文字回傳 None。

    可能大於形狀本身:wrap="none" 往旁邊長、行數超出框高往下長(autofit)。
    """
    tf = shape.text_frame
    text = tf.text
    if not text.strip():
        return None
    ml = (tf.margin_left if tf.margin_left is not None else 91440) / EMU_PER_IN
    mr = (tf.margin_right if tf.margin_right is not None else 91440) / EMU_PER_IN
    mt = (tf.margin_top if tf.margin_top is not None else 45720) / EMU_PER_IN
    mb = (tf.margin_bottom if tf.margin_bottom is not None else 45720) / EMU_PER_IN
    size = _first_run_size_pt(shape)
    avail_w = max(width - ml - mr, 0.05)
    avail_h = max(height - mt - mb, 0.05)
    no_wrap = tf._bodyPr.get("wrap") == "none"
    lines = 0
    widest = 0.0
    for hard_line in text.split("\n"):
        u = _char_units(hard_line) * size / 72.0
        widest = max(widest, u)
        lines += 1 if no_wrap else max(1, min(int(-(-u // avail_w)) if avail_w else 1,
                                             len(hard_line) or 1))
    w = max(widest, avail_w) if no_wrap else avail_w
    h = max(lines * size * LINE_SPACING / 72.0, avail_h)
    align = str(tf.paragraphs[0].alignment) if (tf.paragraphs and tf.paragraphs[0].alignment) else ""
    x0 = left + ml
    if "CENTER" in align:
        x0 -= (w - avail_w) / 2
    elif "RIGHT" in align:
        x0 -= (w - avail_w)
    return (x0, top + mt, x0 + w, top + mt + h)


def overlap_area(a, b) -> float:
    w = min(a[2], b[2]) - max(a[0], b[0])
    h = min(a[3], b[3]) - max(a[1], b[1])
    return w * h if (w > 0 and h > 0) else 0.0


def text_collisions(slide, min_area_in2: float = 0.03):
    """回傳同一頁上文字佔用範圍互相重疊的配對:[(面積, shapeA, shapeB), …] 面積降冪。"""
    items = []
    for s, L, T, W, H in walk_absolute(slide.shapes):
        if not s.has_text_frame:
            continue
        fp = text_footprint(s, L, T, W, H)
        if fp:
            items.append((s, fp))
    hits = []
    for i in range(len(items)):
        for j in range(i + 1, len(items)):
            a = overlap_area(items[i][1], items[j][1])
            if a > min_area_in2:
                hits.append((a, items[i][0], items[j][0]))
    return sorted(hits, key=lambda x: -x[0])
