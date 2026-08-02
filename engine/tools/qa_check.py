#!/usr/bin/env python3
"""qa_check — 產檔後自檢:比對產出 pptx 與 slide_spec.json。

驗證器(validate_slide_spec_gpts.py)擋「產檔前的 spec」;本工具擋「產檔後的
pptx」。兩道都過,才交付給使用者。

用法:
  python qa_check.py --spec /mnt/data/slide_spec.json --pptx /mnt/data/deck.pptx

檢查項目:
  FAIL(必修,exit 1):
    - 頁數與 spec / deck.slide_count 不符
    - spec 槽位文字沒有出現在對應頁(內容遺漏或沒替換到)
    - 殘留 PowerPoint Section
    - 文字壓到別的元素:文字實際佔用範圍互相重疊,判準是**不得比模板原本
      更侵入鄰欄**(模板有幾處刻意疊在一起,絕對零重疊會誤殺)。
      「裝得進自己的框」不等於版面沒壞——wrap="none" 的框往右長、
      autofit 的框往下長,都會壓到鄰欄
  WARN(人工判斷):
    - 字體不在該包 style.allowed_fonts 的 run
    - 頁碼缺漏或不該有頁碼的頁出現頁碼
    - 文字溢出疑慮(CJK 啟發式估算;先報總量與受影響頁,再列最嚴重 5 條,
      完整清單加 --overflow-all)

輸出精簡:只印問題;全部通過就一行 PASS。修法一律是改 render_plan.json /
slide_spec.json 後重跑 render_deck(整檔重生),不要手改 pptx。
"""
from __future__ import annotations

import argparse
import json
import sys
import unicodedata
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches

import text_tools as tt

import pack_loader

# 預設白名單 = light 包值(解不到模板包時的內建後備);正常路徑改讀選定包
# manifest 的 style.allowed_fonts。'+' 開頭為佈景主題字型參照(+mn-ea 等),
# 由模板主題解析,一律放行。微軟正黑體 = JhengHei 中文名。
ALLOWED_FONTS = {"Microsoft JhengHei", "微軟正黑體", "Helvetica", "Noto Sans TC"}
# 頁碼偵測窗預設(吋;light 幾何)——解不到包時後備
DETECT_ZONE = {"left": 11.0, "top": 6.3}


def _font_ok(name: str, allowed=ALLOWED_FONTS) -> bool:
    return name in allowed or name.startswith("+")


def norm(s: str) -> str:
    s = unicodedata.normalize("NFKC", str(s))
    return "".join(ch for ch in s if not ch.isspace())


def collect_strings(val, out):
    """遞迴收集 slots 內所有字串(檢查它們必須出現在產出頁上)。"""
    if isinstance(val, str):
        if val.strip():
            out.append(val)
    elif isinstance(val, list):
        for v in val:
            collect_strings(v, out)
    elif isinstance(val, dict):
        for v in val.values():
            collect_strings(v, out)


def slide_all_text(slide) -> str:
    parts = []
    for shp in tt.iter_text_shapes(slide.shapes):
        parts.append(tt.shape_text(shp))
    for shp in slide.shapes:  # 表格文字 + 圖表數據(chart op 填入者)
        if getattr(shp, "has_table", False):
            for row in shp.table.rows:
                for cell in row.cells:
                    parts.append(cell.text)
        if getattr(shp, "has_chart", False):
            # 讀 chart XML 的 <c:v> 原文(系列名/類別/數值的快取字面值)——
            # 走 plots API 會把數值轉 float,尾零消失("22.0"→22.0→"22")
            # 導致 spec 字串比對失敗
            parts.extend(v.text or "" for v in shp.chart._chartSpace.iter(qn("c:v")))
    return norm("".join(parts))


def has_sections(prs) -> bool:
    for extLst in prs._element.findall(qn("p:extLst")):
        for ext in extLst.iter():
            if ext.tag.endswith("}sectionLst"):
                return True
    return False


def main(argv):
    ap = argparse.ArgumentParser()
    ap.add_argument("--spec", required=True)
    ap.add_argument("--pptx", required=True)
    ap.add_argument("--template-pack", help="模板包 id 或目錄(省略=spec deck.template→light)")
    ap.add_argument("--packs-root", help="模板包根目錄(預設=tools 上層的 templates/)")
    ap.add_argument("--overflow-all", action="store_true",
                    help="列出全部溢出疑慮(預設只列最嚴重 5 條 + 總量)")
    args = ap.parse_args(argv)

    spec = json.loads(Path(args.spec).read_text(encoding="utf-8-sig"))  # 容忍 BOM
    # 模板包感知:字體白名單與頁碼偵測窗讀選定包 manifest;明確指定卻載不到
    # → exit 2;預設 light 載不到 → 內建常數後備(舊部署相容)
    allowed_fonts, detect_zone, pack_note = ALLOWED_FONTS, DETECT_ZONE, ""
    pack = None  # 預設包載不到時走內建常數後備;先綁定,否則後備路徑 UnboundLocalError
    try:
        pack = pack_loader.load_pack(pack_arg=args.template_pack,
                                     spec_deck=spec.get("deck"),
                                     packs_root=args.packs_root,
                                     load_bindings=False)
        allowed_fonts = set((pack.manifest.get("style") or {}).get("allowed_fonts")
                            or ALLOWED_FONTS)
        detect_zone = (pack.manifest.get("page_number") or {}).get("detect_zone_in", DETECT_ZONE)
        pack_note = f"{pack.id}@{pack.version}"
    except pack_loader.PackError as e:
        explicit = args.template_pack or (spec.get("deck") or {}).get("template")
        if explicit:
            print(f"✗ 模板包載入失敗:{e}")
            return 2
    prs = Presentation(args.pptx)
    fails, warns = [], []
    slide_w = (prs.slide_width or 0) / 914400
    slide_h = (prs.slide_height or 0) / 914400

    spec_slides = sorted(spec["slides"], key=lambda s: s["number"])
    declared = spec.get("deck", {}).get("slide_count")

    # 頁數
    if len(prs.slides) != len(spec_slides):
        fails.append(f"頁數不符:pptx {len(prs.slides)} 頁,spec {len(spec_slides)} 頁")
    if declared is not None and declared != len(prs.slides):
        fails.append(f"deck.slide_count={declared} 與 pptx {len(prs.slides)} 頁不符")

    # Section
    if has_sections(prs):
        fails.append("殘留 PowerPoint Section 分組(style_guide 禁止)")

    # autofit 框的容許高度:模板同一個框裝設計師原文時需要多高。
    # light 有多處版位存檔框高小於原文所需(靠 PowerPoint 自動長高排版),
    # 不放寬的話這些框每次都會誤報溢出——狼來了會蓋掉真的問題。
    tpl_need, tpl_collide, tpl_size, tpl_bounds = {}, {}, {}, {}
    if pack is not None:
        try:
            tpl = Presentation(str(pack.resolve_template()))
            for pt, ent in pack.page_types.items():
                pg = ent.get("template_page")
                if ent.get("mode") != "fill" or not pg or pg > len(tpl.slides):
                    continue
                tpl_need[pt] = {s.shape_id: tt.estimate_overflow(s)["needed_pt"]
                                for s in tt.iter_text_shapes(tpl.slides[pg - 1].shapes)
                                if tt.shape_text(s).strip()}
                # 碰撞基準:模板原文在同一對形狀上就已經重疊多少。判準是
                # 「不得比設計師自己的版面更糟」——模板有幾處刻意疊在一起
                # (p47 徽章與副字、p29 的 44pt 數字),絕對零重疊會誤殺。
                tpl_collide[pt] = {frozenset((a.shape_id, b.shape_id)): area
                                   for area, a, b in tt.text_collisions(tpl.slides[pg - 1])}
                # 設計字級基準:字級是設計過的,產出不得比模板小(見下方檢查)。
                tpl_size[pt] = {s.shape_id: tt._first_run_size_pt(s)
                                for s in tt.iter_text_shapes(tpl.slides[pg - 1].shapes)
                                if tt.shape_text(s).strip()}
                # 出界基準:模板自己也有框稍微掛在頁緣外(p47 欄三標題右緣 13.45 吋),
                # 判準同碰撞——不得比設計師的版面更糟,不是絕對零出界。
                tpl_bounds[pt] = {s.shape_id: (L, T, W, H)
                                  for s, L, T, W, H in tt.walk_absolute(tpl.slides[pg - 1].shapes)}
        except Exception:
            tpl_need, tpl_collide, tpl_size, tpl_bounds = {}, {}, {}, {}  # 取不到模板就退回原本判準

    # 逐頁:槽位文字覆蓋 / 頁碼 / 字體 / 溢出
    overflow_all = []
    for i, spec_slide in enumerate(spec_slides):
        if i >= len(prs.slides):
            break
        slide = prs.slides[i]
        num = spec_slide["number"]
        page_text = slide_all_text(slide)

        wanted = []
        collect_strings(spec_slide.get("slots", {}), wanted)
        missing = [w for w in wanted if norm(w) not in page_text]
        for m in missing:
            fails.append(f"p{num}: 內容未出現「{m[:30]}{'…' if len(m) > 30 else ''}」")

        digit_shapes = [s for s in tt.iter_text_shapes(slide.shapes)
                        if tt.shape_text(s).strip().isdigit()
                        and len(tt.shape_text(s).strip()) <= 3
                        and (s.top or 0) > Inches(detect_zone["top"])
                        and (s.left or 0) > Inches(detect_zone["left"])]
        if spec_slide.get("render_page_number"):
            if not any(tt.shape_text(s).strip() == str(num) for s in digit_shapes):
                warns.append(f"p{num}: 找不到右下角頁碼 {num}")
        elif digit_shapes:
            warns.append(f"p{num}: 此頁型不應有頁碼,但右下角有數字框")

        for shp in tt.iter_text_shapes(slide.shapes):
            for para in shp.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.name and not _font_ok(run.font.name, allowed_fonts):
                        warns.append(f"p{num}: 字體 {run.font.name!r}(shape id={shp.shape_id})")
                        break

        # 出界:文字跑到投影片外。舊版三道檢查全看不見這件事——溢出是比對
        # 「文字 vs 自己的框」、碰撞是比對「文字 vs 鄰居」,一個被推到頁面外的
        # 框兩者都不成立(2026-08-03:綁定 resize 誤把絕對座標寫進群組子座標,
        # 第三個 KPI 數字落在 x=12.60、頁寬只有 13.33,而 qa 印 PASS)。
        for shp, L, T, W, H in tt.walk_absolute(slide.shapes):
            if not getattr(shp, "has_text_frame", False) or not tt.shape_text(shp).strip():
                continue
            out = max(-L, -T, L + W - slide_w, T + H - slide_h)
            if out <= 0.05:
                continue
            ref = tpl_bounds.get(spec_slide.get("page_type"), {}).get(shp.shape_id)
            if ref:
                rl, rt, rw, rh = ref
                if out <= max(-rl, -rt, rl + rw - slide_w, rt + rh - slide_h) + 0.05:
                    continue        # 模板本來就這樣掛在頁緣,不比它更糟
            fails.append(
                f"p{num}: 形狀跑出投影片 {out:.2f} 吋"
                f"(x {L:.2f}–{L + W:.2f}、y {T:.2f}–{T + H:.2f};頁面 {slide_w:.2f}×{slide_h:.2f})"
                f"(shape id={shp.shape_id}「{tt.shape_text(shp)[:12]}」)")

        # 字級:設計師定調「字級是被設計過的,塞不下要改稿或換頁型,不准縮字」。
        # 渲染器仍保留 shrink_to_fit 當最後手段,但它**一縮就是靜默的**——
        # 2026-07-26 事故正是「99 個框被縮、所有自動檢查全綠、目檢才發現」。
        # 這裡把它變成 FAIL:縮字代表該槽位的容量上限量錯了,要回頭跑 fit,
        # 不是讓一份字級不一致的檔案出門。
        for shp in tt.iter_text_shapes(slide.shapes):
            if not tt.shape_text(shp).strip():
                continue
            design = tpl_size.get(spec_slide.get("page_type"), {}).get(shp.shape_id)
            got = tt._first_run_size_pt(shp)
            if design and got and got < design - 0.5:
                fails.append(
                    f"p{num}: 字級被縮小 {design:.0f}pt → {got:.0f}pt"
                    f"(shape id={shp.shape_id}「{tt.shape_text(shp)[:12]}」)"
                    "——該槽位容量上限量錯了,跑 template_admin.py fit 重量")

        for shp in tt.iter_text_shapes(slide.shapes):
            if not tt.shape_text(shp).strip():
                continue
            est = tt.estimate_overflow(shp)
            if est["fits"]:
                continue
            allow = est["avail_pt"]
            if tt.has_autofit(shp):     # 框會自動長高:放寬到模板原文需要的高度
                allow = max(allow, tpl_need.get(spec_slide.get("page_type"), {})
                                            .get(shp.shape_id, 0))
            if est["needed_pt"] <= allow * 1.08:
                continue
            overflow_all.append((est["needed_pt"] / max(allow, 1), num,
                                 shp.shape_id, tt.shape_text(shp)[:20]))

        # 文字碰撞:真正的破版不是「裝不進自己的框」,而是 autofit/wrap=none
        # 之後文字跑到不該出現的位置(壓到鄰欄)。「裝得下」全綠但版面已壞的
        # 案例是設計師目檢 golden 抓到的(2026-07-26,light p10/p22/p30)。
        base = tpl_collide.get(spec_slide.get("page_type"), {})
        for area, a, b in tt.text_collisions(slide):
            ref = base.get(frozenset((a.shape_id, b.shape_id)), 0.0)
            if area > max(ref * 1.10, 0.03):
                fails.append(
                    f"p{num}: 文字壓到別的元素 {area:.2f} 平方吋"
                    f"(id={a.shape_id}「{tt.shape_text(a)[:12]}」 ×"
                    f" id={b.shape_id}「{tt.shape_text(b)[:12]}」)"
                    + (f",模板基準僅 {ref:.2f}" if ref else ""))

    warns = list(dict.fromkeys(warns))  # 去重(字體警告易重複)
    # 溢出:先報總量與受影響頁,再列最嚴重前 5。**不可只印前 5 就當全部**——
    # 舊版靜默截斷曾讓 79 條溢出對外顯示成 5 條,PASS 報告因此嚴重低估問題規模。
    if overflow_all:
        pages = sorted({num for _, num, _, _ in overflow_all})
        warns.append(f"溢出疑慮共 {len(overflow_all)} 條,分布於 {len(pages)} 頁"
                     f"(p{', p'.join(map(str, pages))})— 沙箱估算,請逐頁開檔確認")
    shown = sorted(overflow_all, reverse=True)
    if not args.overflow_all:
        shown = shown[:5]
    for ratio, num, sid, txt in shown:
        warns.append(f"p{num}: 溢出疑慮 x{ratio:.1f}(id={sid}「{txt}…」)— 沙箱估算,請開檔確認")
    if len(overflow_all) > len(shown):
        warns.append(f"(另有 {len(overflow_all) - len(shown)} 條較輕微溢出未列出,"
                     f"完整清單請加 --overflow-all)")

    if fails:
        print(f"✗ FAIL ({len(fails)})")
        for f in fails:
            print(f"   [F] {f}")
    if warns:
        print(f"⚠ WARN ({len(warns)})")
        for w in warns:
            print(f"   [W] {w}")
    if fails:
        print("結果:FAIL — 修 render_plan.json / slide_spec.json 後重跑 render_deck,勿手改 pptx")
        return 1
    print(f"結果:PASS({len(prs.slides)} 頁"
          f"{',' + str(len(warns)) + ' 個警告' if warns else ''}"
          f"{',模板包 ' + pack_note if pack_note else ''})— 可交付")
    return 0


if __name__ == "__main__":
    sys.exit(main(sys.argv[1:]))
