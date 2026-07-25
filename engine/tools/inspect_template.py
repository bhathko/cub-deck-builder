#!/usr/bin/env python3
"""inspect_template — 模板盤點器:導出每頁形狀樹的精簡 JSON。

用途:
1. 寫 render_plan 前,查某參考頁有哪些文字框(拿 id / 現有文字當替換錨點)。
2. 一次性排查哪些頁含 python-pptx 動不了的元素(SmartArt 等)→ 那些頁只能走
   幾何重建或回報使用者。

用法(省 token 的推薦姿勢):
  python inspect_template.py --pptx /mnt/data/light_template.pptx --summary
      # 全簡報一頁一行:形狀數 + 風險旗標。先跑這個。
  python inspect_template.py --pptx ... --page 17
      # 只印第 17 頁(1-based)的形狀樹 JSON。寫該頁 plan 前跑這個。
  python inspect_template.py --pptx ... --all --out /mnt/data/template_report.json
      # 全量導出到檔案(不印),之後用 --page 或讀檔查詢。

輸出欄位(精簡鍵名):
  id=shape_id nm=name t=型別 x,y,w,h=吋(2位) txt=文字(截80字) fs=首run字級
  pic/tbl/cht/smart/grp=旗標;grp 時 ch=子形狀陣列
"""
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

EMU_PER_IN = 914400


def _inches(v) -> float:
    return round(v / EMU_PER_IN, 2) if v is not None else None


def _shape_info(shp, with_text=True) -> dict:
    info = {
        "id": shp.shape_id,
        "nm": shp.name,
        "t": str(shp.shape_type).split(" ")[0],
        "x": _inches(shp.left),
        "y": _inches(shp.top),
        "w": _inches(shp.width),
        "h": _inches(shp.height),
    }
    st = shp.shape_type
    if st == MSO_SHAPE_TYPE.PICTURE:
        info["pic"] = 1
    if getattr(shp, "has_table", False):
        info["tbl"] = 1
    if getattr(shp, "has_chart", False):
        info["cht"] = 1
    if st == MSO_SHAPE_TYPE.GROUP:
        info["grp"] = 1
        info["ch"] = [_shape_info(c, with_text) for c in shp.shapes]
    elif "GraphicFrame" in info["t"] and not info.get("tbl") and not info.get("cht"):
        info["smart"] = 1  # 非表格非圖表的 graphicFrame:SmartArt/OLE,pptx 動不了
    if with_text and getattr(shp, "has_text_frame", False):
        txt = shp.text_frame.text
        if txt.strip():
            info["txt"] = txt[:80]
            for para in shp.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size is not None:
                        info["fs"] = run.font.size.pt
                        break
                if "fs" in info:
                    break
    return info


def _flags(shapes) -> dict:
    out = {"n": 0, "txt": 0, "pic": 0, "tbl": 0, "cht": 0, "smart": 0}
    def walk(ss):
        for s in ss:
            out["n"] += 1
            if s.shape_type == MSO_SHAPE_TYPE.GROUP:
                walk(s.shapes)
                continue
            if s.shape_type == MSO_SHAPE_TYPE.PICTURE:
                out["pic"] += 1
            if getattr(s, "has_table", False):
                out["tbl"] += 1
            if getattr(s, "has_chart", False):
                out["cht"] += 1
            if ("GraphicFrame" in str(s.shape_type) and not getattr(s, "has_table", False)
                    and not getattr(s, "has_chart", False)):
                out["smart"] += 1
            if getattr(s, "has_text_frame", False) and s.text_frame.text.strip():
                out["txt"] += 1
    walk(shapes)
    return out


def _verify_pack(pack_dir):
    """比對包內 inventory.json 與現模板檔:id 增刪/幾何漂移即 exit 1
    (「模板改版必重盤點」的機器化;流程見 templates/TEMPLATE_LIFECYCLE.md)。"""
    import hashlib
    pack_dir = Path(pack_dir)
    inv = json.loads((pack_dir / "inventory.json").read_text(encoding="utf-8-sig"))
    manifest = json.loads((pack_dir / "manifest.json").read_text(encoding="utf-8-sig"))
    tpl = pack_dir / manifest.get("template_file", "template.pptx")
    if not tpl.exists():  # light 特例:模板檔可能在包外(prepare_env/沙箱佈局)
        print(f"✗ 找不到模板檔 {tpl};用 --pptx 指定路徑後手動比對")
        return 2
    sha = hashlib.sha256(tpl.read_bytes()).hexdigest()
    drift = []
    if sha != inv.get("template_sha256"):
        drift.append(f"pptx sha256 不符 inventory(檔案已改版)")
    prs = Presentation(str(tpl))
    for page_str, snap in inv.get("pages", {}).items():
        pg = int(page_str)
        if pg > len(prs.slides):
            drift.append(f"p{pg}: 頁不存在(模板只剩 {len(prs.slides)} 頁)")
            continue
        now = [_shape_info(s) for s in prs.slides[pg - 1].shapes]

        def flat(items, acc):
            for it in items:
                acc[it["id"]] = it
                flat(it.get("ch", []), acc)
            return acc
        old_by_id, new_by_id = flat(snap, {}), flat(now, {})
        for sid in sorted(set(old_by_id) - set(new_by_id)):
            drift.append(f"p{pg}: shape id {sid} 消失")
        for sid in sorted(set(new_by_id) - set(old_by_id)):
            drift.append(f"p{pg}: 新增 shape id {sid}")
        for sid in sorted(set(old_by_id) & set(new_by_id)):
            o, n = old_by_id[sid], new_by_id[sid]
            geo = [(k, o.get(k), n.get(k)) for k in ("x", "y", "w", "h")
                   if o.get(k) != n.get(k)]
            if geo:
                drift.append(f"p{pg}: id {sid} 幾何漂移 " +
                             " ".join(f"{k}:{a}→{b}" for k, a, b in geo))
    if drift:
        print(f"✗ 盤點漂移 {len(drift)} 項(重盤點並核對綁定後更新 inventory/manifest):")
        for d in drift:
            print(f"   {d}")
        return 1
    print(f"盤點一致:{len(inv.get('pages', {}))} 個綁定頁 shape 樹與 inventory 相符,sha 相符")
    return 0


def main(argv):
    ap = argparse.ArgumentParser()
    ap.add_argument("--pptx")
    ap.add_argument("--page", type=int, help="只看這一頁(1-based)")
    ap.add_argument("--summary", action="store_true", help="全簡報一頁一行摘要")
    ap.add_argument("--all", action="store_true", help="全量導出(配 --out)")
    ap.add_argument("--out", help="輸出 JSON 檔路徑")
    ap.add_argument("--verify", metavar="PACK_DIR",
                    help="比對模板包 inventory 與現模板檔(漂移 exit 1)")
    args = ap.parse_args(argv)

    if args.verify:
        return _verify_pack(args.verify)
    if not args.pptx:
        print("✗ 需要 --pptx(或 --verify <包目錄>)")
        return 2

    prs = Presentation(args.pptx)
    n = len(prs.slides)

    if args.page:
        if not (1 <= args.page <= n):
            print(f"✗ 頁碼超出範圍 1–{n}")
            return 2
        slide = prs.slides[args.page - 1]
        data = {"page": args.page, "shapes": [_shape_info(s) for s in slide.shapes]}
        text = json.dumps(data, ensure_ascii=False, separators=(",", ":"))
        if args.out:
            with open(args.out, "w", encoding="utf-8") as f:
                f.write(text)
            print(f"寫入 {args.out}")
        else:
            print(text)
        return 0

    if args.all:
        data = [{"page": i + 1, "shapes": [_shape_info(s) for s in sl.shapes]}
                for i, sl in enumerate(prs.slides)]
        out = args.out or "template_report.json"
        with open(out, "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, separators=(",", ":"))
        print(f"共 {n} 頁,導出 {out}(之後用 --page N 或讀檔查單頁,不要整檔印出)")
        return 0

    # 預設 / --summary:一頁一行
    for i, sl in enumerate(prs.slides):
        f = _flags(sl.shapes)
        risk = " ⚠SMARTART" if f["smart"] else ""
        print(f"p{i+1:>2}: shapes={f['n']:>3} text={f['txt']:>3} pic={f['pic']} "
              f"tbl={f['tbl']} cht={f['cht']}{risk}")
    print(f"-- 共 {n} 頁;⚠SMARTART 頁無法用「複製改字」,需幾何重建或回報使用者")
    return 0


if __name__ == "__main__":
    sys.exit(main(sys.argv[1:]))
