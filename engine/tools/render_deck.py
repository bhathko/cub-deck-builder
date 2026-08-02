#!/usr/bin/env python3
"""render_deck — 通用 renderer:spec(+選配 plan)→ .pptx(冪等,全檔重生)。

核心原則:**修計畫,不修產出**。每次執行都從模板重新生成整份簡報;
產出不對就改輸入(spec 或 render_plan.json)再跑一次,禁止對產出 pptx
做局部修補(那正是「刪一個壞的再產一個壞的」循環的來源)。

多模板架構(docs/ARCHITECTURE.md):頁型的自動產出實作(fill 填充)住在
模板包 `templates/<id>/bindings.json`,本工具經 pack_loader 載入選定包後
dispatch;預設包 = light(spec 未寫 deck.template 時)。

★ 註冊頁型完全自動,不需要 plan:
   fill = clone 模板實頁再照綁定填字,綁定在該包 bindings.json。
   → 整份 spec 都是註冊頁型時,直接跑本工具即可,LLM 不需要產任何計畫。
   哪些頁型屬於哪一級**依模板包而定**,跑 `make_skeleton.py --list` 查當下實況;
   這裡刻意不列清單/數量,寫死了就會隨模板加開頁型而過期。

★ 只有「未涵蓋頁型」(page_types.md 頁型庫其他頁型)才需要 plan 條目:
   clone 模板頁 + 文字替換編輯清單。plan 也可覆寫自動頁(以 plan 為準)。

用法:
  python render_deck.py --spec /mnt/data/slide_spec.json \
      [--plan /mnt/data/render_plan.json] \
      [--template-pack light] [--packs-root /mnt/data/templates] \
      [--template <pptx 路徑覆寫>] \
      --asset-dir /mnt/data --out /mnt/data/deck.pptx
  模板檔預設 = 選定包的 template.pptx(包目錄優先、asset-dir 兜底);
  --template 為相容別名/試模板覆寫(WORKLOG §9 工作流)。

render_plan.json 格式(僅未涵蓋頁需要;內容文字必須逐字取自 spec):
{
  "slides": [
    {"number": 5, "mode": "clone", "template_page": 35,
     "edits": [
       {"match": {"id": 23}, "text": "..."},                  // shape_id 最精準
       {"match": {"contains": "模板現有文字"}, "text": "..."},
       {"match": {"name": "Title 1"}, "text": "...", "nth": 0} // 撞多筆用 nth
     ],
     "delete": [ {"match": {"id": 45}} ],
     "shrink": true}
  ]
}

行為保證:
- 有 UNMATCHED / AMBIGUOUS / FillError → 印出清單、exit 1;檔案仍會存
  (方便檢查)但**不得交付**,修輸入後整檔重跑。
- 頁碼規則自動:內容頁右下 28pt #344252;封面/封底無;clone/fill 頁會先清掉
  模板原本右下角的純數字頁碼框。
- 產出只含 spec 的頁:模板原頁全刪、Section 全清、順序=spec 順序。
- 模板檔與選定包 manifest 的 sha 不符時印警告(自訂模板測試情境;正式發版
  前必須照 templates/TEMPLATE_LIFECYCLE.md 重盤點)。
"""
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

import fill_helpers
import pack_loader
import pptx_toolkit as tk
import text_tools as tt

DARK = RGBColor(*tk.COLOR_DARK)


# ---------------------------------------------------------------------------
# 頁碼(幾何讀選定包 manifest 的 page_number;預設值 = light 幾何)
# ---------------------------------------------------------------------------
_PN_DEFAULT = {"box_in": [12.30, 6.72, 0.70, 0.50], "size_pt": 28, "color": "344252",
               "clear_zone_in": {"left": 11.2, "top": 6.3}}


def _pn_cfg(pack):
    cfg = dict(_PN_DEFAULT)
    cfg.update(pack.manifest.get("page_number") or {})
    return cfg


def _tight_margins(tf):
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)


def _textbox(slide, text, x, y, w, h, size, bold=False, color=DARK, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    _tight_margins(tf)
    for i, line in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if align is not None:
            p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = tk.FONT_ZH
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def _page_number(slide, number: int, cfg=_PN_DEFAULT):
    x, y, w, h = cfg["box_in"]
    _textbox(slide, str(number), x, y, w, h, size=cfg["size_pt"], bold=True,
             color=RGBColor.from_string(cfg["color"]), align=PP_ALIGN.RIGHT)


def _finalize_page_number(slide, spec_slide, cfg=_PN_DEFAULT):
    """清掉模板殘留的右下角純數字頁碼,依 spec 決定補標準頁碼。"""
    zone = cfg["clear_zone_in"]
    for s in list(tt.iter_text_shapes(slide.shapes)):
        txt = tt.shape_text(s).strip()
        if (txt.isdigit() and len(txt) <= 3 and s.top and s.left
                and s.top > Inches(zone["top"]) and s.left > Inches(zone["left"])):
            s._element.getparent().remove(s._element)
    if spec_slide.get("render_page_number"):
        _page_number(slide, spec_slide["number"], cfg)


# ---------------------------------------------------------------------------
# plan 驅動的 clone 頁(未涵蓋頁型用)
# ---------------------------------------------------------------------------
def _match_shapes(slide, match: dict):
    if "id" in match or ("name" in match and "contains" not in match):
        found = []

        def walk(shapes):
            for s in shapes:
                if "id" in match and s.shape_id == match["id"]:
                    found.append(s)
                elif "name" in match and s.name == match["name"]:
                    found.append(s)
                if getattr(s, "shape_type", None) is not None and hasattr(s, "shapes"):
                    walk(s.shapes)

        walk(slide.shapes)
        return found
    return tt.find_text_shapes(slide, contains=match.get("contains"),
                               name=match.get("name"))


def apply_clone_plan(prs, spec_slide, page_plan, template_page_count, problems,
                     pn_cfg=_PN_DEFAULT):
    tpl = page_plan.get("template_page")
    num = spec_slide["number"]
    if not tpl or not (1 <= tpl <= template_page_count):
        problems.append(f"p{num}: template_page 缺少或超出範圍 1–{template_page_count}")
        return None
    slide = tk.clone_slide(prs, tpl - 1)
    edited = []

    for d in page_plan.get("delete", []):
        hits = _match_shapes(slide, d.get("match", {}))
        if not hits:
            problems.append(f"p{num}: delete UNMATCHED {d.get('match')}")
            continue
        for s in hits:
            s._element.getparent().remove(s._element)

    for e in page_plan.get("edits", []):
        hits = sorted(_match_shapes(slide, e.get("match", {})),
                      key=lambda s: (s.top or 0, s.left or 0))
        nth = e.get("nth")
        if not hits:
            problems.append(f"p{num}: edit UNMATCHED {e.get('match')}")
            continue
        if len(hits) > 1 and nth is None:
            problems.append(f"p{num}: edit AMBIGUOUS x{len(hits)} {e.get('match')} → 加 nth 或改用 id")
            continue
        target = hits[nth or 0] if nth is not None else hits[0]
        if not getattr(target, "has_text_frame", False):
            problems.append(f"p{num}: edit 目標不是文字框 {e.get('match')}")
            continue
        tt.set_text_keep_style(target, e["text"])
        edited.append(target)

    if page_plan.get("shrink", True):
        for s in edited:
            tt.shrink_to_fit(s, min_pt=12)
    _finalize_page_number(slide, spec_slide, pn_cfg)
    return slide


# ---------------------------------------------------------------------------
# 主流程
# ---------------------------------------------------------------------------
def main(argv):
    ap = argparse.ArgumentParser()
    ap.add_argument("--spec", required=True)
    ap.add_argument("--plan", help="選配:僅未涵蓋頁型需要")
    ap.add_argument("--template", help="模板 pptx 路徑(相容別名;省略=選定包的模板檔)")
    ap.add_argument("--template-pack", help="模板包 id 或目錄(省略=spec 的 deck.template,再省略=light)")
    ap.add_argument("--packs-root", help="模板包根目錄(預設=tools 上層的 templates/)")
    ap.add_argument("--asset-dir", default="/mnt/data")
    ap.add_argument("--out", required=True)
    args = ap.parse_args(argv)

    # utf-8-sig:容忍 Windows 工具(記事本等)寫入的 BOM,使用者上傳的 JSON 常見
    spec = json.loads(Path(args.spec).read_text(encoding="utf-8-sig"))
    plan_by_num = {}
    if args.plan:
        plan = json.loads(Path(args.plan).read_text(encoding="utf-8-sig"))
        plan_by_num = {p["number"]: p for p in plan.get("slides", [])}
    asset_dir = Path(args.asset_dir)

    try:
        pack = pack_loader.load_pack(pack_arg=args.template_pack,
                                     spec_deck=spec.get("deck"),
                                     packs_root=args.packs_root)
    except pack_loader.PackError as e:
        print(f"✗ 模板包載入失敗:{e}")
        return 2
    template_path = args.template or pack.resolve_template(asset_dir)
    if not template_path or not Path(template_path).exists():
        print(f"✗ 找不到模板檔:{template_path or pack.dir / pack.manifest.get('template_file', 'template.pptx')}"
              f"(模板包 {pack.id};可用 --template 指定路徑)")
        return 2
    if not pack.template_hash_matches(template_path):
        if args.template:
            print(f"⚠ 模板檔與模板包 {pack.id}@{pack.version} 的 manifest sha 不符"
                  "(--template 明示指定,視為自訂模板測試;正式發版前必須照"
                  " templates/TEMPLATE_LIFECYCLE.md 重盤點)")
        else:
            print(f"✗ 模板檔與模板包 {pack.id}@{pack.version} 的 manifest sha 不符:"
                  f"{template_path}\n"
                  "  多半是沙箱殘留的舊版模板或包沒同步乾淨:重新解壓該包,"
                  "或改過模板就先跑 freeze。要用自訂模板測試請以 --template 明示指定。")
            return 2
    pn_cfg = _pn_cfg(pack)

    spec_by_num = {s["number"]: s for s in spec["slides"]}
    problems = []
    modes = []

    prs = Presentation(str(template_path))
    n_template = len(prs.slides)

    for num in sorted(spec_by_num):
        spec_slide = spec_by_num[num]
        pt = spec_slide["page_type"]
        page_plan = plan_by_num.get(num)

        if page_plan:  # plan 條目優先(可覆寫自動頁)
            mode = page_plan.get("mode")
            if mode == "clone":
                apply_clone_plan(prs, spec_slide, page_plan, n_template, problems, pn_cfg)
                modes.append(f"p{num}:clone{page_plan.get('template_page')}")
            else:
                problems.append(f"p{num}: 未知 mode {mode!r}(僅支援 clone)")
        elif pt in pack.fills:
            tpl_page, fn = pack.fills[pt]
            slide = tk.clone_slide(prs, tpl_page - 1)
            ctx = fill_helpers.Ctx(slide)
            try:
                fn(ctx, spec_slide)
                ctx.shrink_edited(min_pt=12)
            except fill_helpers.FillError as e:
                problems.append(f"p{num}: {e}(模板包 {pack.id}@{pack.version})")
            _finalize_page_number(slide, spec_slide, pn_cfg)
            modes.append(f"p{num}:auto{tpl_page}")
        else:
            problems.append(
                f"p{num}: 頁型 {pt!r} 無自動支援 → 需要 plan 條目"
                f"(mode=clone + template_page + edits,參考 page_types.md 與"
                f"模板包 page_map.md)")

    tk.delete_slides(prs, range(n_template))
    sections_removed = tk.clear_sections(prs)
    prs.save(args.out)

    print(f"產出:{args.out}  頁數:{len(spec_by_num)}  清除Section:{sections_removed}")
    print(f"模板包:{pack.id}@{pack.version}")
    print("模式:" + " ".join(modes))
    if problems:
        print(f"✗ 問題 {len(problems)} 項(修 spec/plan 後整檔重跑,勿手改 pptx):")
        for p in problems:
            print(f"   {p}")
        return 1
    print("結果:OK — 接著跑 qa_check.py")
    return 0


if __name__ == "__main__":
    sys.exit(main(sys.argv[1:]))
