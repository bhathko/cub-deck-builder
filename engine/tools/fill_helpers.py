#!/usr/bin/env python3
"""fill_helpers — 填充引擎共用件(模板無關;各模板包 bindings 與引擎共用)。

自原 fills.py 抽出:FillError / index_shapes / Ctx / fill_rows(原 _fill_rows
公開化,列高參數化),並把 p33/p54 兩處內聯 textbox 樣板整併為
add_styled_textbox。本檔不得含任何模板專屬常數(shape id、座標、色碼);
那些屬於 engine/templates/<id>/ 的綁定資料,詳見 docs/ARCHITECTURE.md。
"""
from __future__ import annotations

from pathlib import Path

from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Pt

import text_tools as tt


def resolve_asset(rel, asset_dir, pack_dir, pack_first=False):
    """素材路徑解析鏈:asset_dir 與模板包目錄二擇一命中(順序由 pack_first 決定)。

    light 包 = asset_dir 優先(舊 spec 的 assets/... 路徑不變);非 light 包
    預設包目錄優先,防止沙箱裡 light 的 assets/ 遮蔽新包素材(ARCHITECTURE §5)。
    全部不存在時回傳第一候選(讓呼叫端的存在性檢查印出主要路徑)。
    """
    candidates = [Path(pack_dir) / rel, Path(asset_dir) / rel] if pack_first \
        else [Path(asset_dir) / rel, Path(pack_dir) / rel]
    for c in candidates:
        if c.exists():
            return c
    return candidates[0]


class FillError(Exception):
    pass


def index_shapes(slide) -> dict:
    """shape_id → shape(含群組內層)。clone 保留原模板 id,故可寫死引用。"""
    idx = {}

    def walk(shapes):
        for s in shapes:
            idx[s.shape_id] = s
            if s.shape_type == MSO_SHAPE_TYPE.GROUP:
                walk(s.shapes)

    walk(slide.shapes)
    return idx


class Ctx:
    """單頁填充的工作環境:查 id、換字、刪形狀,並記錄動過的框供縮字。"""

    def __init__(self, slide):
        self.slide = slide
        self.idx = index_shapes(slide)
        self.edited = []

    def shape(self, sid):
        if sid not in self.idx:
            raise FillError(f"shape id {sid} 不存在(模板改版?請重新盤點並更新綁定)")
        return self.idx[sid]

    def set(self, sid, text):
        shp = self.shape(sid)
        tt.set_text_keep_style(shp, text)
        self.edited.append(shp)

    def delete(self, *sids):
        for sid in sids:
            shp = self.idx.pop(sid, None)
            if shp is not None:
                el = shp._element
                parent = el.getparent()
                if parent is not None:
                    parent.remove(el)

    def shrink_edited(self, min_pt=12.0):
        for shp in self.edited:
            try:
                tt.shrink_to_fit(shp, min_pt=min_pt)
            except Exception:
                pass  # 縮字失敗不致命,交給 qa_check 的溢出警告


def fill_rows(ctx, points, row_text_ids, sep_ids, merge_height_in=0.80):
    """把 points 填進固定列(不足刪列與分隔線,超出併入最後一列並加高)。"""
    n_rows = len(row_text_ids)
    for i, sid in enumerate(row_text_ids):
        if i < len(points):
            text = points[i]
            if i == n_rows - 1 and len(points) > n_rows:  # 溢出併入最後一列
                text = "\n".join(points[i:])
                ctx.shape(sid).height = Inches(merge_height_in)
            ctx.set(sid, text)
        else:
            ctx.delete(sid)
            if i - 1 < len(sep_ids):
                ctx.delete(sep_ids[i - 1])


def add_styled_textbox(slide, text, x, y, w, h, *, font, size_pt, bold=False):
    """動態補建單段文字框(模板缺佔位符時用;字型/字級由綁定端提供)。"""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size_pt)
    if bold:
        run.font.bold = True
    return box
