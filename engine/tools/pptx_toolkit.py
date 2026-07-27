#!/usr/bin/env python3
"""pptx_toolkit — 投影片層級的機械操作(複製/刪除/排序/清 Section)。

這些操作 python-pptx 沒有官方 API,社群 recipe 又常漏掉 relationship 重映射
(症狀:複製出來的頁圖片破圖)。統一封裝在這裡,渲染流程一律 import 使用,
禁止在對話中現寫等價程式碼。

依賴:python-pptx(ChatGPT Code Interpreter 內建)。
"""
from __future__ import annotations

import copy

from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn

R_NS = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"

# 頁碼文字框用的兩個常數(唯一消費者是 render_deck)。其餘品牌色不放這裡:
# 色碼是模板專屬知識,真相在各包 manifest 的 style.colors 與 style_guide.md
# ——原本並列的 6 個 accent 常數在 2026-07-26 已無任何消費者,一併移除。
FONT_ZH = "Microsoft JhengHei"
COLOR_DARK = (0x34, 0x42, 0x52)


def _remap_rids(element, id_map: dict) -> None:
    """把元素樹中所有 r: 命名空間屬性(embed/link/id…)的舊 rId 換成新 rId。"""
    if not id_map:
        return
    for el in element.iter():
        for attr, val in list(el.attrib.items()):
            if attr.startswith(R_NS) and val in id_map:
                el.set(attr, id_map[val])


def _clone_chart_part(src_part):
    """深複製 chart part(含內嵌 xlsx);colors/style/userShapes 等樣式部件共用。

    背景:clone_slide 對 relationship 目標預設「共用」(圖片共用無害),但圖表
    會被 fills_engine 的 chart op 改寫——共用會讓同一參考頁 clone 出的多頁
    圖表互相覆寫(Phase 4 實測:min/max 兩頁都變成最後一次 replace_data 的
    數據)。故 chart part 與其內嵌 xlsx 必須逐頁複製。
    """
    package = src_part.package
    new_part = type(src_part).load(
        package.next_partname("/ppt/charts/chart%d.xml"),
        src_part.content_type, package, src_part.blob)
    id_map = {}
    for rid, rel in src_part.rels.items():
        if rel.is_external:
            new_rid = new_part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        elif rel.reltype == RT.PACKAGE:  # 內嵌 xlsx:可變件,一併複製
            tgt = rel.target_part
            xlsx = type(tgt).load(
                package.next_partname("/ppt/embeddings/Microsoft_Excel_Sheet%d.xlsx"),
                tgt.content_type, package, tgt.blob)
            new_rid = new_part.relate_to(xlsx, rel.reltype)
        else:
            new_rid = new_part.relate_to(rel.target_part, rel.reltype)
        id_map[rid] = new_rid
    _remap_rids(new_part._element, id_map)
    return new_part


def clone_slide(prs, src_index: int):
    """複製第 src_index 頁(0-based)成新投影片,附加在簡報最後。

    含:形狀樹 deepcopy、relationship(圖片等)複製與 rId 重映射、
    chart part 深複製(見 _clone_chart_part)、頁面背景(p:bg)複製。
    回傳新 slide 物件。
    """
    src = prs.slides[src_index]
    dest = prs.slides.add_slide(src.slide_layout)

    # 清掉 layout 帶進來的預設佔位符
    for shp in list(dest.shapes):
        shp._element.getparent().remove(shp._element)

    # 複製 relationships(略過 layout / notes),建立舊→新 rId 對照
    id_map = {}
    for rid, rel in src.part.rels.items():
        if rel.reltype in (RT.SLIDE_LAYOUT, RT.NOTES_SLIDE):
            continue
        if rel.is_external:
            new_rid = dest.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        elif rel.reltype == RT.CHART:
            new_rid = dest.part.relate_to(_clone_chart_part(rel.target_part), rel.reltype)
        else:
            new_rid = dest.part.relate_to(rel.target_part, rel.reltype)
        id_map[rid] = new_rid

    # 複製形狀
    for shp in src.shapes:
        el = copy.deepcopy(shp._element)
        _remap_rids(el, id_map)
        dest.shapes._spTree.append(el)

    # 複製頁面背景(若有)
    src_csld = src._element.find(qn("p:cSld"))
    src_bg = src_csld.find(qn("p:bg")) if src_csld is not None else None
    if src_bg is not None:
        dest_csld = dest._element.find(qn("p:cSld"))
        old_bg = dest_csld.find(qn("p:bg"))
        if old_bg is not None:
            dest_csld.remove(old_bg)
        new_bg = copy.deepcopy(src_bg)
        _remap_rids(new_bg, id_map)
        dest_csld.insert(0, new_bg)

    return dest


def delete_slides(prs, indices) -> int:
    """刪除多頁(0-based index,以呼叫當下的順序計)。回傳刪除數。"""
    sldIdLst = prs.slides._sldIdLst
    entries = list(sldIdLst)
    n = 0
    for i in sorted(set(indices), reverse=True):
        if i < 0 or i >= len(entries):
            continue
        rId = entries[i].get(qn("r:id"))
        prs.part.drop_rel(rId)
        sldIdLst.remove(entries[i])
        n += 1
    return n


def reorder_slides(prs, new_order) -> None:
    """依 new_order(0-based index 的新排列)重排投影片。"""
    sldIdLst = prs.slides._sldIdLst
    entries = list(sldIdLst)
    if sorted(new_order) != list(range(len(entries))):
        raise ValueError(f"new_order 必須是 0..{len(entries)-1} 的排列")
    for entry in entries:
        sldIdLst.remove(entry)
    for i in new_order:
        sldIdLst.append(entries[i])


def clear_sections(prs) -> int:
    """移除 PowerPoint Section 分組(style_guide 硬規定)。回傳移除的 ext 數。"""
    pres_el = prs._element  # noqa: SLF001 — 讀私有屬性,無公開 API
    removed = 0
    for extLst in pres_el.findall(qn("p:extLst")):
        for ext in list(extLst):
            has_section = any(
                child.tag.endswith("}sectionLst") or child.tag.endswith("}sectionPr")
                for child in ext.iter()
            )
            if has_section:
                extLst.remove(ext)
                removed += 1
        if len(extLst) == 0:
            pres_el.remove(extLst)
    return removed

